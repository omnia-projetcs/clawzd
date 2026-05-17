"""
Clawzd — Presentation & Rich Document Generation Tools.
Supports generating PowerPoint, PDF, and Word documents from structured JSON layouts.
"""

from fastapi import APIRouter, Request, HTTPException
import os
import uuid
import base64
import requests
from io import BytesIO
from datetime import datetime
from config import DATA_DIR
import logging

logger = logging.getLogger("clawzd.tools_presentation")

router = APIRouter()

def get_documents_dir():
    docs_dir = os.path.join(DATA_DIR, "documents")
    os.makedirs(docs_dir, exist_ok=True)
    return docs_dir

def get_presentations_dir():
    pres_dir = os.path.join(DATA_DIR, "presentations")
    os.makedirs(pres_dir, exist_ok=True)
    return pres_dir

def _resolve_image(src: str):
    """Resolve an image source (URL, relative path, or base64) to a BytesIO object."""
    if src.startswith("data:image"):
        header, encoded = src.split(",", 1)
        return BytesIO(base64.b64decode(encoded))
    elif src.startswith("http"):
        resp = requests.get(src)
        if resp.status_code == 200:
            return BytesIO(resp.content)
    elif src.startswith("/data/images/") or src.startswith("/data/screenshots/"):
        # Local file
        filename = src.split("/")[-1]
        folder = "screenshots" if "screenshots" in src else "images"
        filepath = os.path.join(DATA_DIR, folder, filename)
        if os.path.exists(filepath):
            with open(filepath, "rb") as f:
                return BytesIO(f.read())
    return None

def _hex_to_rgb(hex_color: str):
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    if not hex_color:
        return (0, 0, 0)
    try:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except (ValueError, IndexError):
        return (0, 0, 0)

def _generate_powerpoint(pages: list, filepath: str, canvas_width: int, canvas_height: int):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        prs = Presentation()
        # Set slide width/height
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        scale_x = prs.slide_width / canvas_width if canvas_width > 0 else 1
        scale_y = prs.slide_height / canvas_height if canvas_height > 0 else 1

        blank_slide_layout = prs.slide_layouts[6]

        for page in pages:
            slide = prs.slides.add_slide(blank_slide_layout)
            elements = page.get("elements", [])
            for el in elements:
                x = el.get("x", 0) * scale_x
                y = el.get("y", 0) * scale_y
                w = el.get("width", 100) * scale_x
                h = el.get("height", 50) * scale_y
                
                el_type = el.get("type", "text")
                if el_type == "text":
                    txBox = slide.shapes.add_textbox(x, y, w, h)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    is_list = el.get("isList", False)
                    lines = str(el.get("content", "")).split('\n')
                    
                    for i, line in enumerate(lines):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                            
                        p.text = line
                        if is_list:
                            p.level = 1
                            
                        if el.get("fontFamily"):
                            font_name = el.get("fontFamily").split(",")[0].replace("'", "").replace('"', '').strip()
                            p.font.name = font_name
                            
                        font_size = el.get("fontSize", 16)
                        p.font.size = Pt(font_size * 0.75) 
                        
                        color = _hex_to_rgb(el.get("color", "#000000"))
                        p.font.color.rgb = RGBColor(color[0], color[1], color[2])
                        
                        align = el.get("textAlign", "left")
                        if align == "center":
                            p.alignment = PP_ALIGN.CENTER
                        elif align == "right":
                            p.alignment = PP_ALIGN.RIGHT
                            
                        if el.get("isBold"): p.font.bold = True
                        if el.get("isItalic"): p.font.italic = True
                        if el.get("isUnderline"): p.font.underline = True
                        
                    bg = el.get("backgroundColor")
                    if bg and bg != "transparent":
                        c = _hex_to_rgb(bg)
                        txBox.fill.solid()
                        txBox.fill.fore_color.rgb = RGBColor(c[0], c[1], c[2])
                    
                    border = el.get("borderColor")
                    bw = el.get("borderWidth", 0)
                    if border and border != "transparent" and bw > 0:
                        c = _hex_to_rgb(border)
                        txBox.line.color.rgb = RGBColor(c[0], c[1], c[2])
                        txBox.line.width = Pt(bw)
                    
                elif el_type == "image":
                    src = el.get("src", "")
                    if src:
                        img_stream = _resolve_image(src)
                        if img_stream:
                            try:
                                op = el.get("opacity", 100)
                                if op < 100:
                                    from PIL import Image
                                    img_pil = Image.open(img_stream).convert("RGBA")
                                    alpha = img_pil.split()[3]
                                    alpha = alpha.point(lambda p: int(p * (op / 100.0)))
                                    img_pil.putalpha(alpha)
                                    out_stream = BytesIO()
                                    img_pil.save(out_stream, format="PNG")
                                    out_stream.seek(0)
                                    img_stream = out_stream
                                    
                                slide.shapes.add_picture(img_stream, x, y, w, h)
                            except Exception as e:
                                logger.warning(f"Failed to add image to slide: {e}")
                                
                elif el_type == "mermaid":
                    # Fallback to text for mermaid in PPTX
                    txBox = slide.shapes.add_textbox(x, y, w, h)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = "[Mermaid Diagram]\n" + str(el.get("content", ""))
                    p.font.size = Pt(12)
                    color = _hex_to_rgb(el.get("color", "#000000"))
                    p.font.color.rgb = RGBColor(color[0], color[1], color[2])
                    
                elif el_type == "shape":
                    from pptx.enum.shapes import MSO_SHAPE
                    shape_type = el.get("shapeType", "rect")
                    if shape_type == "circle":
                        st = MSO_SHAPE.OVAL
                    elif shape_type == "triangle":
                        st = MSO_SHAPE.ISOSCELES_TRIANGLE
                    elif shape_type == "hexagon":
                        st = MSO_SHAPE.HEXAGON
                    elif shape_type == "arrow":
                        st = MSO_SHAPE.RIGHT_ARROW
                    else:
                        st = MSO_SHAPE.RECTANGLE
                        
                    shp = slide.shapes.add_shape(st, x, y, w, h)
                    
                    bg = el.get("backgroundColor", "#000000")
                    op = el.get("opacity", 100)
                    if bg and bg != "transparent":
                        c = _hex_to_rgb(bg)
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = RGBColor(c[0], c[1], c[2])
                        if op < 100:
                            try:
                                from pptx.oxml.xmlchemy import OxmlElement
                                solidFill = shp.fill._fill._solidFill
                                if hasattr(solidFill, "srgbClr"):
                                    alpha = OxmlElement('a:alpha')
                                    alpha.set('val', str(int(op * 1000)))
                                    solidFill.srgbClr.append(alpha)
                            except Exception as e:
                                logger.warning(f"Failed to apply opacity to PPTX shape fill: {e}")
                    else:
                        shp.fill.background()
                        
                    border = el.get("borderColor", "#ffffff")
                    bw = el.get("borderWidth", 0)
                    if border and border != "transparent" and bw > 0:
                        c = _hex_to_rgb(border)
                        shp.line.color.rgb = RGBColor(c[0], c[1], c[2])
                        shp.line.width = Pt(bw)
                        if op < 100:
                            try:
                                from pptx.oxml.xmlchemy import OxmlElement
                                solidFill = shp.line.fill._fill._solidFill
                                if hasattr(solidFill, "srgbClr"):
                                    alpha = OxmlElement('a:alpha')
                                    alpha.set('val', str(int(op * 1000)))
                                    solidFill.srgbClr.append(alpha)
                            except Exception as e:
                                logger.warning(f"Failed to apply opacity to PPTX shape border: {e}")
                    else:
                        shp.line.fill.background()
                                
        prs.save(filepath)
        return filepath
    except ImportError:
        raise Exception("python-pptx is not installed.")

def _generate_pdf(pages: list, filepath: str, canvas_width: int, canvas_height: int):
    try:
        from fpdf import FPDF
        
        px_to_mm = 0.264583
        pdf_w = canvas_width * px_to_mm
        pdf_h = canvas_height * px_to_mm
        
        orientation = 'L' if pdf_w > pdf_h else 'P'
        # FPDF expects the format tuple as (portrait_width, portrait_height)
        format_tuple = (min(pdf_w, pdf_h), max(pdf_w, pdf_h))
        pdf = FPDF(orientation=orientation, unit='mm', format=format_tuple)
        pdf.set_auto_page_break(auto=False)
        
        for page in pages:
            pdf.add_page()
            elements = page.get("elements", [])
            for el in elements:
                x = el.get("x", 0) * px_to_mm
                y = el.get("y", 0) * px_to_mm
                w = el.get("width", 100) * px_to_mm
                h = el.get("height", 50) * px_to_mm
                
                el_type = el.get("type", "text")
                if el_type == "text":
                    color = _hex_to_rgb(el.get("color", "#000000"))
                    pdf.set_text_color(color[0], color[1], color[2])
                    
                    font_size_pt = el.get("fontSize", 16) * 0.75
                    pdf.set_font("Arial", size=font_size_pt)
                    
                    align_map = {"left": "L", "center": "C", "right": "R"}
                    align = align_map.get(el.get("textAlign", "left"), "L")
                    
                    bg = el.get("backgroundColor", "transparent")
                    fill = False
                    if bg != "transparent":
                        c = _hex_to_rgb(bg)
                        pdf.set_fill_color(c[0], c[1], c[2])
                        fill = True
                    
                    border = el.get("borderColor", "transparent")
                    bw = el.get("borderWidth", 0)
                    draw_border = 0
                    if border != "transparent" and bw > 0:
                        c = _hex_to_rgb(border)
                        pdf.set_draw_color(c[0], c[1], c[2])
                        pdf.set_line_width(bw * px_to_mm)
                        draw_border = 1
                    
                    pdf.set_xy(x, y)
                    content = el.get("content", "").encode('latin-1', 'replace').decode('latin-1')
                    
                    if el.get("isList"):
                        lines = content.split('\n')
                        content = '\n'.join([f"    • {line}" for line in lines])
                        
                    font_style = ""
                    if el.get("isBold"): font_style += "B"
                    if el.get("isItalic"): font_style += "I"
                    if el.get("isUnderline"): font_style += "U"
                    
                    font_family = "Arial"
                    if el.get("fontFamily"):
                        f_name = el.get("fontFamily").split(",")[0].replace("'", "").replace('"', '').strip().lower()
                        if f_name in ['serif', 'times new roman', 'georgia']:
                            font_family = 'Times'
                        elif f_name in ['monospace', 'courier new', 'courier', 'consolas']:
                            font_family = 'Courier'
                        elif f_name in ['webdings', 'wingdings']:
                            font_family = 'Symbol'
                    
                    try: pdf.set_font(font_family, style=font_style, size=font_size_pt)
                    except RuntimeError: pass
                    
                    pdf.multi_cell(w, font_size_pt * 0.35, txt=content, align=align, border=draw_border, fill=fill)
                    
                    try: pdf.set_font("Arial", style="", size=font_size_pt)
                    except RuntimeError: pass
                    
                elif el_type == "image":
                    src = el.get("src", "")
                    if src:
                        import tempfile
                        img_stream = _resolve_image(src)
                        if img_stream:
                            op = el.get("opacity", 100)
                            if op < 100:
                                try:
                                    from PIL import Image
                                    img_pil = Image.open(img_stream).convert("RGBA")
                                    alpha = img_pil.split()[3]
                                    alpha = alpha.point(lambda p: int(p * (op / 100.0)))
                                    img_pil.putalpha(alpha)
                                    out_stream = BytesIO()
                                    img_pil.save(out_stream, format="PNG")
                                    out_stream.seek(0)
                                    img_stream = out_stream
                                except Exception as e:
                                    logger.warning(f"Failed to apply opacity to PDF image: {e}")
                                    
                            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                                tmp.write(img_stream.read())
                                tmp_path = tmp.name
                            try:
                                pdf.image(tmp_path, x=x, y=y, w=w, h=h)
                            except Exception as e:
                                logger.warning(f"Failed to add image to PDF: {e}")
                            finally:
                                os.unlink(tmp_path)
                                
                elif el_type == "mermaid":
                    # Fallback to text for mermaid in PDF
                    try: pdf.set_font("Courier", style="", size=10)
                    except RuntimeError: pass
                    pdf.set_text_color(0, 0, 0)
                    pdf.set_xy(x, y)
                    content = "[Mermaid Diagram]\n" + str(el.get("content", "")).encode('latin-1', 'replace').decode('latin-1')
                    pdf.multi_cell(w, 4, txt=content, align='L', border=1, fill=False)

                elif el_type == "shape":
                    bg = el.get("backgroundColor", "transparent")
                    border = el.get("borderColor", "transparent")
                    bw = el.get("borderWidth", 0)
                    style = ''
                    if bg != "transparent":
                        c = _hex_to_rgb(bg)
                        pdf.set_fill_color(c[0], c[1], c[2])
                        style += 'F'
                    if border != "transparent" and bw > 0:
                        c = _hex_to_rgb(border)
                        pdf.set_draw_color(c[0], c[1], c[2])
                        pdf.set_line_width(bw * px_to_mm)
                        style += 'D'
                    if not style:
                        continue
                        
                    shape_type = el.get("shapeType", "rect")
                    op = el.get("opacity", 100)
                    
                    def draw_shape():
                        if shape_type == "rect":
                            pdf.rect(x, y, w, h, style=style)
                        elif shape_type == "circle":
                            pdf.ellipse(x, y, w, h, style=style)
                        elif shape_type == "triangle":
                            pdf.polygon([(x+w/2, y), (x+w, y+h), (x, y+h)], style=style)
                        elif shape_type == "hexagon":
                            pdf.polygon([(x+w/4, y), (x+w*3/4, y), (x+w, y+h/2), (x+w*3/4, y+h), (x+w/4, y+h), (x, y+h/2)], style=style)
                        elif shape_type == "arrow":
                            headStart = w * 0.6
                            bodyH = h * 0.4
                            bodyY1 = y + (h - bodyH) / 2
                            bodyY2 = bodyY1 + bodyH
                            pdf.polygon([(x, bodyY1), (x+headStart, bodyY1), (x+headStart, y), (x+w, y+h/2), (x+headStart, y+h), (x+headStart, bodyY2), (x, bodyY2)], style=style)
                            
                    if op < 100:
                        try:
                            with pdf.local_context(fill_opacity=op/100.0, stroke_opacity=op/100.0):
                                draw_shape()
                        except AttributeError:
                            # FPDF2 version might not support local_context properly, fallback to normal
                            draw_shape()
                    else:
                        draw_shape()

        pdf.output(filepath)
        return filepath
    except ImportError:
        raise Exception("fpdf2 is not installed.")

def _generate_word(pages: list, filepath: str, canvas_width: int, canvas_height: int):
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.section import WD_ORIENT
        
        doc = Document()
        
        if canvas_width > canvas_height:
            for section in doc.sections:
                section.orientation = WD_ORIENT.LANDSCAPE
                new_width, new_height = section.page_height, section.page_width
                section.page_width = new_width
                section.page_height = new_height
        
        for idx, page in enumerate(pages):
            if idx > 0:
                doc.add_page_break()
                
            elements = page.get("elements", [])
            sorted_elements = sorted(elements, key=lambda e: e.get("y", 0))
            
            for el in sorted_elements:
                el_type = el.get("type", "text")
                if el_type == "text":
                    p = doc.add_paragraph()
                    
                    align = el.get("textAlign", "left")
                    if align == "center":
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    elif align == "right":
                        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        
                    run = p.add_run(el.get("content", ""))
                    run.font.size = Pt(el.get("fontSize", 16) * 0.75)
                    color = _hex_to_rgb(el.get("color", "#000000"))
                    from docx.shared import RGBColor
                    run.font.color.rgb = RGBColor(color[0], color[1], color[2])
                    
                elif el_type == "image":
                    src = el.get("src", "")
                    if src:
                        img_stream = _resolve_image(src)
                        if img_stream:
                            try:
                                doc.add_picture(img_stream, width=Inches(el.get("width", 200) / 96.0))
                            except Exception as e:
                                logger.warning(f"Failed to add image to Word: {e}")
                                
                elif el_type == "mermaid":
                    # Fallback to text for mermaid in Word
                    p = doc.add_paragraph()
                    run = p.add_run("[Mermaid Diagram]\n" + str(el.get("content", "")))
                    run.font.size = Pt(10)
                    run.font.name = "Courier New"
                                
        doc.save(filepath)
        return filepath
    except ImportError:
        raise Exception("python-docx is not installed.")

def _generate_pngs(pages: list, base_filepath: str, canvas_width: int, canvas_height: int, scale_factor: float = 1.0):
    from PIL import Image, ImageDraw, ImageFont
    import tempfile
    
    generated_files = []
    
    # Apply scale factor for high-DPI export
    render_w = int(canvas_width * scale_factor)
    render_h = int(canvas_height * scale_factor)
    sf = scale_factor
    
    # Try to load a generic TrueType font, otherwise fallback to default
    try:
        # Use a generic sans-serif font available on linux
        base_font = ImageFont.truetype("DejaVuSans.ttf", int(24 * sf))
        has_ttf = True
    except IOError:
        base_font = ImageFont.load_default()
        has_ttf = False

    for idx, page in enumerate(pages):
        img = Image.new('RGBA', (render_w, render_h), color='white')
        draw = ImageDraw.Draw(img, 'RGBA')
        
        elements = page.get("elements", [])
        # Sort so images are behind text if possible, or preserve z-index 
        # (elements array is already the z-index order)
        for el in elements:
            x = int(el.get("x", 0) * sf)
            y = int(el.get("y", 0) * sf)
            w = int(el.get("width", 100) * sf)
            h = int(el.get("height", 50) * sf)
            
            el_type = el.get("type", "text")
            if el_type == "text":
                content = el.get("content", "")
                color = _hex_to_rgb(el.get("color", "#000000"))
                font_size = int(el.get("fontSize", 24) * sf)
                
                if has_ttf:
                    try:
                        font = ImageFont.truetype("DejaVuSans.ttf", font_size)
                    except (IOError, OSError):
                        font = base_font
                else:
                    font = base_font
                    
                # Text background/border
                bg = el.get("backgroundColor", "transparent")
                border = el.get("borderColor", "transparent")
                bw = el.get("borderWidth", 0)
                op = el.get("opacity", 100)
                
                alpha = int((op / 100) * 255)
                
                fill_color = None
                if bg != "transparent":
                    c = _hex_to_rgb(bg)
                    fill_color = (c[0], c[1], c[2], alpha)
                
                outline_color = None
                if border != "transparent" and bw > 0:
                    c = _hex_to_rgb(border)
                    outline_color = (c[0], c[1], c[2], alpha)
                    
                if fill_color or outline_color:
                    # Estimate bbox. Pillow textbbox is better but simplified here.
                    # Since we don't have perfect text wrap in Pillow here, we just use the w,h
                    draw.rectangle([x, y, x+w, y+h], fill=fill_color, outline=outline_color, width=bw)

                align = el.get("textAlign", "left")
                text_color = (color[0], color[1], color[2], alpha)
                content = str(el.get("content", ""))
                is_bold = el.get("isBold", False)
                is_underline = el.get("isUnderline", False)
                is_strike = el.get("isStrikethrough", False)
                stroke_width = 1 if is_bold else 0
                
                import textwrap
                lines = []
                avg_char_width = font_size * 0.55
                max_chars = max(1, int(w / avg_char_width))
                
                for orig_line in content.split('\n'):
                    if el.get("isList") and orig_line.strip():
                        orig_line = f"    • {orig_line}"
                    wrapped = textwrap.wrap(orig_line, width=max_chars) if orig_line.strip() else [""]
                    lines.extend(wrapped)
                
                current_y = y + 5
                for line in lines:
                    try:
                        bbox = draw.textbbox((0, 0), line, font=font)
                        line_w = bbox[2] - bbox[0]
                        line_h = bbox[3] - bbox[1]
                    except (AttributeError, TypeError):
                        line_w = len(line) * font_size * 0.55
                        line_h = font_size
                    
                    if align == "center":
                        line_x = x + (w - line_w) / 2
                    elif align == "right":
                        line_x = x + w - line_w - 5
                    else:
                        line_x = x + 5
                        
                    draw.text((line_x, current_y), line, fill=text_color, font=font, stroke_width=stroke_width, stroke_fill=text_color)
                    
                    if is_underline:
                        draw.line([(line_x, current_y + line_h + 2), (line_x + line_w, current_y + line_h + 2)], fill=text_color, width=max(1, int(font_size/10)))
                    if is_strike:
                        draw.line([(line_x, current_y + line_h/2), (line_x + line_w, current_y + line_h/2)], fill=text_color, width=max(1, int(font_size/10)))
                        
                    current_y += line_h + (font_size * 0.2)
                
            elif el_type == "image":
                src = el.get("src", "")
                if src:
                    img_stream = _resolve_image(src)
                    if img_stream:
                        try:
                            el_img = Image.open(img_stream).convert("RGBA")
                            el_img = el_img.resize((w, h), Image.Resampling.LANCZOS)
                            
                            op = el.get("opacity", 100)
                            if op < 100:
                                alpha_channel = el_img.getchannel('A')
                                alpha_channel = alpha_channel.point(lambda i: int(i * op / 100))
                                el_img.putalpha(alpha_channel)
                                
                            # create a mask if RGBA
                            if el_img.mode == 'RGBA':
                                img.paste(el_img, (x, y), el_img)
                            else:
                                img.paste(el_img, (x, y))
                        except Exception as e:
                            logger.warning(f"Failed to add image to PNG: {e}")
            elif el_type == "mermaid":
                # Fallback to text for mermaid in PNG
                content = "[Mermaid Diagram]\n" + str(el.get("content", ""))
                try: m_font = ImageFont.truetype("DejaVuSansMono.ttf", int(14 * sf))
                except (IOError, OSError): m_font = base_font
                draw.rectangle([x, y, x+w, y+h], outline=(150, 150, 150, 255), width=2)
                draw.text((x + 10, y + 10), content, fill=(0, 0, 0, 255), font=m_font)
            elif el_type == "shape":
                bg = el.get("backgroundColor", "transparent")
                border = el.get("borderColor", "transparent")
                bw = el.get("borderWidth", 0)
                op = el.get("opacity", 100)
                alpha = int((op / 100) * 255)
                
                fill_color = None
                if bg != "transparent":
                    c = _hex_to_rgb(bg)
                    fill_color = (c[0], c[1], c[2], alpha)
                    
                outline_color = None
                if border != "transparent" and bw > 0:
                    c = _hex_to_rgb(border)
                    outline_color = (c[0], c[1], c[2], alpha)
                    
                shape_type = el.get("shapeType", "rect")
                
                def draw_shape_on(target_draw):
                    if shape_type == "rect":
                        target_draw.rectangle([x, y, x+w, y+h], fill=fill_color, outline=outline_color, width=bw)
                    elif shape_type == "circle":
                        target_draw.ellipse([x, y, x+w, y+h], fill=fill_color, outline=outline_color, width=bw)
                    elif shape_type == "triangle":
                        target_draw.polygon([(x+w/2, y), (x+w, y+h), (x, y+h)], fill=fill_color, outline=outline_color)
                    elif shape_type == "hexagon":
                        target_draw.polygon([(x+w/4, y), (x+w*3/4, y), (x+w, y+h/2), (x+w*3/4, y+h), (x+w/4, y+h), (x, y+h/2)], fill=fill_color, outline=outline_color)
                    elif shape_type == "arrow":
                        headStart = w * 0.6
                        bodyH = h * 0.4
                        bodyY1 = y + (h - bodyH) / 2
                        bodyY2 = bodyY1 + bodyH
                        target_draw.polygon([(x, bodyY1), (x+headStart, bodyY1), (x+headStart, y), (x+w, y+h/2), (x+headStart, y+h), (x+headStart, bodyY2), (x, bodyY2)], fill=fill_color, outline=outline_color)
                
                if op < 100:
                    overlay = Image.new('RGBA', img.size, (0,0,0,0))
                    overlay_draw = ImageDraw.Draw(overlay, 'RGBA')
                    draw_shape_on(overlay_draw)
                    img = Image.alpha_composite(img, overlay)
                    draw = ImageDraw.Draw(img, 'RGBA') # Rebind draw to the new composited image
                else:
                    draw_shape_on(draw)
                            
        # Convert RGBA back to RGB for PNG saving without alpha channel issues if needed, 
        # or keep RGBA. PNG supports RGBA.
        out_path = base_filepath.replace(".png", f"_slide_{idx+1}.png") if len(pages) > 1 else base_filepath
        img.save(out_path)
        generated_files.append(out_path)
        
    return generated_files

@router.post("/export")
async def export_presentation(request: Request):
    """Export a rich presentation layout to PPTX, PDF, DOCX, PNG, or CMYK PDF."""
    data = await request.json()
    format_type = data.get("format", "pptx").lower()
    pages = data.get("pages", [])
    canvas_width = data.get("canvas_width", 960)
    canvas_height = data.get("canvas_height", 540)
    
    if not pages:
        raise HTTPException(400, "Pages data is required")

    docs_dir = get_documents_dir()
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    file_id = uuid.uuid4().hex[:6]
    
    ext = "pdf" if format_type == "pdf-cmyk" else format_type
    filename = f"presentation_{timestamp}_{file_id}.{ext}"
    filepath = os.path.join(docs_dir, filename)

    try:
        if format_type == "pptx":
            _generate_powerpoint(pages, filepath, canvas_width, canvas_height)
        elif format_type == "pdf":
            _generate_pdf(pages, filepath, canvas_width, canvas_height)
        elif format_type == "pdf-cmyk":
            _generate_pdf_cmyk(pages, filepath, canvas_width, canvas_height)
        elif format_type == "docx":
            _generate_word(pages, filepath, canvas_width, canvas_height)
        elif format_type == "png":
            png_files = _generate_pngs(pages, filepath, canvas_width, canvas_height)
            if len(png_files) == 1:
                filepath = png_files[0]
                filename = os.path.basename(filepath)
                ext = "png"
            else:
                import zipfile
                zip_filename = f"presentation_{timestamp}_{file_id}.zip"
                zip_filepath = os.path.join(docs_dir, zip_filename)
                with zipfile.ZipFile(zip_filepath, 'w') as zf:
                    for f in png_files:
                        zf.write(f, os.path.basename(f))
                        os.unlink(f) # cleanup individual pngs
                filepath = zip_filepath
                filename = zip_filename
                ext = "zip"
        else:
            raise HTTPException(400, f"Unsupported format: {format_type}")

        return {"status": "ok", "filename": filename, "format": ext, "path": filepath, "url": f"/data/documents/{filename}"}
    except Exception as e:
        logger.error(f"Error generating presentation: {e}")
        raise HTTPException(500, str(e))

@router.post("/import")
async def import_presentation(request: Request):
    """Import a PDF or PPTX file and convert it to presentation pages."""
    from fastapi import UploadFile, File
    import tempfile

    form = await request.form()
    uploaded_file = form.get("file")
    if not uploaded_file:
        raise HTTPException(400, "No file uploaded")

    filename = uploaded_file.filename or "unknown"
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ("pdf", "pptx"):
        raise HTTPException(400, f"Unsupported format: {ext}. Only PDF and PPTX are supported.")

    content = await uploaded_file.read()
    title = filename.rsplit(".", 1)[0] if "." in filename else filename

    try:
        if ext == "pdf":
            pages, cw, ch = _import_pdf(content)
        else:
            pages, cw, ch = _import_pptx(content)

        return {
            "status": "ok",
            "title": title,
            "pages": pages,
            "canvas_width": cw,
            "canvas_height": ch,
        }
    except ImportError as e:
        logger.error(f"Missing dependency for import: {e}")
        raise HTTPException(500, f"Missing dependency: {e}")
    except Exception as e:
        logger.error(f"Error importing presentation: {e}")
        raise HTTPException(500, str(e))


def _import_pdf(content: bytes):
    """Convert PDF pages to image-based presentation slides."""
    import tempfile

    images_dir = os.path.join(DATA_DIR, "images")
    os.makedirs(images_dir, exist_ok=True)

    canvas_w, canvas_h = 960, 540

    # Try pdf2image (poppler) first, fallback to PyMuPDF (fitz)
    pages = []
    try:
        from pdf2image import convert_from_bytes
        pil_images = convert_from_bytes(content, dpi=150)
        for idx, img in enumerate(pil_images):
            img_id = f"import_{uuid.uuid4().hex[:8]}"
            img_filename = f"{img_id}.png"
            img_path = os.path.join(images_dir, img_filename)
            img.save(img_path, "PNG")

            # Scale to canvas
            iw, ih = img.size
            scale = min(canvas_w / iw, canvas_h / ih)
            el_w = int(iw * scale)
            el_h = int(ih * scale)
            el_x = (canvas_w - el_w) // 2
            el_y = (canvas_h - el_h) // 2

            pages.append({
                "elements": [{
                    "id": f"el_{uuid.uuid4().hex[:8]}",
                    "type": "image",
                    "src": f"/data/images/{img_filename}",
                    "x": el_x, "y": el_y,
                    "width": el_w, "height": el_h,
                    "opacity": 100,
                }]
            })
    except ImportError:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            for idx in range(len(doc)):
                page = doc[idx]
                mat = fitz.Matrix(2, 2)  # 2x zoom for decent quality
                pix = page.get_pixmap(matrix=mat)
                img_id = f"import_{uuid.uuid4().hex[:8]}"
                img_filename = f"{img_id}.png"
                img_path = os.path.join(images_dir, img_filename)
                pix.save(img_path)

                iw, ih = pix.width, pix.height
                scale = min(canvas_w / iw, canvas_h / ih)
                el_w = int(iw * scale)
                el_h = int(ih * scale)
                el_x = (canvas_w - el_w) // 2
                el_y = (canvas_h - el_h) // 2

                pages.append({
                    "elements": [{
                        "id": f"el_{uuid.uuid4().hex[:8]}",
                        "type": "image",
                        "src": f"/data/images/{img_filename}",
                        "x": el_x, "y": el_y,
                        "width": el_w, "height": el_h,
                        "opacity": 100,
                    }]
                })
            doc.close()
        except ImportError:
            raise ImportError(
                "PDF import requires either 'pdf2image' (with poppler) or 'PyMuPDF' (fitz). "
                "Install with: pip install pdf2image  or  pip install PyMuPDF"
            )

    if not pages:
        pages = [{"elements": []}]

    return pages, canvas_w, canvas_h


def _import_pptx(content: bytes):
    """Convert PPTX slides to presentation pages with editable elements."""
    from pptx import Presentation as PptxPresentation
    from pptx.util import Emu

    prs = PptxPresentation(BytesIO(content))
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    canvas_w = 960
    canvas_h = int(canvas_w * slide_h_emu / slide_w_emu) if slide_w_emu else 540

    scale_x = canvas_w / slide_w_emu if slide_w_emu else 1
    scale_y = canvas_h / slide_h_emu if slide_h_emu else 1

    images_dir = os.path.join(DATA_DIR, "images")
    os.makedirs(images_dir, exist_ok=True)

    pages = []
    for slide in prs.slides:
        elements = []
        for shape in slide.shapes:
            left = int((shape.left or 0) * scale_x)
            top = int((shape.top or 0) * scale_y)
            width = int((shape.width or 100) * scale_x)
            height = int((shape.height or 50) * scale_y)

            el_id = f"el_{uuid.uuid4().hex[:8]}"

            if shape.has_text_frame:
                lines = []
                font_size = 16
                font_bold = False
                font_italic = False
                font_color = "#000000"
                text_align = "left"

                for para in shape.text_frame.paragraphs:
                    line_text = para.text
                    if not line_text.strip():
                        lines.append("")
                        continue
                    lines.append(line_text)

                    # Extract formatting from first run
                    if para.runs:
                        run = para.runs[0]
                        if run.font.size:
                            font_size = max(8, int(run.font.size / Emu(12700)))
                        if run.font.bold:
                            font_bold = True
                        if run.font.italic:
                            font_italic = True
                        if run.font.color and run.font.color.rgb:
                            font_color = f"#{run.font.color.rgb}"

                    # Alignment
                    from pptx.enum.text import PP_ALIGN
                    if para.alignment == PP_ALIGN.CENTER:
                        text_align = "center"
                    elif para.alignment == PP_ALIGN.RIGHT:
                        text_align = "right"

                content = "\n".join(lines)
                if content.strip():
                    elements.append({
                        "id": el_id,
                        "type": "text",
                        "content": content,
                        "x": left, "y": top,
                        "width": width, "height": height,
                        "fontSize": font_size,
                        "color": font_color,
                        "backgroundColor": "transparent",
                        "isBold": font_bold,
                        "isItalic": font_italic,
                        "textAlign": text_align,
                        "opacity": 100,
                    })

            elif shape.shape_type and hasattr(shape, "image"):
                try:
                    img_blob = shape.image.blob
                    img_ext = shape.image.content_type.split("/")[-1] if shape.image.content_type else "png"
                    if img_ext == "jpeg":
                        img_ext = "jpg"
                    img_id = f"import_{uuid.uuid4().hex[:8]}"
                    img_filename = f"{img_id}.{img_ext}"
                    img_path = os.path.join(images_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(img_blob)

                    elements.append({
                        "id": el_id,
                        "type": "image",
                        "src": f"/data/images/{img_filename}",
                        "x": left, "y": top,
                        "width": width, "height": height,
                        "opacity": 100,
                    })
                except Exception as e:
                    logger.warning(f"Failed to extract image from PPTX shape: {e}")

            elif hasattr(shape, "shape_type"):
                # Generic shape (rectangle, etc.)
                bg_color = "#cccccc"
                try:
                    if shape.fill and shape.fill.fore_color:
                        bg_color = f"#{shape.fill.fore_color.rgb}"
                except Exception:
                    pass

                elements.append({
                    "id": el_id,
                    "type": "shape",
                    "shapeType": "rect",
                    "x": left, "y": top,
                    "width": width, "height": height,
                    "backgroundColor": bg_color,
                    "borderColor": "transparent",
                    "borderWidth": 0,
                    "opacity": 100,
                })

        pages.append({"elements": elements})

    if not pages:
        pages = [{"elements": []}]

    return pages, canvas_w, canvas_h

@router.post("/save")
async def save_presentation(request: Request):
    data = await request.json()
    pres_id = data.get("id") or uuid.uuid4().hex[:8]
    data["id"] = pres_id
    
    pages = data.get("pages", [])
    if pages:
        first_page = [pages[0]]
        cw = data.get("canvas_width", 960)
        ch = data.get("canvas_height", 540)
        docs_dir = get_documents_dir()
        thumb_name = f"thumb_{pres_id}.png"
        thumb_path = os.path.join(docs_dir, thumb_name)
        try:
            _generate_pngs(first_page, thumb_path, cw, ch)
            data["thumbnail"] = f"/data/documents/{thumb_name}"
        except Exception as e:
            logger.warning(f"Failed to generate thumbnail: {e}")
            
    pres_dir = get_presentations_dir()
    filepath = os.path.join(pres_dir, f"{pres_id}.json")
    
    import json
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)
        
    return {"status": "ok", "id": pres_id}

@router.get("/list")
async def list_presentations():
    pres_dir = get_presentations_dir()
    import glob, json
    files = glob.glob(os.path.join(pres_dir, "*.json"))
    results = []
    for f in files:
        pres_id = os.path.basename(f).replace(".json", "")
        # read basic info
        try:
            with open(f, "r", encoding="utf-8") as file:
                d = json.load(file)
                results.append({
                    "id": pres_id,
                    "title": d.get("title", f"Presentation {pres_id}"),
                    "thumbnail": d.get("thumbnail"),
                    "updated_at": os.path.getmtime(f)
                })
        except (json.JSONDecodeError, OSError) as e:
            logger.warning("Failed to read presentation %s: %s", pres_id, e)
    results.sort(key=lambda x: x["updated_at"], reverse=True)
    return {"presentations": results}

@router.get("/load/{pres_id}")
async def load_presentation(pres_id: str):
    pres_dir = get_presentations_dir()
    # basic security
    safe_id = "".join([c for c in pres_id if c.isalnum() or c == '-'])
    filepath = os.path.join(pres_dir, f"{safe_id}.json")
    if not os.path.exists(filepath):
        raise HTTPException(404, "Not found")
    
    import json
    with open(filepath, "r", encoding="utf-8") as f:
        return json.load(f)

@router.delete("/delete/{pres_id}")
async def delete_presentation(pres_id: str):
    pres_dir = get_presentations_dir()
    safe_id = "".join([c for c in pres_id if c.isalnum() or c == '-'])
    filepath = os.path.join(pres_dir, f"{safe_id}.json")
    if os.path.exists(filepath):
        try:
            import json
            with open(filepath, "r", encoding="utf-8") as f:
                d = json.load(f)
                thumb = d.get("thumbnail")
                if thumb and thumb.startswith("/data/documents/"):
                    thumb_filename = thumb.split("?")[0].split("/")[-1]
                    thumb_path = os.path.join(get_documents_dir(), thumb_filename)
                    if os.path.exists(thumb_path):
                        os.unlink(thumb_path)
            os.unlink(filepath)
            return {"status": "ok"}
        except Exception as e:
            logger.error(f"Error deleting presentation: {e}")
            raise HTTPException(500, str(e))
    raise HTTPException(404, "Not found")


# ---------------------------------------------------------------------------
# CMYK PDF Export (High Quality, 300 DPI)
# ---------------------------------------------------------------------------
def _generate_pdf_cmyk(pages: list, filepath: str, canvas_width: int, canvas_height: int):
    """Generate a high-quality CMYK PDF by rendering pages as 300 DPI images and converting to CMYK."""
    import tempfile
    from PIL import Image
    
    # First generate high-res PNGs
    dpi = 300
    scale = dpi / 96  # standard screen is ~96 DPI
    hi_res_w = int(canvas_width * scale)
    hi_res_h = int(canvas_height * scale)
    
    # Generate PNG pages at high resolution
    tmp_pngs = []
    for i, page in enumerate(pages):
        tmp_path = os.path.join(tempfile.gettempdir(), f"cmyk_page_{i}.png")
        _generate_pngs([page], tmp_path, canvas_width, canvas_height, scale_factor=scale)
        if os.path.exists(tmp_path):
            tmp_pngs.append(tmp_path)
    
    # Convert each PNG to CMYK and embed in PDF
    try:
        cmyk_images = []
        for png_path in tmp_pngs:
            img = Image.open(png_path).convert("RGB")
            # Simple mathematical RGB → CMYK conversion
            cmyk_img = _rgb_to_cmyk(img)
            cmyk_images.append(cmyk_img)
        
        # Save as multi-page PDF
        if cmyk_images:
            first = cmyk_images[0]
            rest = cmyk_images[1:]
            first.save(filepath, "PDF", resolution=dpi, save_all=True, append_images=rest)
            logger.info(f"CMYK PDF generated: {filepath} ({len(cmyk_images)} pages, {dpi} DPI)")
        
    except Exception as e:
        logger.warning(f"Failed to create CMYK PDF: {e}. Falling back to standard high-DPI PDF")
        # Fallback: standard high-resolution PDF without CMYK
        images_rgb = []
        for png_path in tmp_pngs:
            img = Image.open(png_path).convert("RGB")
            images_rgb.append(img)
        if images_rgb:
            images_rgb[0].save(filepath, "PDF", resolution=dpi, save_all=True, 
                             append_images=images_rgb[1:])
    finally:
        # Cleanup temp files
        for p in tmp_pngs:
            try: os.unlink(p)
            except OSError: pass


def _rgb_to_cmyk(img):
    """Convert an RGB PIL Image to CMYK using mathematical conversion."""
    from PIL import Image
    import numpy as np
    
    rgb = np.array(img, dtype=np.float64) / 255.0
    r, g, b = rgb[:,:,0], rgb[:,:,1], rgb[:,:,2]
    
    k = 1 - np.maximum(np.maximum(r, g), b)
    k_inv = np.where(k < 1, 1 - k, 1)  # avoid division by zero
    
    c = (1 - r - k) / k_inv
    m = (1 - g - k) / k_inv
    y_ch = (1 - b - k) / k_inv
    
    c = np.clip(c, 0, 1)
    m = np.clip(m, 0, 1)
    y_ch = np.clip(y_ch, 0, 1)
    k = np.clip(k, 0, 1)
    
    cmyk = np.stack([c, m, y_ch, k], axis=-1)
    cmyk_uint8 = (cmyk * 255).astype(np.uint8)
    
    return Image.fromarray(cmyk_uint8, mode="CMYK")


# ---------------------------------------------------------------------------
# AI Presentation Generation (2-phase: LLM content → Python layout)
# ---------------------------------------------------------------------------

# Pre-defined color themes
_THEMES = {
    "midnight": {
        "bg": "#0f0f1a", "surface": "#1a1a2e", "accent": "#7c5cfc",
        "accent2": "#f472b6", "text": "#ffffff", "muted": "#a0a0b8",
        "shape1": "#7c5cfc", "shape2": "#f472b6", "shape3": "#34d399",
    },
    "ocean": {
        "bg": "#0a192f", "surface": "#112240", "accent": "#38bdf8",
        "accent2": "#818cf8", "text": "#e2e8f0", "muted": "#8892b0",
        "shape1": "#38bdf8", "shape2": "#818cf8", "shape3": "#34d399",
    },
    "sunset": {
        "bg": "#1a0a2e", "surface": "#2d1b4e", "accent": "#f59e0b",
        "accent2": "#ef4444", "text": "#fef3c7", "muted": "#d4a574",
        "shape1": "#f59e0b", "shape2": "#ef4444", "shape3": "#f472b6",
    },
    "forest": {
        "bg": "#0a1a0f", "surface": "#132a1a", "accent": "#34d399",
        "accent2": "#a3e635", "text": "#ecfdf5", "muted": "#86efac",
        "shape1": "#34d399", "shape2": "#a3e635", "shape3": "#38bdf8",
    },
    "corporate": {
        "bg": "#111827", "surface": "#1f2937", "accent": "#6366f1",
        "accent2": "#06b6d4", "text": "#f9fafb", "muted": "#9ca3af",
        "shape1": "#6366f1", "shape2": "#06b6d4", "shape3": "#f472b6",
    },
    "minimal_light": {
        "bg": "#ffffff", "surface": "#f8fafc", "accent": "#2563eb",
        "accent2": "#059669", "text": "#0f172a", "muted": "#64748b",
        "shape1": "#2563eb", "shape2": "#059669", "shape3": "#8b5cf6",
    },
    "pastel_blue": {
        "bg": "#f0f9ff", "surface": "#e0f2fe", "accent": "#0ea5e9",
        "accent2": "#8b5cf6", "text": "#0f172a", "muted": "#475569",
        "shape1": "#0ea5e9", "shape2": "#8b5cf6", "shape3": "#f43f5e",
    },
    "warm_light": {
        "bg": "#fffbeb", "surface": "#fef3c7", "accent": "#d97706",
        "accent2": "#dc2626", "text": "#451a03", "muted": "#78350f",
        "shape1": "#d97706", "shape2": "#dc2626", "shape3": "#059669",
    },
    "cyberpunk": {
        "bg": "#120422", "surface": "#230841", "accent": "#00ffcc",
        "accent2": "#ff00ff", "text": "#ffffff", "muted": "#a0aab8",
        "shape1": "#00ffcc", "shape2": "#ff00ff", "shape3": "#fcd34d",
    },
    "nordic": {
        "bg": "#eceff4", "surface": "#e5e9f0", "accent": "#5e81ac",
        "accent2": "#81a1c1", "text": "#2e3440", "muted": "#4c566a",
        "shape1": "#5e81ac", "shape2": "#88c0d0", "shape3": "#b48ead",
    },
    "dracula": {
        "bg": "#282a36", "surface": "#44475a", "accent": "#ff79c6",
        "accent2": "#bd93f9", "text": "#f8f8f2", "muted": "#6272a4",
        "shape1": "#ff79c6", "shape2": "#bd93f9", "shape3": "#8be9fd",
    },
    "monochrome": {
        "bg": "#ffffff", "surface": "#f5f5f5", "accent": "#000000",
        "accent2": "#404040", "text": "#171717", "muted": "#737373",
        "shape1": "#000000", "shape2": "#525252", "shape3": "#a3a3a3",
    },
    "solarized_dark": {
        "bg": "#002b36", "surface": "#073642", "accent": "#b58900",
        "accent2": "#2aa198", "text": "#839496", "muted": "#586e75",
        "shape1": "#cb4b16", "shape2": "#268bd2", "shape3": "#859900",
    },
    "ruby": {
        "bg": "#1a0505", "surface": "#330a0a", "accent": "#ff4d4d",
        "accent2": "#ff9999", "text": "#ffffff", "muted": "#cc8888",
        "shape1": "#ff4d4d", "shape2": "#ff9999", "shape3": "#ffb3b3",
    },
    "emerald": {
        "bg": "#f0fdf4", "surface": "#dcfce7", "accent": "#16a34a",
        "accent2": "#22c55e", "text": "#14532d", "muted": "#166534",
        "shape1": "#16a34a", "shape2": "#86efac", "shape3": "#34d399",
    },
    "sapphire": {
        "bg": "#020617", "surface": "#0f172a", "accent": "#2563eb",
        "accent2": "#3b82f6", "text": "#f8fafc", "muted": "#94a3b8",
        "shape1": "#2563eb", "shape2": "#3b82f6", "shape3": "#60a5fa",
    },
    "golden": {
        "bg": "#171717", "surface": "#262626", "accent": "#eab308",
        "accent2": "#facc15", "text": "#fafafa", "muted": "#a3a3a3",
        "shape1": "#eab308", "shape2": "#facc15", "shape3": "#fef08a",
    },
    "lavender": {
        "bg": "#faf5ff", "surface": "#f3e8ff", "accent": "#9333ea",
        "accent2": "#a855f7", "text": "#3b0764", "muted": "#6b21a8",
        "shape1": "#9333ea", "shape2": "#c084fc", "shape3": "#d8b4fe",
    },
    "desert": {
        "bg": "#fff7ed", "surface": "#ffedd5", "accent": "#ea580c",
        "accent2": "#f97316", "text": "#431407", "muted": "#9a3412",
        "shape1": "#ea580c", "shape2": "#fb923c", "shape3": "#fdba74",
    },
    "hacker": {
        "bg": "#0d1117", "surface": "#161b22", "accent": "#00ff00",
        "accent2": "#33ff33", "text": "#c9d1d9", "muted": "#8b949e",
        "shape1": "#00ff00", "shape2": "#2ea043", "shape3": "#00ff00",
    },
}

# ---------------------------------------------------------------------------
# Open-Source SVG Illustration Catalog
# ---------------------------------------------------------------------------
_SVG_ILLUSTRATIONS = {
    "rocket": '<path d="M100 20 C100 20 70 60 70 110 C70 140 85 160 100 170 C115 160 130 140 130 110 C130 60 100 20 100 20Z" fill="{c1}" opacity="0.9"/><circle cx="100" cy="100" r="15" fill="{c2}"/><path d="M75 140 L55 170 L80 155Z" fill="{c1}" opacity="0.7"/><path d="M125 140 L145 170 L120 155Z" fill="{c1}" opacity="0.7"/><path d="M90 170 L100 195 L110 170" fill="{c2}" opacity="0.6"/>',
    "growth": '<rect x="30" y="140" width="25" height="40" rx="4" fill="{c1}" opacity="0.5"/><rect x="65" y="110" width="25" height="70" rx="4" fill="{c1}" opacity="0.65"/><rect x="100" y="80" width="25" height="100" rx="4" fill="{c1}" opacity="0.8"/><rect x="135" y="45" width="25" height="135" rx="4" fill="{c2}"/><path d="M35 130 L80 95 L115 70 L150 35" stroke="{c2}" stroke-width="3" fill="none" stroke-linecap="round"/><circle cx="150" cy="35" r="6" fill="{c2}"/>',
    "target": '<circle cx="100" cy="100" r="70" fill="none" stroke="{c1}" stroke-width="4" opacity="0.3"/><circle cx="100" cy="100" r="50" fill="none" stroke="{c1}" stroke-width="4" opacity="0.5"/><circle cx="100" cy="100" r="30" fill="none" stroke="{c2}" stroke-width="4" opacity="0.8"/><circle cx="100" cy="100" r="10" fill="{c2}"/><path d="M100 20 L100 40" stroke="{c1}" stroke-width="2"/><path d="M100 160 L100 180" stroke="{c1}" stroke-width="2"/>',
    "lightbulb": '<path d="M100 30 C65 30 45 55 45 85 C45 110 60 125 70 140 L130 140 C140 125 155 110 155 85 C155 55 135 30 100 30Z" fill="{c1}" opacity="0.85"/><rect x="75" y="145" width="50" height="8" rx="4" fill="{c2}"/><rect x="80" y="158" width="40" height="8" rx="4" fill="{c2}" opacity="0.7"/><rect x="85" y="171" width="30" height="8" rx="4" fill="{c2}" opacity="0.5"/><path d="M85 85 L100 65 L115 85" stroke="{c2}" stroke-width="3" fill="none"/>',
    "shield": '<path d="M100 25 L40 55 L40 105 C40 140 65 170 100 185 C135 170 160 140 160 105 L160 55Z" fill="{c1}" opacity="0.85"/><path d="M90 110 L100 120 L120 85" stroke="{c2}" stroke-width="5" fill="none" stroke-linecap="round" stroke-linejoin="round"/>',
    "cloud": '<circle cx="70" cy="110" r="35" fill="{c1}" opacity="0.8"/><circle cx="110" cy="95" r="40" fill="{c1}" opacity="0.85"/><circle cx="145" cy="110" r="30" fill="{c1}" opacity="0.75"/><rect x="45" y="110" width="130" height="40" rx="10" fill="{c1}" opacity="0.9"/><path d="M95 130 L95 165 M80 150 L95 165 L110 150" stroke="{c2}" stroke-width="3" fill="none" stroke-linecap="round"/>',
    "globe": '<circle cx="100" cy="100" r="65" fill="none" stroke="{c1}" stroke-width="3" opacity="0.8"/><ellipse cx="100" cy="100" rx="30" ry="65" fill="none" stroke="{c1}" stroke-width="2" opacity="0.5"/><path d="M35 100 L165 100" stroke="{c1}" stroke-width="2" opacity="0.4"/><path d="M45 70 L155 70" stroke="{c1}" stroke-width="1.5" opacity="0.3"/><path d="M45 130 L155 130" stroke="{c1}" stroke-width="1.5" opacity="0.3"/><circle cx="130" cy="60" r="8" fill="{c2}" opacity="0.9"/>',
    "team": '<circle cx="100" cy="55" r="20" fill="{c1}"/><circle cx="55" cy="75" r="16" fill="{c1}" opacity="0.6"/><circle cx="145" cy="75" r="16" fill="{c1}" opacity="0.6"/><path d="M65 100 C65 85 135 85 135 100 L135 140 C135 145 65 145 65 140Z" fill="{c2}" opacity="0.8"/><path d="M30 115 C30 100 75 105 75 115 L75 150 C75 155 30 155 30 150Z" fill="{c2}" opacity="0.5"/><path d="M125 115 C125 100 170 105 170 115 L170 150 C170 155 125 155 125 150Z" fill="{c2}" opacity="0.5"/>',
    "brain": '<path d="M100 160 L100 100" stroke="{c1}" stroke-width="3"/><path d="M70 50 C45 50 35 75 50 90 C35 100 45 125 65 120 C65 140 90 145 100 130 C110 145 135 140 135 120 C155 125 165 100 150 90 C165 75 155 50 130 50 C125 35 110 30 100 40 C90 30 75 35 70 50Z" fill="{c1}" opacity="0.8"/><path d="M100 55 L100 130 M80 70 L100 85 M120 70 L100 85 M80 105 L100 95 M120 105 L100 95" stroke="{c2}" stroke-width="2" fill="none" opacity="0.6"/>',
    "chart-up": '<path d="M30 170 L30 30" stroke="{c1}" stroke-width="3"/><path d="M30 170 L170 170" stroke="{c1}" stroke-width="3"/><path d="M50 140 L80 100 L110 120 L140 60 L160 45" stroke="{c2}" stroke-width="3" fill="none" stroke-linecap="round"/><circle cx="160" cy="45" r="5" fill="{c2}"/><path d="M145 45 L160 45 L160 60" stroke="{c2}" stroke-width="2" fill="none"/>',
    "handshake": '<path d="M40 90 L70 70 L95 85 L120 65 L160 90" stroke="{c1}" stroke-width="4" fill="none" stroke-linecap="round"/><path d="M70 70 L50 110 L80 130 L110 110" fill="{c1}" opacity="0.6"/><path d="M120 65 L150 110 L120 130 L90 110" fill="{c2}" opacity="0.6"/><circle cx="100" cy="105" r="8" fill="{c2}"/>',
    "puzzle": '<path d="M40 40 L90 40 C90 30 110 30 110 40 L160 40 L160 90 C170 90 170 110 160 110 L160 160 L110 160 C110 170 90 170 90 160 L40 160 L40 110 C30 110 30 90 40 90Z" fill="{c1}" opacity="0.8"/><path d="M100 40 L100 160 M40 100 L160 100" stroke="{c2}" stroke-width="2" opacity="0.4"/>',
    "database": '<ellipse cx="100" cy="50" rx="60" ry="20" fill="{c1}" opacity="0.9"/><path d="M40 50 L40 150" stroke="{c1}" stroke-width="2"/><path d="M160 50 L160 150" stroke="{c1}" stroke-width="2"/><ellipse cx="100" cy="150" rx="60" ry="20" fill="{c1}" opacity="0.7"/><ellipse cx="100" cy="100" rx="60" ry="20" fill="none" stroke="{c1}" stroke-width="2" opacity="0.5"/><circle cx="100" cy="100" r="8" fill="{c2}"/>',
    "gear": '<circle cx="100" cy="100" r="25" fill="none" stroke="{c1}" stroke-width="4"/><circle cx="100" cy="100" r="10" fill="{c2}"/><path d="M100 35 L95 55 L105 55Z M100 165 L105 145 L95 145Z M35 100 L55 95 L55 105Z M165 100 L145 105 L145 95Z M55 55 L68 70 L75 63Z M145 145 L132 130 L125 137Z M145 55 L130 68 L137 75Z M55 145 L70 132 L63 125Z" fill="{c1}" opacity="0.8"/>',
    "trophy": '<path d="M75 50 L125 50 L120 110 L80 110Z" fill="{c1}" opacity="0.85"/><path d="M65 50 C30 50 30 90 65 90" stroke="{c1}" stroke-width="3" fill="none" opacity="0.6"/><path d="M135 50 C170 50 170 90 135 90" stroke="{c1}" stroke-width="3" fill="none" opacity="0.6"/><rect x="90" y="110" width="20" height="25" fill="{c2}" opacity="0.7"/><rect x="70" y="135" width="60" height="15" rx="4" fill="{c2}"/><path d="M95 65 L100 55 L105 65 L100 75Z" fill="{c2}"/>',
    "money": '<circle cx="100" cy="100" r="60" fill="{c1}" opacity="0.2"/><circle cx="100" cy="100" r="45" fill="{c1}" opacity="0.4"/><text x="100" y="115" text-anchor="middle" font-size="50" font-weight="bold" fill="{c2}">$</text>',
    "lock": '<rect x="65" y="95" width="70" height="60" rx="8" fill="{c1}" opacity="0.85"/><path d="M80 95 L80 70 C80 48 120 48 120 70 L120 95" stroke="{c1}" stroke-width="5" fill="none"/><circle cx="100" cy="120" r="8" fill="{c2}"/><path d="M100 128 L100 140" stroke="{c2}" stroke-width="4"/>',
    "book": '<path d="M100 35 L100 165" stroke="{c1}" stroke-width="3"/><path d="M100 35 C80 30 40 35 35 45 L35 155 C40 148 80 145 100 155" fill="{c1}" opacity="0.7"/><path d="M100 35 C120 30 160 35 165 45 L165 155 C160 148 120 145 100 155" fill="{c2}" opacity="0.5"/><path d="M55 70 L85 70 M55 90 L85 90 M55 110 L80 110" stroke="{c2}" stroke-width="2" opacity="0.5"/>',
    "megaphone": '<path d="M50 80 L50 120 L80 120 L140 155 L140 45 L80 80Z" fill="{c1}" opacity="0.85"/><rect x="30" y="85" width="25" height="30" rx="6" fill="{c2}" opacity="0.7"/><path d="M150 75 L170 65 M150 100 L175 100 M150 125 L170 135" stroke="{c2}" stroke-width="3" stroke-linecap="round"/>',
    "compass": '<circle cx="100" cy="100" r="60" fill="none" stroke="{c1}" stroke-width="3" opacity="0.6"/><circle cx="100" cy="100" r="50" fill="none" stroke="{c1}" stroke-width="1" opacity="0.3"/><polygon points="100,50 108,95 100,85 92,95" fill="{c2}"/><polygon points="100,150 92,105 100,115 108,105" fill="{c1}" opacity="0.6"/><circle cx="100" cy="100" r="5" fill="{c2}"/>',
    "atom": '<circle cx="100" cy="100" r="10" fill="{c2}"/><ellipse cx="100" cy="100" rx="60" ry="20" fill="none" stroke="{c1}" stroke-width="2" opacity="0.7"/><ellipse cx="100" cy="100" rx="60" ry="20" fill="none" stroke="{c1}" stroke-width="2" opacity="0.7" transform="rotate(60 100 100)"/><ellipse cx="100" cy="100" rx="60" ry="20" fill="none" stroke="{c1}" stroke-width="2" opacity="0.7" transform="rotate(-60 100 100)"/>',
    "code": '<rect x="30" y="40" width="140" height="120" rx="8" fill="{c1}" opacity="0.2"/><path d="M70 80 L50 100 L70 120" stroke="{c2}" stroke-width="4" fill="none" stroke-linecap="round" stroke-linejoin="round"/><path d="M130 80 L150 100 L130 120" stroke="{c2}" stroke-width="4" fill="none" stroke-linecap="round" stroke-linejoin="round"/><path d="M110 70 L90 130" stroke="{c1}" stroke-width="3" opacity="0.6"/>',
    "calendar": '<rect x="40" y="50" width="120" height="110" rx="8" fill="{c1}" opacity="0.2"/><rect x="40" y="50" width="120" height="30" rx="8" fill="{c1}" opacity="0.8"/><rect x="55" y="40" width="8" height="20" rx="3" fill="{c2}"/><rect x="137" y="40" width="8" height="20" rx="3" fill="{c2}"/><rect x="55" y="95" width="18" height="18" rx="3" fill="{c2}" opacity="0.5"/><rect x="85" y="95" width="18" height="18" rx="3" fill="{c2}"/><rect x="115" y="95" width="18" height="18" rx="3" fill="{c2}" opacity="0.5"/><rect x="55" y="125" width="18" height="18" rx="3" fill="{c2}" opacity="0.3"/>',
    "network": '<circle cx="100" cy="60" r="12" fill="{c2}"/><circle cx="55" cy="130" r="12" fill="{c1}" opacity="0.8"/><circle cx="145" cy="130" r="12" fill="{c1}" opacity="0.8"/><circle cx="35" cy="85" r="8" fill="{c1}" opacity="0.5"/><circle cx="165" cy="85" r="8" fill="{c1}" opacity="0.5"/><path d="M100 72 L55 118 M100 72 L145 118 M55 130 L145 130 M43 85 L88 60 M157 85 L112 60" stroke="{c1}" stroke-width="2" opacity="0.5"/>',
    "innovation": '<circle cx="100" cy="70" r="35" fill="{c1}" opacity="0.3"/><path d="M85 55 L100 35 L115 55" stroke="{c2}" stroke-width="3" fill="none"/><path d="M100 35 L100 100" stroke="{c2}" stroke-width="2"/><path d="M70 120 L100 100 L130 120" stroke="{c1}" stroke-width="2" fill="none"/><circle cx="70" cy="130" r="8" fill="{c2}" opacity="0.6"/><circle cx="100" cy="140" r="8" fill="{c2}" opacity="0.8"/><circle cx="130" cy="130" r="8" fill="{c2}" opacity="0.6"/>',
    "analytics": '<rect x="30" y="40" width="140" height="120" rx="8" fill="{c1}" opacity="0.15"/><circle cx="85" cy="100" r="35" fill="none" stroke="{c1}" stroke-width="3" opacity="0.7"/><path d="M85 65 L85 100 L110 100" stroke="{c2}" stroke-width="3" fill="none"/><path d="M85 65 A35 35 0 0 1 120 100" fill="{c2}" opacity="0.3"/><rect x="130" y="55" width="30" height="6" rx="3" fill="{c1}" opacity="0.5"/><rect x="130" y="70" width="25" height="6" rx="3" fill="{c2}" opacity="0.6"/><rect x="130" y="85" width="20" height="6" rx="3" fill="{c1}" opacity="0.4"/>',
    "flag": '<path d="M50 30 L50 175" stroke="{c1}" stroke-width="4" stroke-linecap="round"/><path d="M50 35 L150 55 L150 100 L50 80Z" fill="{c2}" opacity="0.8"/><path d="M50 35 L150 55 L150 100 L50 80Z" fill="none" stroke="{c1}" stroke-width="2" opacity="0.3"/>',
    "microscope": '<circle cx="85" cy="55" r="25" fill="none" stroke="{c1}" stroke-width="3" opacity="0.7"/><path d="M105 72 L130 100" stroke="{c1}" stroke-width="4" stroke-linecap="round"/><rect x="75" y="100" width="50" height="8" rx="4" fill="{c2}"/><rect x="95" y="108" width="10" height="40" fill="{c1}" opacity="0.8"/><rect x="70" y="148" width="60" height="12" rx="6" fill="{c2}"/><circle cx="85" cy="55" r="8" fill="{c2}" opacity="0.3"/>',
    "chat": '<path d="M40 45 L160 45 C165 45 170 50 170 55 L170 115 C170 120 165 125 160 125 L110 125 L80 155 L85 125 L45 125 C40 125 35 120 35 115 L35 55 C35 50 40 45 40 45Z" fill="{c1}" opacity="0.8"/><rect x="60" y="70" width="40" height="6" rx="3" fill="{c2}" opacity="0.6"/><rect x="60" y="85" width="60" height="6" rx="3" fill="{c2}" opacity="0.4"/><rect x="60" y="100" width="30" height="6" rx="3" fill="{c2}" opacity="0.3"/>',
}

# Categories for the SVG catalog (for frontend browsing)
_SVG_CATEGORIES = {
    "business": ["growth", "target", "handshake", "trophy", "money", "flag"],
    "technology": ["rocket", "code", "gear", "database", "atom", "network"],
    "data": ["chart-up", "analytics", "calendar"],
    "communication": ["team", "megaphone", "chat", "globe"],
    "education": ["lightbulb", "brain", "book", "microscope"],
    "security": ["shield", "lock"],
    "other": ["cloud", "puzzle", "compass", "innovation"],
}

def _get_svg_illustration(name: str, color1: str, color2: str, w: int = 200, h: int = 200) -> str:
    """Generate an inline SVG illustration using theme colors."""
    template = _SVG_ILLUSTRATIONS.get(name)
    if not template:
        return ""
    svg_body = template.replace("{c1}", color1).replace("{c2}", color2)
    return f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 200" width="{w}" height="{h}">{svg_body}</svg>'

def _get_svg_data_uri(name: str, color1: str, color2: str, w: int = 200, h: int = 200) -> str:
    """Generate a data URI for an SVG illustration."""
    svg = _get_svg_illustration(name, color1, color2, w, h)
    if not svg:
        return ""
    import urllib.parse
    return "data:image/svg+xml," + urllib.parse.quote(svg, safe='')


def _build_slide_layout(slide_data: dict, cw: int, ch: int, theme: dict, font: str, header_text: str = "", footer_text: str = "", page_number_position: str = "bottom-right") -> dict:
    """Build a pixel-perfect slide from structured content using predefined layouts."""
    slide_type = slide_data.get("type", "content")
    elements = []
    pad = int(cw * 0.06)  # 6% padding
    
    # --- Full-slide background rectangle ---
    elements.append({
        "type": "shape", "shapeType": "rect",
        "x": 0, "y": 0, "width": cw, "height": ch,
        "backgroundColor": theme["bg"], "borderWidth": 0, "borderColor": "transparent",
    })
    
    if slide_type == "title":
        # Decorative accent bar at top
        elements.append({
            "type": "shape", "shapeType": "rect",
            "x": 0, "y": 0, "width": cw, "height": int(ch * 0.02),
            "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
        })
        # Decorative top-right shape
        elements.append({
            "type": "shape", "shapeType": theme.get("shapeType", "circle"),
            "x": int(cw * 0.75), "y": int(ch * 0.05),
            "width": int(ch * 0.35), "height": int(ch * 0.35),
            "backgroundColor": theme["shape1"], "borderWidth": 0, "borderColor": "transparent",
            "opacity": 15,
        })
        # Decorative bottom-left shape
        elements.append({
            "type": "shape", "shapeType": theme.get("shapeType", "circle"),
            "x": int(cw * 0.03), "y": int(ch * 0.6),
            "width": int(ch * 0.25), "height": int(ch * 0.25),
            "backgroundColor": theme["shape2"], "borderWidth": 0, "borderColor": "transparent",
            "opacity": 12,
        })
        # Title text
        title = slide_data.get("title", "")
        title_fs = int(min(cw, ch) * 0.08)
        elements.append({
            "type": "text", "content": title,
            "x": pad, "y": int(ch * 0.30), "width": cw - pad * 2, "height": int(ch * 0.20),
            "fontSize": title_fs, "fontFamily": font, "color": theme["text"],
            "textAlign": "center", "isBold": True, "backgroundColor": "transparent",
        })
        # Subtitle
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            elements.append({
                "type": "text", "content": subtitle,
                "x": int(cw * 0.15), "y": int(ch * 0.55), "width": int(cw * 0.7), "height": int(ch * 0.10),
                "fontSize": int(title_fs * 0.5), "fontFamily": font, "color": theme["muted"],
                "textAlign": "center", "isBold": False, "backgroundColor": "transparent",
            })
        # Bottom accent line
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": int(cw * 0.35), "y": int(ch * 0.50), "width": int(cw * 0.30), "height": 4,
            "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
        })
    
    elif slide_type == "section":
        # Section divider slide
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": 0, "y": 0, "width": int(cw * 0.4), "height": ch,
            "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
        })
        # Section number / icon
        section_num = slide_data.get("section_num", "")
        if section_num:
            elements.append({
                "type": "text", "content": str(section_num),
                "x": int(cw * 0.10), "y": int(ch * 0.30),
                "width": int(cw * 0.20), "height": int(ch * 0.25),
                "fontSize": int(min(cw, ch) * 0.15), "fontFamily": font,
                "color": "#ffffff", "textAlign": "center", "isBold": True,
                "backgroundColor": "transparent",
            })
        # Section title
        elements.append({
            "type": "text", "content": slide_data.get("title", ""),
            "x": int(cw * 0.47), "y": int(ch * 0.30),
            "width": int(cw * 0.48), "height": int(ch * 0.20),
            "fontSize": int(min(cw, ch) * 0.06), "fontFamily": font,
            "color": theme["text"], "textAlign": "left", "isBold": True,
            "backgroundColor": "transparent",
        })
        sub = slide_data.get("subtitle", "")
        if sub:
            elements.append({
                "type": "text", "content": sub,
                "x": int(cw * 0.47), "y": int(ch * 0.55),
                "width": int(cw * 0.48), "height": int(ch * 0.15),
                "fontSize": int(min(cw, ch) * 0.03), "fontFamily": font,
                "color": theme["muted"], "textAlign": "left", "isBold": False,
                "backgroundColor": "transparent",
            })
    
    elif slide_type == "bullets":
        # Heading
        elements.append({
            "type": "text", "content": slide_data.get("title", ""),
            "x": pad, "y": pad, "width": cw - pad * 2, "height": int(ch * 0.12),
            "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
            "color": theme["accent"], "textAlign": "left", "isBold": True,
            "backgroundColor": "transparent",
        })
        # Accent underline
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": pad, "y": int(ch * 0.16), "width": int(cw * 0.15), "height": 3,
            "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
        })
        # Bullet points
        points = slide_data.get("points", [])
        bullet_y = int(ch * 0.22)
        bullet_h = int((ch * 0.70) / max(len(points), 1))
        fs = int(min(cw, ch) * 0.033)
        for i, pt in enumerate(points[:6]):
            # Bullet dot shape
            dot_size = int(fs * 0.5)
            elements.append({
                "type": "shape", "shapeType": theme.get("shapeType", "circle"),
                "x": pad + 5, "y": bullet_y + int(bullet_h * 0.2),
                "width": dot_size, "height": dot_size,
                "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
            })
            elements.append({
                "type": "text", "content": pt,
                "x": pad + dot_size + 18, "y": bullet_y,
                "width": cw - pad * 2 - dot_size - 20, "height": bullet_h - 4,
                "fontSize": fs, "fontFamily": font, "color": theme["text"],
                "textAlign": "left", "isBold": False, "backgroundColor": "transparent",
            })
            bullet_y += bullet_h
    
    elif slide_type == "two_column":
        # Heading
        elements.append({
            "type": "text", "content": slide_data.get("title", ""),
            "x": pad, "y": pad, "width": cw - pad * 2, "height": int(ch * 0.12),
            "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
            "color": theme["accent"], "textAlign": "left", "isBold": True,
            "backgroundColor": "transparent",
        })
        # Accent underline
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": pad, "y": int(ch * 0.16), "width": int(cw * 0.15), "height": 3,
            "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
        })
        col_w = int((cw - pad * 3) / 2)
        col_top = int(ch * 0.22)
        col_h = int(ch * 0.65)
        # Left column card
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": pad, "y": col_top, "width": col_w, "height": col_h,
            "backgroundColor": theme["surface"], "borderWidth": 1, "borderColor": theme["accent"] + "40",
            "shadowLevel": 8, "borderRadius": 12,
        })
        left_title = slide_data.get("left_title", "")
        left_body = slide_data.get("left_body", "")
        if left_title:
            elements.append({
                "type": "text", "content": left_title,
                "x": pad + 15, "y": col_top + 15, "width": col_w - 30, "height": int(col_h * 0.15),
                "fontSize": int(min(cw, ch) * 0.035), "fontFamily": font,
                "color": theme["accent"], "textAlign": "left", "isBold": True,
                "backgroundColor": "transparent",
            })
        elements.append({
            "type": "text", "content": left_body,
            "x": pad + 15, "y": col_top + int(col_h * 0.20), "width": col_w - 30, "height": int(col_h * 0.75),
            "fontSize": int(min(cw, ch) * 0.028), "fontFamily": font,
            "color": theme["text"], "textAlign": "left", "isBold": False,
            "backgroundColor": "transparent",
        })
        # Right column card
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": pad * 2 + col_w, "y": col_top, "width": col_w, "height": col_h,
            "backgroundColor": theme["surface"], "borderWidth": 1, "borderColor": theme["shape2"] + "40",
            "shadowLevel": 8, "borderRadius": 12,
        })
        right_title = slide_data.get("right_title", "")
        right_body = slide_data.get("right_body", "")
        if right_title:
            elements.append({
                "type": "text", "content": right_title,
                "x": pad * 2 + col_w + 15, "y": col_top + 15, "width": col_w - 30, "height": int(col_h * 0.15),
                "fontSize": int(min(cw, ch) * 0.035), "fontFamily": font,
                "color": theme["shape2"], "textAlign": "left", "isBold": True,
                "backgroundColor": "transparent",
            })
        elements.append({
            "type": "text", "content": right_body,
            "x": pad * 2 + col_w + 15, "y": col_top + int(col_h * 0.20), "width": col_w - 30, "height": int(col_h * 0.75),
            "fontSize": int(min(cw, ch) * 0.028), "fontFamily": font,
            "color": theme["text"], "textAlign": "left", "isBold": False,
            "backgroundColor": "transparent",
        })
    
    elif slide_type == "quote":
        # Large quote mark
        elements.append({
            "type": "text", "content": "\u201c",
            "x": int(cw * 0.08), "y": int(ch * 0.15),
            "width": int(cw * 0.15), "height": int(ch * 0.25),
            "fontSize": int(min(cw, ch) * 0.2), "fontFamily": "Serif",
            "color": theme["accent"] + "50", "textAlign": "left", "isBold": True,
            "backgroundColor": "transparent",
        })
        elements.append({
            "type": "text", "content": slide_data.get("quote", ""),
            "x": int(cw * 0.12), "y": int(ch * 0.30),
            "width": int(cw * 0.76), "height": int(ch * 0.30),
            "fontSize": int(min(cw, ch) * 0.04), "fontFamily": font,
            "color": theme["text"], "textAlign": "center", "isBold": False,
            "isItalic": True, "backgroundColor": "transparent",
        })
        author = slide_data.get("author", "")
        if author:
            elements.append({
                "type": "text", "content": f"— {author}",
                "x": int(cw * 0.3), "y": int(ch * 0.65),
                "width": int(cw * 0.4), "height": int(ch * 0.08),
                "fontSize": int(min(cw, ch) * 0.025), "fontFamily": font,
                "color": theme["accent"], "textAlign": "center", "isBold": True,
                "backgroundColor": "transparent",
            })
    
    elif slide_type == "big_number":
        title = slide_data.get("title", "")
        if title:
            elements.append({
                "type": "text", "content": title,
                "x": pad, "y": pad, "width": cw - pad * 2, "height": int(ch * 0.12),
                "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
                "color": theme["accent"], "textAlign": "center", "isBold": True,
                "backgroundColor": "transparent",
            })
            
        icon_name = slide_data.get("icon", "")
        if icon_name:
            icon_size = int(ch * 0.15)
            elements.append({
                "type": "icon", "svgContent": f'<svg width="100%" height="100%"><use href="#icon-{icon_name}"></use></svg>',
                "x": int(cw / 2) - int(icon_size / 2), "y": int(ch * 0.20), "width": icon_size, "height": icon_size,
                "borderColor": theme["accent2"], "borderWidth": 2, "backgroundColor": "transparent"
            })
            
        number_y = int(ch * 0.35) if icon_name else int(ch * 0.25)
        elements.append({
            "type": "text", "content": slide_data.get("number", "100%"),
            "x": pad, "y": number_y, "width": cw - pad * 2, "height": int(ch * 0.3),
            "fontSize": int(min(cw, ch) * 0.25), "fontFamily": font,
            "color": theme["accent2"], "textAlign": "center", "isBold": True,
            "backgroundColor": "transparent",
        })
        elements.append({
            "type": "text", "content": slide_data.get("description", ""),
            "x": int(cw * 0.15), "y": int(ch * 0.7), "width": int(cw * 0.7), "height": int(ch * 0.2),
            "fontSize": int(min(cw, ch) * 0.04), "fontFamily": font,
            "color": theme["text"], "textAlign": "center", "isBold": False,
            "backgroundColor": "transparent",
        })

    elif slide_type == "comparison":
        elements.append({
            "type": "text", "content": slide_data.get("title", "Comparison"),
            "x": pad, "y": pad, "width": cw - pad * 2, "height": int(ch * 0.12),
            "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
            "color": theme["accent"], "textAlign": "center", "isBold": True,
            "backgroundColor": "transparent",
        })
        col_w = int((cw - pad * 3) / 2)
        col_top = int(ch * 0.20)
        
        # Left Side
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": pad, "y": col_top, "width": col_w, "height": int(ch * 0.65),
            "backgroundColor": theme["surface"], "borderWidth": 2, "borderColor": theme["shape1"],
            "shadowLevel": 10, "borderRadius": 12,
        })
        elements.append({
            "type": "text", "content": slide_data.get("left_title", ""),
            "x": pad + 15, "y": col_top + 15, "width": col_w - 30, "height": int(ch * 0.1),
            "fontSize": int(min(cw, ch) * 0.04), "fontFamily": font,
            "color": theme["shape1"], "textAlign": "center", "isBold": True,
            "backgroundColor": "transparent",
        })
        left_body = "\n".join(slide_data.get("left_points", []))
        elements.append({
            "type": "text", "content": left_body,
            "x": pad + 15, "y": col_top + int(ch * 0.15), "width": col_w - 30, "height": int(ch * 0.45),
            "fontSize": int(min(cw, ch) * 0.03), "fontFamily": font,
            "color": theme["text"], "textAlign": "left", "isBold": False,
            "isList": True, "backgroundColor": "transparent",
        })
        
        # Right Side
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": pad * 2 + col_w, "y": col_top, "width": col_w, "height": int(ch * 0.65),
            "backgroundColor": theme["surface"], "borderWidth": 2, "borderColor": theme["shape2"],
            "shadowLevel": 10, "borderRadius": 12,
        })
        elements.append({
            "type": "text", "content": slide_data.get("right_title", ""),
            "x": pad * 2 + col_w + 15, "y": col_top + 15, "width": col_w - 30, "height": int(ch * 0.1),
            "fontSize": int(min(cw, ch) * 0.04), "fontFamily": font,
            "color": theme["shape2"], "textAlign": "center", "isBold": True,
            "backgroundColor": "transparent",
        })
        right_body = "\n".join(slide_data.get("right_points", []))
        elements.append({
            "type": "text", "content": right_body,
            "x": pad * 2 + col_w + 15, "y": col_top + int(ch * 0.15), "width": col_w - 30, "height": int(ch * 0.45),
            "fontSize": int(min(cw, ch) * 0.03), "fontFamily": font,
            "color": theme["text"], "textAlign": "left", "isBold": False,
            "isList": True, "backgroundColor": "transparent",
        })

    elif slide_type == "image_focus":
        elements.append({
            "type": "text", "content": slide_data.get("title", ""),
            "x": pad, "y": pad, "width": int(cw * 0.4), "height": int(ch * 0.15),
            "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
            "color": theme["accent"], "textAlign": "left", "isBold": True,
            "backgroundColor": "transparent",
        })
        elements.append({
            "type": "text", "content": slide_data.get("caption", ""),
            "x": pad, "y": int(ch * 0.3), "width": int(cw * 0.4), "height": int(ch * 0.5),
            "fontSize": int(min(cw, ch) * 0.035), "fontFamily": font,
            "color": theme["text"], "textAlign": "left", "isBold": False,
            "backgroundColor": "transparent",
        })
        # Image area with SVG illustration fallback
        illustration = slide_data.get("illustration", "")
        img_x = int(cw * 0.5)
        img_w = int(cw * 0.5) - pad
        img_h = ch - pad * 2
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": img_x, "y": pad, "width": img_w, "height": img_h,
            "backgroundColor": theme["muted"], "borderWidth": 0, "borderColor": "transparent",
            "opacity": 15, "shadowLevel": 6, "borderRadius": 16,
        })
        if illustration and illustration in _SVG_ILLUSTRATIONS:
            svg_uri = _get_svg_data_uri(illustration, theme["accent"], theme["accent2"], img_w - 40, img_h - 40)
            if svg_uri:
                elements.append({
                    "type": "image", "src": svg_uri,
                    "x": img_x + 20, "y": pad + 20, "width": img_w - 40, "height": img_h - 40,
                    "opacity": 90,
                })
        else:
            elements.append({
                "type": "text", "content": "IMAGE",
                "x": img_x, "y": int(ch * 0.45), "width": img_w, "height": int(ch * 0.1),
                "fontSize": int(min(cw, ch) * 0.03), "fontFamily": font,
                "color": theme["muted"], "textAlign": "center", "isBold": True,
                "backgroundColor": "transparent", "opacity": 40,
            })

    elif slide_type == "mermaid":
        elements.append({
            "type": "text", "content": slide_data.get("title", "Diagram"),
            "x": pad, "y": pad, "width": cw - pad * 2, "height": int(ch * 0.12),
            "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
            "color": theme["accent"], "textAlign": "center", "isBold": True,
            "backgroundColor": "transparent",
        })
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": pad, "y": int(ch * 0.2), "width": cw - pad * 2, "height": int(ch * 0.7),
            "backgroundColor": theme["surface"], "borderWidth": 1, "borderColor": theme["accent"] + "40",
        })
        elements.append({
            "type": "mermaid", "content": slide_data.get("mermaid_code", "graph TD;\nA-->B;"),
            "x": pad + 10, "y": int(ch * 0.2) + 10, "width": cw - pad * 2 - 20, "height": int(ch * 0.7) - 20,
            "backgroundColor": "transparent", "color": theme["text"],
        })

    elif slide_type == "three_column":
        elements.append({
            "type": "text", "content": slide_data.get("title", ""),
            "x": pad, "y": pad, "width": cw - pad * 2, "height": int(ch * 0.12),
            "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
            "color": theme["accent"], "textAlign": "center", "isBold": True,
            "backgroundColor": "transparent",
        })
        col_w = int((cw - pad * 4) / 3)
        col_top = int(ch * 0.25)
        col_h = int(ch * 0.60)
        
        cols = [
            ("col1_title", "col1_body", "col1_icon", theme["shape1"]),
            ("col2_title", "col2_body", "col2_icon", theme["accent"]),
            ("col3_title", "col3_body", "col3_icon", theme["shape2"])
        ]
        
        for i, (t_key, b_key, i_key, col_color) in enumerate(cols):
            x_pos = pad + i * (col_w + pad)
            icon_name = slide_data.get(i_key, "")
            
            if icon_name:
                elements.append({
                    "type": "icon", "svgContent": f'<svg width="100%" height="100%"><use href="#icon-{icon_name}"></use></svg>',
                    "x": x_pos + int(col_w/2) - int(ch*0.06), "y": col_top - int(ch*0.06), 
                    "width": int(ch*0.12), "height": int(ch*0.12),
                    "borderColor": col_color, "borderWidth": 2, "backgroundColor": theme["surface"]
                })
            else:
                elements.append({
                    "type": "shape", "shapeType": theme.get("shapeType", "circle"),
                    "x": x_pos + int(col_w/2) - int(ch*0.06), "y": col_top - int(ch*0.06), 
                    "width": int(ch*0.12), "height": int(ch*0.12),
                    "backgroundColor": col_color, "borderWidth": 0, "borderColor": "transparent",
                })
                
            elements.append({
                "type": "shape", "shapeType": "rounded-rect",
                "x": x_pos, "y": col_top + int(ch*0.06), "width": col_w, "height": col_h - int(ch*0.06),
                "backgroundColor": theme["surface"], "borderWidth": 1, "borderColor": col_color + "40",
            })
            elements.append({
                "type": "text", "content": slide_data.get(t_key, ""),
                "x": x_pos + 10, "y": col_top + int(ch*0.10), "width": col_w - 20, "height": int(ch * 0.10),
                "fontSize": int(min(cw, ch) * 0.035), "fontFamily": font,
                "color": col_color, "textAlign": "center", "isBold": True,
                "backgroundColor": "transparent",
            })
            elements.append({
                "type": "text", "content": slide_data.get(b_key, ""),
                "x": x_pos + 15, "y": col_top + int(ch * 0.22), "width": col_w - 30, "height": int(ch * 0.35),
                "fontSize": int(min(cw, ch) * 0.025), "fontFamily": font,
                "color": theme["text"], "textAlign": "center", "isBold": False,
                "backgroundColor": "transparent",
            })

    elif slide_type == "illustrated":
        # Slide with SVG illustration + text
        illustration = slide_data.get("illustration", "rocket")
        position = slide_data.get("illustration_position", "right")
        title = slide_data.get("title", "")
        body = slide_data.get("body", "")
        
        # Title
        elements.append({
            "type": "text", "content": title,
            "x": pad, "y": pad, "width": cw - pad * 2, "height": int(ch * 0.12),
            "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
            "color": theme["accent"], "textAlign": "left", "isBold": True,
            "backgroundColor": "transparent",
        })
        # Accent underline
        elements.append({
            "type": "shape", "shapeType": "rounded-rect",
            "x": pad, "y": int(ch * 0.16), "width": int(cw * 0.15), "height": 3,
            "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
        })
        
        illust_size = int(min(cw, ch) * 0.55)
        text_w = cw - illust_size - pad * 3
        
        if position == "left":
            illust_x = pad
            text_x = pad * 2 + illust_size
        else:
            text_x = pad
            illust_x = cw - pad - illust_size
        
        illust_y = int(ch * 0.22)
        
        # Illustration background glow
        elements.append({
            "type": "shape", "shapeType": "circle",
            "x": illust_x + int(illust_size * 0.1), "y": illust_y + int(illust_size * 0.1),
            "width": int(illust_size * 0.8), "height": int(illust_size * 0.8),
            "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
            "opacity": 10,
        })
        
        # SVG illustration
        if illustration and illustration in _SVG_ILLUSTRATIONS:
            svg_uri = _get_svg_data_uri(illustration, theme["accent"], theme["accent2"], illust_size, illust_size)
            if svg_uri:
                elements.append({
                    "type": "image", "src": svg_uri,
                    "x": illust_x, "y": illust_y, "width": illust_size, "height": illust_size,
                    "opacity": 90,
                })
        
        # Body text
        elements.append({
            "type": "text", "content": body,
            "x": text_x, "y": int(ch * 0.25), "width": text_w, "height": int(ch * 0.60),
            "fontSize": int(min(cw, ch) * 0.03), "fontFamily": font,
            "color": theme["text"], "textAlign": "left", "isBold": False,
            "backgroundColor": "transparent",
        })

    elif slide_type == "closing":
        # Decorative shapes
        elements.append({
            "type": "shape", "shapeType": theme.get("shapeType", "circle"),
            "x": int(cw * 0.05), "y": int(ch * 0.1),
            "width": int(ch * 0.2), "height": int(ch * 0.2),
            "backgroundColor": theme["shape1"], "borderWidth": 0, "borderColor": "transparent",
            "opacity": 18,
        })
        elements.append({
            "type": "shape", "shapeType": theme.get("shapeType", "circle"),
            "x": int(cw * 0.80), "y": int(ch * 0.65),
            "width": int(ch * 0.30), "height": int(ch * 0.30),
            "backgroundColor": theme["shape2"], "borderWidth": 0, "borderColor": "transparent",
            "opacity": 12,
        })
        elements.append({
            "type": "text", "content": slide_data.get("title", "Thank You"),
            "x": pad, "y": int(ch * 0.30), "width": cw - pad * 2, "height": int(ch * 0.18),
            "fontSize": int(min(cw, ch) * 0.08), "fontFamily": font,
            "color": theme["text"], "textAlign": "center", "isBold": True,
            "backgroundColor": "transparent",
        })
        sub = slide_data.get("subtitle", "")
        if sub:
            elements.append({
                "type": "text", "content": sub,
                "x": int(cw * 0.2), "y": int(ch * 0.55),
                "width": int(cw * 0.6), "height": int(ch * 0.10),
                "fontSize": int(min(cw, ch) * 0.03), "fontFamily": font,
                "color": theme["muted"], "textAlign": "center", "isBold": False,
                "backgroundColor": "transparent",
            })
        # Bottom accent bar
        elements.append({
            "type": "shape", "shapeType": "rect",
            "x": 0, "y": ch - int(ch * 0.02), "width": cw, "height": int(ch * 0.02),
            "backgroundColor": theme["accent"], "borderWidth": 0, "borderColor": "transparent",
        })
    
    else:  # generic "content" fallback
        elements.append({
            "type": "text", "content": slide_data.get("title", ""),
            "x": pad, "y": pad, "width": cw - pad * 2, "height": int(ch * 0.12),
            "fontSize": int(min(cw, ch) * 0.05), "fontFamily": font,
            "color": theme["accent"], "textAlign": "left", "isBold": True,
            "backgroundColor": "transparent",
        })
        body = slide_data.get("body", "")
        elements.append({
            "type": "text", "content": body,
            "x": pad, "y": int(ch * 0.20), "width": cw - pad * 2, "height": int(ch * 0.70),
            "fontSize": int(min(cw, ch) * 0.03), "fontFamily": font,
            "color": theme["text"], "textAlign": "left", "isBold": False,
            "backgroundColor": "transparent",
        })
    
    fs_small = int(min(cw, ch) * 0.02)
    # Header
    if header_text and slide_type != "title":
        elements.append({
            "type": "text", "content": header_text,
            "x": pad, "y": int(ch * 0.02), "width": cw - 2*pad, "height": 20,
            "fontSize": fs_small, "fontFamily": font,
            "color": theme["muted"], "textAlign": "left", "isBold": False,
            "backgroundColor": "transparent",
        })
        
    # Footer
    if footer_text and slide_type != "title":
        elements.append({
            "type": "text", "content": footer_text,
            "x": pad, "y": ch - int(ch * 0.05), "width": cw - 2*pad, "height": 20,
            "fontSize": fs_small, "fontFamily": font,
            "color": theme["muted"], "textAlign": "left", "isBold": False,
            "backgroundColor": "transparent",
        })

    # Page number
    page_num = slide_data.get("page_num", "")
    if page_num and page_number_position != "none":
        px, py, palign = cw - pad - 30, ch - int(ch * 0.05), "right"
        if page_number_position == "bottom-center":
            px, py, palign = cw // 2 - 15, ch - int(ch * 0.05), "center"
        elif page_number_position == "top-right":
            px, py, palign = cw - pad - 30, int(ch * 0.02), "right"
            
        elements.append({
            "type": "text", "content": str(page_num),
            "x": px, "y": py, "width": 30, "height": 20,
            "fontSize": fs_small, "fontFamily": font,
            "color": theme["muted"], "textAlign": palign, "isBold": False,
            "backgroundColor": "transparent",
        })
    
    return {"elements": elements}

_PREDEFINED_TEMPLATES = {
    "pitch_deck": {
        "title": "Startup Pitch Deck",
        "theme": "sapphire",
        "slides": [
            {"type": "title", "title": "Our Vision", "subtitle": "Disrupting the industry"},
            {"type": "bullets", "title": "The Problem", "points": ["Current solutions are slow", "High costs", "Poor user experience"]},
            {"type": "two_column", "title": "Our Solution", "left_title": "Faster", "left_body": "10x faster than competitors", "right_title": "Cheaper", "right_body": "50% lower cost"},
            {"type": "two_column", "title": "Market Size", "left_title": "TAM", "left_body": "$10 Billion total addressable market", "right_title": "SAM", "right_body": "$2 Billion serviceable addressable market"},
            {"type": "closing", "title": "Join Us", "subtitle": "Contact: founders@example.com"}
        ]
    },
    "portfolio": {
        "title": "Creative Portfolio",
        "theme": "monochrome",
        "slides": [
            {"type": "title", "title": "My Portfolio", "subtitle": "Creative Designer & Developer"},
            {"type": "two_column", "title": "About Me", "left_title": "Background", "left_body": "5+ years in digital design", "right_title": "Skills", "right_body": "UI/UX, Branding, Web Dev"},
            {"type": "bullets", "title": "Featured Projects", "points": ["E-commerce Redesign", "Brand Identity for Startup", "Mobile App UI"]},
            {"type": "closing", "title": "Let's Connect", "subtitle": "hello@myportfolio.com"}
        ]
    },
    "course_module": {
        "title": "Course Module",
        "theme": "nordic",
        "slides": [
            {"type": "title", "title": "Introduction to React", "subtitle": "Module 1: Fundamentals"},
            {"type": "bullets", "title": "Learning Objectives", "points": ["Understand components", "Learn state management", "Master props"]},
            {"type": "two_column", "title": "Core Concepts", "left_title": "JSX", "left_body": "JavaScript XML syntax", "right_title": "Virtual DOM", "right_body": "Efficient UI rendering"},
            {"type": "closing", "title": "Homework", "subtitle": "Build a simple counter app"}
        ]
    },
    "status_report": {
        "title": "Weekly Status Report",
        "theme": "corporate",
        "slides": [
            {"type": "title", "title": "Weekly Sync", "subtitle": "Project Alpha Status"},
            {"type": "bullets", "title": "Executive Summary", "points": ["Milestone 1 completed", "Budget is on track", "Minor delay in QQA"]},
            {"type": "two_column", "title": "Metrics", "left_title": "Completed", "left_body": "45 tasks finished this week", "right_title": "Pending", "right_body": "12 tasks in backlog"},
            {"type": "bullets", "title": "Roadblocks", "points": ["Waiting on API access", "Design approvals delayed"]},
            {"type": "closing", "title": "Next Steps", "subtitle": "Focus on backend integration"}
        ]
    },
    "marketing_plan": {
        "title": "Marketing Strategy",
        "theme": "emerald",
        "slides": [
            {"type": "title", "title": "Q3 Marketing Plan", "subtitle": "Growth & Acquisition"},
            {"type": "bullets", "title": "Target Audience", "points": ["Gen Z professionals", "Tech enthusiasts", "Small business owners"]},
            {"type": "two_column", "title": "Channels", "left_title": "Social Media", "left_body": "Instagram, TikTok, LinkedIn", "right_title": "Paid Search", "right_body": "Google Ads, Bing"},
            {"type": "closing", "title": "Q&A", "subtitle": "Open floor for questions"}
        ]
    },
    "event_planning": {
        "title": "Event Logistics",
        "theme": "lavender",
        "slides": [
            {"type": "title", "title": "Annual Conference", "subtitle": "Planning & Logistics"},
            {"type": "bullets", "title": "Schedule Overview", "points": ["Day 1: Keynotes", "Day 2: Workshops", "Day 3: Networking"]},
            {"type": "two_column", "title": "Venue", "left_title": "Location", "left_body": "Downtown Convention Center", "right_title": "Capacity", "right_body": "Up to 500 attendees"},
            {"type": "closing", "title": "Action Items", "subtitle": "Finalize catering by Friday"}
        ]
    }
}

@router.get("/templates")
async def list_templates():
    templates = []
    for t_id, t_data in _PREDEFINED_TEMPLATES.items():
        theme_name = t_data["theme"]
        theme_colors = _THEMES.get(theme_name, _THEMES["midnight"])
        templates.append({
            "id": t_id,
            "title": t_data["title"],
            "theme": theme_name,
            "colors": theme_colors
        })
    return {"templates": templates}

@router.get("/template/{template_id}")
async def get_template(template_id: str):
    t_data = _PREDEFINED_TEMPLATES.get(template_id)
    if not t_data:
        raise HTTPException(404, "Template not found")
        
    theme_name = t_data["theme"]
    theme = _THEMES.get(theme_name, _THEMES["midnight"])
    font = "Sans-Serif"
    cw, ch = 960, 540
    
    pages = []
    for i, slide in enumerate(t_data["slides"]):
        slide["page_num"] = i + 1
        page = _build_slide_layout(slide, cw, ch, theme, font)
        pages.append(page)
        
    return {
        "title": t_data["title"], 
        "pages": pages, 
        "canvas_width": cw, 
        "canvas_height": ch
    }


@router.post("/ai-generate")
async def ai_generate_presentation(request: Request):
    """Generate a presentation: Phase 1 = LLM content plan, Phase 2 = Python layout engine."""
    import httpx, json, re
    from config import OLLAMA_HOST, OLLAMA_MODEL
    
    data = await request.json()
    prompt = data.get("prompt", "").strip()
    canvas_width = data.get("canvas_width", 960)
    canvas_height = data.get("canvas_height", 540)
    
    if not prompt:
        raise HTTPException(400, "Prompt is required")
        
    # Extract URLs and scrape content
    url_matches = re.findall(r"https?://[^\s]+", prompt)
    scraped_text = ""
    if url_matches:
        from bs4 import BeautifulSoup
        for url in url_matches:
            try:
                async with httpx.AsyncClient(timeout=10.0, follow_redirects=True) as client:
                    res = await client.get(url)
                    if res.status_code == 200:
                        soup = BeautifulSoup(res.text, "html.parser")
                        for script in soup(["script", "style"]):
                            script.decompose()
                        text = soup.get_text(separator=" ")
                        text = re.sub(r'\s+', ' ', text).strip()
                        scraped_text += f"\n\nContent from {url}:\n{text[:2500]}"
            except Exception as e:
                logger.warning(f"Failed to scrape {url}: {e}")
                
    if scraped_text:
        prompt += f"\n\nUse the following scraped context to build the presentation:{scraped_text}"
    
    # Phase 1: Ask LLM for CONTENT ONLY (no coordinates)
    svg_names = ", ".join(sorted(_SVG_ILLUSTRATIONS.keys()))
    system_prompt = f"""You are an expert presentation content strategist. Generate a structured JSON content plan.
OUTPUT ONLY valid JSON. No markdown fences, no explanation, no extra text.

The JSON structure must be:
{{
  "title": "Presentation Title",
  "header_text": "Optional header text for all slides",
  "footer_text": "Optional footer text for all slides",
  "page_number_position": "bottom-right|bottom-center|top-right|none",
  "theme": {{
    "bg": "#0f0f1a",
    "surface": "#1a1a2e",
    "accent": "#7c5cfc",
    "accent2": "#f472b6",
    "text": "#ffffff",
    "muted": "#a0a0b8",
    "shape1": "#7c5cfc",
    "shape2": "#f472b6",
    "shapeType": "circle|diamond|hexagon|rounded-rect|star|cloud|parallelogram|pentagon"
  }},
  "font": "Sans-Serif|Serif|Monospace",
  "slides": [
    {{"type": "title", "title": "...", "subtitle": "..."}},
    {{"type": "section", "section_num": "01", "title": "...", "subtitle": "..."}},
    {{"type": "bullets", "title": "...", "points": ["point 1", "point 2", "point 3"]}},
    {{"type": "two_column", "title": "...", "left_title": "...", "left_body": "...", "right_title": "...", "right_body": "..."}},
    {{"type": "three_column", "title": "...", "col1_icon": "activity", "col1_title": "...", "col1_body": "...", "col2_icon": "shield", "col2_title": "...", "col2_body": "...", "col3_icon": "sparkles", "col3_title": "...", "col3_body": "..."}},
    {{"type": "comparison", "title": "...", "left_title": "...", "left_points": ["..."], "right_title": "...", "right_points": ["..."]}},
    {{"type": "big_number", "title": "...", "number": "99%", "description": "...", "icon": "bar-chart"}},
    {{"type": "image_focus", "title": "...", "caption": "...", "illustration": "rocket"}},
    {{"type": "illustrated", "title": "...", "body": "...", "illustration": "growth", "illustration_position": "right"}},
    {{"type": "quote", "quote": "...", "author": "..."}},
    {{"type": "mermaid", "title": "Architecture", "mermaid_code": "graph TD;\\nA-->B;"}},
    {{"type": "closing", "title": "Thank You", "subtitle": "contact info or closing message"}}
  ]
}}

SLIDE TYPES available:
- "title": Opening slide
- "section": Section divider
- "bullets": Bullet points
- "two_column": Two side-by-side cards with shadow effects
- "three_column": Three side-by-side columns/steps
- "comparison": Side-by-side pros/cons or before/after comparison with shadow effects
- "big_number": A single giant metric/number with description
- "image_focus": A text side with a large illustration. Use "illustration" field to pick an SVG.
- "illustrated": A slide with an SVG illustration + text body. Use "illustration" for the SVG name and "illustration_position" (left or right).
- "mermaid": A slide containing a Mermaid.js diagram (generate actual mermaid code in mermaid_code)
- "quote": Highlighted quote
- "closing": Final thank-you

AVAILABLE SVG ILLUSTRATIONS for "illustrated" and "image_focus" slides:
{svg_names}

RULES:
- Always start with "title", end with "closing".
- Adjust the number of slides to what the user requested. If not specified, default to 5-8 slides total.
- Keep text short and impactful.
- Design a beautiful custom color palette in the "theme" object. CRITICAL: You must choose between a LIGHT theme (white/light background with dark text) or a DARK theme (dark background with light text) based on the mood. Do not just use dark themes. Make sure text contrast is high!
- Pick a shapeType that matches the mood.
- Use at least one "illustrated" or "image_focus" slide with a relevant SVG illustration to make slides visually rich.
- You can optionally use UI icons in 'three_column' and 'big_number' slides. Available icons: activity, bar-chart, bolt, camera, check, clock, cloud, cpu, heart, image, layers, monitor, palette, play, settings, shield, smartphone, sparkles, star, terminal, video.
- If the user asks for diagrams, graphs, or architecture, USE the "mermaid" slide type with valid mermaid.js code.
"""
    
    payload = {
        "model": OLLAMA_MODEL,
        "prompt": f"Create a presentation about: {prompt}",
        "system": system_prompt,
        "stream": False,
        "options": {"temperature": 0.7, "num_predict": 4096, "num_ctx": 8192},
    }
    
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            resp = await client.post(f"{OLLAMA_HOST}/api/generate", json=payload)
            resp.raise_for_status()
            result = resp.json()
            original_raw = result.get("response", "")
            raw = original_raw
        
        # Clean LLM output: handle unclosed think blocks
        if "<think>" in raw and "</think>" not in raw:
            raw += "</think>"
        raw = re.sub(r"<think>.*?</think>", "", raw, flags=re.DOTALL).strip()
        
        json_str = None
        # First try to find markdown JSON block
        md_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", raw, flags=re.DOTALL)
        if md_match:
            json_str = md_match.group(1)
        else:
            # Fallback to brace matching
            start_idx = raw.find('{')
            if start_idx != -1:
                brace_count = 0
                end_idx = -1
                in_string = False
                escape = False
                for i in range(start_idx, len(raw)):
                    char = raw[i]
                    if escape:
                        escape = False
                    elif char == '\\':
                        escape = True
                    elif char == '"':
                        in_string = not in_string
                    elif not in_string:
                        if char == '{':
                            brace_count += 1
                        elif char == '}':
                            brace_count -= 1
                            if brace_count == 0:
                                end_idx = i
                                break
                if end_idx != -1:
                    json_str = raw[start_idx:end_idx+1]

        if not json_str:
            logger.error(f"AI raw output (no JSON found):\n---ORIGINAL---\n{original_raw[:1000]}\n---CLEANED---\n{raw[:500]}")
            raise HTTPException(500, "AI did not produce valid JSON content plan")
        
        content_plan = json.loads(json_str)
        
        # Phase 2: Build pixel-perfect layouts from content plan
        title = content_plan.get("title", "AI Presentation")
        theme_input = content_plan.get("theme")
        theme_name = "custom"
        if isinstance(theme_input, dict):
            theme = theme_input
        else:
            theme_name = theme_input if isinstance(theme_input, str) else "midnight"
            theme = _THEMES.get(theme_name, _THEMES["midnight"])
        
        # Ensure fallback keys
        default_theme = _THEMES["midnight"]
        for k, v in default_theme.items():
            if k not in theme:
                theme[k] = v
        font = content_plan.get("font", "Sans-Serif")
        header_text = content_plan.get("header_text", "")
        footer_text = content_plan.get("footer_text", "")
        page_number_position = content_plan.get("page_number_position", "bottom-right")
        slides = content_plan.get("slides", [])
        
        if not slides:
            raise HTTPException(500, "AI generated empty slide plan")
        
        pages = []
        for i, slide in enumerate(slides):
            slide["page_num"] = i + 1
            page = _build_slide_layout(slide, canvas_width, canvas_height, theme, font, header_text, footer_text, page_number_position)
            pages.append(page)
        
        logger.info(f"AI presentation built: '{title}' — {len(pages)} slides, theme={theme_name}")
        return {"title": title, "pages": pages}
        
    except httpx.TimeoutException:
        raise HTTPException(504, "AI generation timed out. Try a simpler prompt.")
    except json.JSONDecodeError as e:
        logger.error(f"AI JSON parse error: {e}\nRaw: {raw[:500]}")
        raise HTTPException(500, f"AI produced invalid JSON: {str(e)[:200]}")
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"AI generation error: {e}")
        raise HTTPException(500, str(e))

@router.post("/ai-enrich")
async def ai_enrich_presentation(request: Request):
    import httpx, json
    from config import OLLAMA_HOST, OLLAMA_MODEL
    
    data = await request.json()
    prompt = data.get("prompt", "").strip()
    mode = data.get("mode", "enhance")
    
    if not prompt:
        raise HTTPException(400, "Prompt is required")
        
    system_prompt = "You are an expert presentation designer and copywriter. Keep responses short and directly usable."
    if mode == "enhance":
        user_prompt = f"Improve, rewrite, or generate presentation content based on this request. If there is existing content, modify it according to the task. Return ONLY the final text or list, no markdown formatting like backticks, no quotes, no conversational filler:\n\n{prompt}"
    else:
        user_prompt = f"Based on this slide text: '{prompt}', suggest 2-3 brief, impactful bullet points to add. Return ONLY the bullet points separated by new lines, no intro/outro, no markdown."
        
    payload = {
        "model": OLLAMA_MODEL,
        # /no_think disables Qwen 3 reasoning chain for fast output
        "prompt": f"/no_think\n{user_prompt}",
        "system": system_prompt,
        "stream": False,
        "options": {"temperature": 0.6, "num_predict": 4096},
    }
    
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            resp = await client.post(f"{OLLAMA_HOST}/api/generate", json=payload)
            resp.raise_for_status()
            result = resp.json()
            raw = result.get("response", "").strip()
        
        # Clean output
        import re
        raw = re.sub(r"<think>.*?</think>", "", raw, flags=re.DOTALL).strip()
        
        return {"result": raw}
    except Exception as e:
        logger.error(f"AI enrich error: {e}")
        raise HTTPException(500, str(e))


@router.get("/svg-catalog")
async def get_svg_catalog():
    """Return the full SVG illustration catalog with categories and rendered previews."""
    catalog = []
    for name, template in _SVG_ILLUSTRATIONS.items():
        # Use neutral colors for preview
        preview_svg = _get_svg_illustration(name, "#6366f1", "#f472b6", 80, 80)
        catalog.append({
            "name": name,
            "svg": preview_svg,
        })
    return {
        "illustrations": catalog,
        "categories": _SVG_CATEGORIES,
    }


@router.get("/stock-images")
async def get_stock_images(q: str = "", category: str = ""):
    """Return stock images from the local open-source image library."""
    catalog_path = os.path.join("static", "img", "stock", "catalog.json")
    
    if not os.path.exists(catalog_path):
        return {
            "images": [],
            "categories": [],
            "message": "Stock image library not installed. Run: python scripts/download_stock_images.py"
        }
    
    import json as _json
    with open(catalog_path, "r") as f:
        catalog = _json.load(f)
    
    images = catalog.get("images", [])
    categories = catalog.get("categories", [])
    
    # Filter by category
    if category:
        images = [img for img in images if img.get("category") == category]
    
    # Filter by search query
    if q:
        q_lower = q.lower()
        images = [
            img for img in images
            if q_lower in img.get("description", "").lower()
            or q_lower in img.get("category", "").lower()
            or q_lower in img.get("filename", "").lower()
        ]
    
    return {
        "images": images,
        "categories": categories,
        "source": "local",
        "total": len(images),
    }
