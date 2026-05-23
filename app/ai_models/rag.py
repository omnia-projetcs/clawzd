"""
Clawzd — RAG (Retrieval-Augmented Generation) module.
Uses ChromaDB for vector storage and sentence-transformers for embeddings.
Supports multi-format documents: md, csv, docx, pptx, xlsx, txt, pdf, code files.
"""
import os
import logging
import hashlib
from datetime import datetime, timezone
from fastapi import APIRouter, UploadFile, File, HTTPException, Form
from typing import List
from config import CHROMA_DB_PATH, RAG_DIR

router = APIRouter()
logger = logging.getLogger("clawzd.rag")

# Lazy-loaded globals to avoid slow startup if not used
_client = None
_collection = None
_encoder = None

# Supported file extensions by category
_TEXT_EXTENSIONS = {".txt", ".md", ".log", ".env", ".toml", ".ini", ".yaml", ".yml", ".xml", ".html", ".css", ".json"}
_CODE_EXTENSIONS = {".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp",
                    ".rb", ".php", ".sh", ".bash", ".sql", ".r", ".swift", ".kt", ".scala", ".lua",
                    ".pl", ".m", ".cs", ".dart", ".zig", ".vue", ".svelte"}
_OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx"}
_DATA_EXTENSIONS = {".csv"}
_PDF_EXTENSIONS = {".pdf"}
_EBOOK_EXTENSIONS = {".epub"}
_ARCHIVE_EXTENSIONS = {".zip", ".tar", ".gz", ".tgz"}
_ALL_SUPPORTED = _TEXT_EXTENSIONS | _CODE_EXTENSIONS | _OFFICE_EXTENSIONS | _DATA_EXTENSIONS | _PDF_EXTENSIONS | _EBOOK_EXTENSIONS | _ARCHIVE_EXTENSIONS

# Path for locally cached embedding model (avoids HuggingFace Hub calls on every startup)
_EMBEDDING_CACHE_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "models", "embeddings")

# Track file hashes for incremental indexing
_indexed_hashes: dict[str, str] = {}


def _get_rag():
    """Lazy-initialize ChromaDB client and encoder."""
    global _client, _collection, _encoder
    if _client is None:
        try:
            import chromadb
            from chromadb.config import Settings
            from sentence_transformers import SentenceTransformer
        except (ImportError, ModuleNotFoundError) as e:
            logger.warning("RAG dependencies are not installed on this system: %s. RAG features are disabled.", e)
            raise HTTPException(503, "RAG dependencies (chromadb/sentence-transformers) are not installed on this system.")

        _client = chromadb.PersistentClient(
            path=CHROMA_DB_PATH,
            settings=Settings(anonymized_telemetry=False),
        )
        _collection = _client.get_or_create_collection("knowledge_base")

        # Force offline mode if model is already cached — no HF Hub requests
        model_name = "sentence-transformers/all-MiniLM-L6-v2"
        os.makedirs(_EMBEDDING_CACHE_DIR, exist_ok=True)
        cached_marker = os.path.join(_EMBEDDING_CACHE_DIR, "models--sentence-transformers--all-MiniLM-L6-v2")
        if os.path.isdir(cached_marker):
            os.environ.setdefault("HF_HUB_OFFLINE", "1")
            logger.info("RAG: Using cached embedding model (offline mode)")

        _encoder = SentenceTransformer(model_name, cache_folder=_EMBEDDING_CACHE_DIR)
        logger.info("RAG initialized: ChromaDB at %s", CHROMA_DB_PATH)
    return _collection, _encoder


# ---------------------------------------------------------------------------
# Multi-format text extraction
# ---------------------------------------------------------------------------

def _extract_text(content: bytes, filename: str) -> str:
    """Extract text from various file formats."""
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""

    # --- PDF ---
    if ext in _PDF_EXTENSIONS:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=content, filetype="pdf")
            extracted_text = "\n\n".join(page.get_text() for page in doc)
            
            # OCR Fallback: if very little text is found (e.g. scanned image), try OCR
            if len(extracted_text.strip()) < 50 * len(doc):
                try:
                    import pytesseract
                    from PIL import Image
                    import io
                    ocr_parts = []
                    for page in doc:
                        # Extract text first, if it's there
                        page_text = page.get_text().strip()
                        if len(page_text) < 50:
                            # Render page to an image
                            pix = page.get_pixmap(dpi=150)
                            img = Image.open(io.BytesIO(pix.tobytes("png")))
                            page_text = pytesseract.image_to_string(img)
                        ocr_parts.append(page_text)
                    extracted_text = "\n\n".join(ocr_parts)
                except Exception as ocr_e:
                    logger.warning("OCR fallback failed for PDF %s: %s", filename, ocr_e)
            
            return extracted_text
        except ImportError:
            try:
                import pdfplumber, io
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    return "\n\n".join(p.extract_text() or "" for p in pdf.pages)
            except ImportError:
                return content.decode("utf-8", errors="ignore")

    # --- DOCX (Word) ---
    if ext == ".docx":
        try:
            from docx import Document
            import io
            doc = Document(io.BytesIO(content))
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
            return "\n\n".join(parts)
        except Exception as e:
            logger.warning("DOCX extraction failed for %s: %s", filename, e)
            return content.decode("utf-8", errors="ignore")

    # --- PPTX (PowerPoint) ---
    if ext == ".pptx":
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"--- Slide {slide_num} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                slide_texts.append(row_text)
                if len(slide_texts) > 1:  # More than just the header
                    parts.append("\n".join(slide_texts))
            return "\n\n".join(parts)
        except Exception as e:
            logger.warning("PPTX extraction failed for %s: %s", filename, e)
            return content.decode("utf-8", errors="ignore")

    # --- XLSX (Excel) ---
    if ext == ".xlsx":
        try:
            import openpyxl
            import io
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_rows = []
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(c) if c is not None else "" for c in row]
                    if any(v.strip() for v in row_vals):
                        sheet_rows.append(" | ".join(row_vals))
                if sheet_rows:
                    parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(sheet_rows))
            wb.close()
            return "\n\n".join(parts)
        except Exception as e:
            logger.warning("XLSX extraction failed for %s: %s", filename, e)
            return content.decode("utf-8", errors="ignore")

    # --- CSV ---
    if ext == ".csv":
        try:
            text = content.decode("utf-8", errors="ignore")
            lines = text.strip().split("\n")
            # Keep header + rows as pipe-separated for better semantic search
            parts = []
            if lines:
                header = lines[0]
                parts.append(f"Columns: {header}")
                for line in lines[1:]:
                    parts.append(line)
            return "\n".join(parts)
        except Exception:
            return content.decode("utf-8", errors="ignore")

    # --- EPUB ---
    if ext == ".epub":
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            import tempfile
            import os

            with tempfile.NamedTemporaryFile(delete=False, suffix=".epub") as tmp:
                tmp.write(content)
                tmp_path = tmp.name

            try:
                book = epub.read_epub(tmp_path)
                parts = []
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        soup = BeautifulSoup(item.get_body_content(), 'html.parser')
                        text = soup.get_text(separator='\n', strip=True)
                        if text:
                            parts.append(text)
                return "\n\n".join(parts)
            finally:
                os.remove(tmp_path)
        except Exception as e:
            logger.warning("EPUB extraction failed for %s: %s", filename, e)
            return content.decode("utf-8", errors="ignore")

    # --- Code files ---
    if ext in _CODE_EXTENSIONS:
        text = content.decode("utf-8", errors="ignore")
        lang = ext.lstrip(".")
        return f"[Code: {filename} ({lang})]\n{text}"

    # --- Archives ---
    if ext == ".zip":
        import zipfile, io
        texts = []
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            for name in zf.namelist():
                member_ext = ("." + name.rsplit(".", 1)[-1].lower()) if "." in name else ""
                if member_ext in (_TEXT_EXTENSIONS | _CODE_EXTENSIONS | _DATA_EXTENSIONS):
                    member_content = zf.read(name)
                    texts.append(f"--- {name} ---\n" + _extract_text(member_content, name))
        return "\n\n".join(texts)

    if ext in (".tar", ".gz", ".tgz"):
        import tarfile, io
        texts = []
        with tarfile.open(fileobj=io.BytesIO(content)) as tf:
            for member in tf.getmembers():
                if member.isfile():
                    member_ext = ("." + member.name.rsplit(".", 1)[-1].lower()) if "." in member.name else ""
                    if member_ext in (_TEXT_EXTENSIONS | _CODE_EXTENSIONS | _DATA_EXTENSIONS):
                        f = tf.extractfile(member)
                        if f:
                            texts.append(f"--- {member.name} ---\n" + _extract_text(f.read(), member.name))
        return "\n\n".join(texts)

    # --- Default: plain text ---
    return content.decode("utf-8", errors="ignore")


def _chunk_text(text: str, chunk_size: int = 1000) -> list[str]:
    """Split text into semantically-aware chunks."""
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = ""
    for para in paragraphs:
        if len(current_chunk) + len(para) > chunk_size:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            current_chunk = para
        else:
            current_chunk += "\n\n" + para
    if current_chunk.strip():
        chunks.append(current_chunk.strip())
    return chunks


def _file_hash(content: bytes) -> str:
    """Compute MD5 hash for incremental indexing."""
    return hashlib.md5(content).hexdigest()


def _get_file_type(filename: str) -> str:
    """Return human-readable file type."""
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    if ext in _PDF_EXTENSIONS: return "PDF"
    if ext in _EBOOK_EXTENSIONS: return "eBook"
    if ext == ".docx": return "Word"
    if ext == ".pptx": return "PowerPoint"
    if ext == ".xlsx": return "Excel"
    if ext == ".csv": return "CSV"
    if ext == ".md": return "Markdown"
    if ext == ".txt": return "Text"
    if ext in _CODE_EXTENSIONS: return f"Code ({ext.lstrip('.')})"
    if ext in _ARCHIVE_EXTENSIONS: return "Archive"
    return "Text"


# ---------------------------------------------------------------------------
# Index a single document
# ---------------------------------------------------------------------------

def _index_document(content: bytes, filename: str, source_prefix: str = "") -> dict:
    """Index a document into the knowledge base. Returns stats dict."""
    collection, encoder = _get_rag()
    text = _extract_text(content, filename)
    chunks = _chunk_text(text)

    if not chunks:
        return {"status": "empty", "filename": filename, "chunks": 0}

    source_name = f"{source_prefix}{filename}" if source_prefix else filename
    file_type = _get_file_type(filename)

    for idx, chunk in enumerate(chunks):
        embedding = encoder.encode(chunk).tolist()
        doc_id = f"{source_name}_{idx}"
        collection.upsert(
            documents=[chunk],
            embeddings=[embedding],
            ids=[doc_id],
            metadatas=[{
                "source": source_name,
                "chunk_index": idx,
                "file_type": file_type,
                "indexed_at": datetime.now(timezone.utc).isoformat(),
            }],
        )

    logger.info("Indexed %s: %d chunks (%s)", source_name, len(chunks), file_type)
    return {"status": "indexed", "filename": source_name, "chunks": len(chunks), "file_type": file_type}


# ---------------------------------------------------------------------------
# RAG Folder scanning
# ---------------------------------------------------------------------------

def scan_rag_folder() -> dict:
    """Scan data/rag/ folder and index new/modified files.

    Returns a report: { added: [...], updated: [...], skipped: int, errors: [...] }
    """
    os.makedirs(RAG_DIR, exist_ok=True)
    report = {"added": [], "updated": [], "skipped": 0, "errors": [], "total_scanned": 0}

    # Get already indexed sources
    existing_sources = set()
    try:
        collection, _ = _get_rag()
        all_data = collection.get(include=["metadatas"])
        for m in (all_data.get("metadatas") or []):
            src = (m or {}).get("source", "")
            if src:
                existing_sources.add(src)
    except Exception:
        pass

    for root, dirs, files in os.walk(RAG_DIR):
        # Skip hidden directories
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fname in files:
            if fname.startswith("."):
                continue
            ext = ("." + fname.rsplit(".", 1)[-1].lower()) if "." in fname else ""
            if ext not in _ALL_SUPPORTED:
                continue

            filepath = os.path.join(root, fname)
            rel_path = os.path.relpath(filepath, RAG_DIR)
            source_name = f"rag/{rel_path}"
            report["total_scanned"] += 1

            try:
                with open(filepath, "rb") as f:
                    content = f.read()

                content_hash = _file_hash(content)

                # Check if already indexed with same hash
                if source_name in _indexed_hashes and _indexed_hashes[source_name] == content_hash:
                    report["skipped"] += 1
                    continue

                is_update = source_name in existing_sources

                # Delete old chunks if updating
                if is_update:
                    try:
                        _delete_source_chunks(source_name)
                    except Exception:
                        pass

                result = _index_document(content, fname, source_prefix="rag/")
                _indexed_hashes[source_name] = content_hash

                if is_update:
                    report["updated"].append({"file": rel_path, "chunks": result["chunks"]})
                else:
                    report["added"].append({"file": rel_path, "chunks": result["chunks"]})

            except Exception as e:
                logger.error("RAG scan error for %s: %s", filepath, e)
                report["errors"].append({"file": rel_path, "error": str(e)})

    total = len(report["added"]) + len(report["updated"])
    logger.info("RAG folder scan complete: %d new, %d updated, %d skipped, %d errors",
                len(report["added"]), len(report["updated"]), report["skipped"], len(report["errors"]))
    return report


def _delete_source_chunks(source_name: str):
    """Delete all chunks for a given source."""
    collection, _ = _get_rag()
    all_data = collection.get(include=["metadatas"])
    ids_to_delete = []
    for doc_id, meta in zip(all_data.get("ids", []), all_data.get("metadatas", [])):
        if (meta or {}).get("source") == source_name:
            ids_to_delete.append(doc_id)
    if ids_to_delete:
        collection.delete(ids=ids_to_delete)


# ---------------------------------------------------------------------------
# API Endpoints
# ---------------------------------------------------------------------------

@router.post("/add")
async def add_document(file: UploadFile = File(...)):
    """Upload and index a document into the knowledge base."""
    content = await file.read()
    fname = file.filename or ""
    result = _index_document(content, fname)
    if result["status"] == "empty":
        raise HTTPException(400, "Empty document")
    return result


@router.post("/upload-multi")
async def upload_multi(files: List[UploadFile] = File(...)):
    """Upload and index multiple documents at once."""
    results = []
    for file in files:
        try:
            content = await file.read()
            fname = file.filename or ""
            result = _index_document(content, fname)
            results.append(result)
        except Exception as e:
            results.append({"status": "error", "filename": file.filename, "error": str(e)})
    return {"results": results, "total": len(results), "indexed": sum(1 for r in results if r.get("status") == "indexed")}


@router.post("/scan")
async def scan_folder():
    """Scan the RAG folder and index new/modified files."""
    import asyncio
    report = await asyncio.to_thread(scan_rag_folder)
    return report


@router.get("/search")
async def search(query: str, k: int = 3, hybrid: bool = True):
    """Search the knowledge base using hybrid search (dense + BM25).

    When ``hybrid=True`` (default), combines ChromaDB dense vector search
    with BM25 keyword scoring using Reciprocal Rank Fusion (RRF).
    """
    collection, encoder = _get_rag()

    if not hybrid:
        # Pure dense search
        query_embedding = encoder.encode(query).tolist()
        results = collection.query(query_embeddings=[query_embedding], n_results=k)
        docs = results.get("documents", [[]])[0]
        metadatas = results.get("metadatas", [[]])[0]
        return {"documents": docs, "metadatas": metadatas, "method": "dense"}

    # --- Hybrid search: dense + BM25 with RRF ---
    # 1. Dense search (retrieve more candidates for fusion)
    n_candidates = min(k * 3, 20)
    query_embedding = encoder.encode(query).tolist()
    dense_results = collection.query(
        query_embeddings=[query_embedding], n_results=n_candidates,
        include=["documents", "metadatas"],
    )
    dense_docs = dense_results.get("documents", [[]])[0]
    dense_ids = dense_results.get("ids", [[]])[0]
    dense_metas = dense_results.get("metadatas", [[]])[0]

    if not dense_docs:
        return {"documents": [], "metadatas": [], "method": "hybrid"}

    # 2. BM25 scoring on the dense candidates
    try:
        from rank_bm25 import BM25Okapi
        tokenized_corpus = [doc.lower().split() for doc in dense_docs]
        bm25 = BM25Okapi(tokenized_corpus)
        bm25_scores = bm25.get_scores(query.lower().split())
    except ImportError:
        logger.warning("rank_bm25 not installed — falling back to dense-only search")
        return {"documents": dense_docs[:k], "metadatas": dense_metas[:k], "method": "dense_fallback"}

    # 3. Reciprocal Rank Fusion (RRF)
    rrf_k = 60  # standard RRF constant
    rrf_scores = {}
    for rank, doc_id in enumerate(dense_ids):
        rrf_scores[doc_id] = rrf_scores.get(doc_id, 0) + 1.0 / (rrf_k + rank + 1)

    bm25_ranking = sorted(range(len(bm25_scores)), key=lambda i: bm25_scores[i], reverse=True)
    for rank, idx in enumerate(bm25_ranking):
        doc_id = dense_ids[idx]
        rrf_scores[doc_id] = rrf_scores.get(doc_id, 0) + 1.0 / (rrf_k + rank + 1)

    # Sort by fused score
    sorted_ids = sorted(rrf_scores.keys(), key=lambda x: rrf_scores[x], reverse=True)[:k]

    # Map back to docs and metadata
    id_to_idx = {did: i for i, did in enumerate(dense_ids)}
    docs = [dense_docs[id_to_idx[did]] for did in sorted_ids if did in id_to_idx]
    metas = [dense_metas[id_to_idx[did]] for did in sorted_ids if did in id_to_idx]

    return {"documents": docs, "metadatas": metas, "method": "hybrid"}


@router.get("/stats")
async def rag_stats():
    """Return statistics about the knowledge base (does NOT load the embedding model)."""
    try:
        if _client is not None and _collection is not None:
            # Already initialized — use existing objects
            count = _collection.count()
            # Get unique sources
            all_meta = _collection.get(include=["metadatas"])
            sources = {}
            for m in (all_meta.get("metadatas") or []):
                src = (m or {}).get("source", "unknown")
                ftype = (m or {}).get("file_type", "")
                sources[src] = sources.get(src, 0) + 1
            return {"total_chunks": count, "sources": sources, "source_count": len(sources)}
        # Not yet initialized — open ChromaDB read-only without loading embeddings
        import chromadb
        from chromadb.config import Settings
        client = chromadb.PersistentClient(
            path=CHROMA_DB_PATH,
            settings=Settings(anonymized_telemetry=False),
        )
        col = client.get_or_create_collection("knowledge_base")
        return {"total_chunks": col.count()}
    except Exception:
        return {"total_chunks": 0, "status": "unavailable"}


@router.get("/folder-info")
async def folder_info():
    """Return info about the RAG folder (files present, indexed status)."""
    os.makedirs(RAG_DIR, exist_ok=True)
    files = []
    for root, dirs, fnames in os.walk(RAG_DIR):
        dirs[:] = [d for d in dirs if not d.startswith(".")]
        for fname in fnames:
            if fname.startswith("."):
                continue
            ext = ("." + fname.rsplit(".", 1)[-1].lower()) if "." in fname else ""
            if ext not in _ALL_SUPPORTED:
                continue
            filepath = os.path.join(root, fname)
            rel_path = os.path.relpath(filepath, RAG_DIR)
            source_name = f"rag/{rel_path}"
            stat = os.stat(filepath)
            files.append({
                "name": rel_path,
                "size": stat.st_size,
                "modified": datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc).isoformat(),
                "indexed": source_name in _indexed_hashes,
                "type": _get_file_type(fname),
            })
    return {"path": RAG_DIR, "files": files, "total": len(files)}


# --- RAG Management ---

@router.get("/sources")
async def list_sources():
    """List all indexed document sources with chunk counts and metadata."""
    collection, _ = _get_rag()
    all_data = collection.get(include=["metadatas"])
    sources: dict[str, dict] = {}
    for m in (all_data.get("metadatas") or []):
        src = (m or {}).get("source", "unknown")
        if src not in sources:
            sources[src] = {
                "name": src,
                "chunks": 0,
                "file_type": (m or {}).get("file_type", "unknown"),
                "indexed_at": (m or {}).get("indexed_at", ""),
            }
        sources[src]["chunks"] += 1
    return {"sources": sorted(sources.values(), key=lambda x: x["name"])}


@router.delete("/source/{source_name:path}")
async def delete_source(source_name: str):
    """Delete all chunks from a specific source document."""
    collection, _ = _get_rag()
    # Find all IDs for this source
    all_data = collection.get(include=["metadatas"])
    ids_to_delete = []
    for doc_id, meta in zip(all_data.get("ids", []), all_data.get("metadatas", [])):
        if (meta or {}).get("source") == source_name:
            ids_to_delete.append(doc_id)

    if not ids_to_delete:
        raise HTTPException(404, f"Source '{source_name}' not found")

    collection.delete(ids=ids_to_delete)

    # Remove from hash cache
    _indexed_hashes.pop(source_name, None)

    logger.info("Deleted %d chunks from source '%s'", len(ids_to_delete), source_name)
    return {"status": "deleted", "source": source_name, "chunks_removed": len(ids_to_delete)}


@router.delete("/clear")
async def clear_all():
    """Clear the entire knowledge base."""
    collection, _ = _get_rag()
    # Get all IDs and delete
    all_data = collection.get()
    all_ids = all_data.get("ids", [])
    if all_ids:
        collection.delete(ids=all_ids)
    _indexed_hashes.clear()
    logger.info("Cleared entire knowledge base (%d chunks)", len(all_ids))
    return {"status": "cleared", "chunks_removed": len(all_ids)}


# --- Auto-RAG Context Injection ---

def auto_rag_context(user_message: str, threshold: float = 0.3, k: int = 3) -> str | None:
    """Check if RAG has relevant context for a user message.

    Returns a formatted context string to inject into the LLM prompt,
    or None if no relevant documents are found.
    Only triggers if the knowledge base is initialized and non-empty.
    """
    try:
        # Initialize ChromaDB first, THEN check count
        collection, encoder = _get_rag()
        if collection.count() == 0:
            return None

        query_embedding = encoder.encode(user_message).tolist()
        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=k,
            include=["documents", "metadatas", "distances"],
        )

        docs = results.get("documents", [[]])[0]
        distances = results.get("distances", [[]])[0]
        metadatas = results.get("metadatas", [[]])[0]

        if not docs:
            return None

        # Filter by distance threshold (lower = more similar for cosine)
        relevant = []
        for doc, dist, meta in zip(docs, distances, metadatas):
            if dist < threshold:
                source = (meta or {}).get("source", "?")
                file_type = (meta or {}).get("file_type", "")
                type_tag = f" ({file_type})" if file_type else ""
                relevant.append(f"[Source: {source}{type_tag}]\n{doc}")

        if not relevant:
            return None

        context = "\n\n---\n\n".join(relevant)
        return (
            "📚 **Relevant context from your knowledge base:**\n\n"
            f"{context}\n\n---\n\n"
            "Use the above context to inform your answer if relevant."
        )
    except Exception as e:
        logger.debug("Auto-RAG skipped: %s", e)
        return None


def explicit_rag_search(query: str, k: int = 5) -> str | None:
    """Perform an explicit RAG search with more results and no threshold filtering.

    Used when the user explicitly requests RAG search (via @rag or RAG pill).
    Returns formatted context string or None.
    """
    try:
        # Initialize ChromaDB first, THEN check count
        collection, encoder = _get_rag()
        if collection.count() == 0:
            return "📚 **Knowledge base is empty.** Upload documents in Settings → Knowledge Base, or place files in the `data/rag/` folder."

        query_embedding = encoder.encode(query).tolist()
        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=k,
            include=["documents", "metadatas", "distances"],
        )

        docs = results.get("documents", [[]])[0]
        distances = results.get("distances", [[]])[0]
        metadatas = results.get("metadatas", [[]])[0]

        if not docs:
            return "📚 **No relevant documents found in the knowledge base.**"

        parts = []
        for i, (doc, dist, meta) in enumerate(zip(docs, distances, metadatas), 1):
            source = (meta or {}).get("source", "?")
            file_type = (meta or {}).get("file_type", "")
            type_tag = f" ({file_type})" if file_type else ""
            score = f"{1.0 - dist:.0%}" if dist < 1.0 else "N/A"
            parts.append(f"### Result {i} — {source}{type_tag} [relevance: {score}]\n{doc}")

        context = "\n\n---\n\n".join(parts)
        return (
            f"📚 **Knowledge Base Search Results ({len(parts)} results for: \"{query}\"):**\n\n"
            f"{context}\n\n---\n\n"
            "Use the above documents to answer the user's question comprehensively. "
            "Cite sources when possible."
        )
    except Exception as e:
        logger.debug("Explicit RAG search failed: %s", e)
        return f"📚 **RAG search error:** {e}"