"""
Clawzd — Research Studio module.
Autonomous research bot with iterative improvement, auto-evaluation,
web scraping, asset downloading, and multi-format report export.
V2: Multi-phase process, editable Markdown process, research profiles,
sandbox execution with venv, chat-to-research integration.
"""
import os, json, uuid, logging, asyncio, hashlib, shutil, subprocess, sys
from datetime import datetime, timezone
from fastapi import APIRouter, Request, HTTPException
from fastapi.responses import FileResponse, StreamingResponse
from sse_starlette.sse import EventSourceResponse
from config import DATA_DIR
from app.tools.research_profiles import (
    list_profiles, get_profile, create_profile, delete_profile,
    render_process_md, DEFAULT_PROCESS_MD,
)
from app.tools.research_engine import (
    generate_perspectives, generate_sub_questions,
    deep_research_branch, flatten_branch_results,
    flatten_branch_summaries, evaluate_structured,
    reflect_on_iteration, generate_report_with_citations,
    improve_process_md,
    research_by_perspectives_parallel, synthesize_perspective_branches,
    update_dynamic_outline, outline_to_report_structure,
    reflect_after_search, should_stop_early,
    generate_sectioned_report,
    sanitize_mermaid_in_report, merge_iteration_findings,
)
from app.tools.research_brief import (
    generate_research_brief, brief_to_planning_context,
)
from app.tools.research_scraper import batch_scrape, smart_scrape
from app.tools.research_archive import (
    classify_query_domain, save_strategy, get_best_strategy_for_domain,
    get_archive_stats,
)
from app.tools.research_condensation import (
    should_condense, condense_research_context,
)
from app.tools.research_progress import (
    ResearchProgress, ResearchPhase, ResearchCostTracker,
)
from app.core.tokens import count_tokens

router = APIRouter()
logger = logging.getLogger("clawzd.research")

RESEARCH_DIR = os.path.join(DATA_DIR, "research")
os.makedirs(RESEARCH_DIR, exist_ok=True)

# In-memory state for running research tasks
_running: dict[str, asyncio.Task] = {}
_sse_queues: dict[str, asyncio.Queue] = {}
_stop_requested: set[str] = set()  # Clean pause signaling (Bug 3)


def _fix_orphan_running_statuses():
    """Fix projects stuck in 'running' status after a server crash (Bug 6).

    Called once at import time. If project.json says 'running' but the
    process just started, no asyncio task can be alive — mark as 'paused'.
    """
    if not os.path.isdir(RESEARCH_DIR):
        return
    for name in os.listdir(RESEARCH_DIR):
        pf = os.path.join(RESEARCH_DIR, name, "project.json")
        if not os.path.isfile(pf):
            continue
        try:
            with open(pf) as f:
                d = json.load(f)
            if d.get("status") == "running":
                d["status"] = "paused"
                d["updated_at"] = datetime.now(timezone.utc).isoformat()
                with open(pf, "w") as f:
                    json.dump(d, f, indent=2, ensure_ascii=False)
                logger.info(
                    "Fixed orphan running status for project %s → paused",
                    d.get("id", name),
                )
        except Exception:
            pass


_fix_orphan_running_statuses()


def _get_dev_profile_summary() -> str:
    """Load a condensed dev best practices summary for research script generation."""
    try:
        from app.preprompts import _load_dev_profile
        profile = _load_dev_profile()
        if profile:
            # Return a condensed version to save tokens in the planning prompt
            return profile[:1500]
    except Exception:
        pass
    return ""


# ── Helpers ──

def _proj_dir(pid: str) -> str:
    d = os.path.join(RESEARCH_DIR, pid)
    os.makedirs(d, exist_ok=True)
    os.makedirs(os.path.join(d, "assets"), exist_ok=True)
    return d

def _load(pid: str) -> dict | None:
    p = os.path.join(_proj_dir(pid), "project.json")
    if os.path.exists(p):
        with open(p) as f:
            return json.load(f)
    return None

def _save(proj: dict):
    proj["updated_at"] = datetime.now(timezone.utc).isoformat()
    # Bug 12: Strip full_text from search_results to prevent multi-MB JSON files.
    # The full content is already saved separately as text assets.
    results = proj.get("search_results", [])
    for r in results:
        r.pop("full_text", None)
        r.pop("raw_text", None)
    with open(os.path.join(_proj_dir(proj["id"]), "project.json"), "w") as f:
        json.dump(proj, f, indent=2, ensure_ascii=False)

def _new_project(
    query: str, provider: str = "", model: str = "",
    profile_id: str = "", profile_data: dict | None = None,
) -> dict:
    now = datetime.now(timezone.utc).isoformat()
    pid = uuid.uuid4().hex[:8]
    target_score = 0.7
    max_iterations = 10
    sources = "Web"
    if profile_data:
        target_score = profile_data.get("target_score", target_score)
        max_iterations = profile_data.get("max_iterations", max_iterations)
        sources = profile_data.get("sources", sources)
    proj = {
        "id": pid,
        "title": query[:80],
        "query": query,
        "status": "idle",
        "created_at": now,
        "updated_at": now,
        "iterations": [],
        "current_score": 0.0,
        "target_score": target_score,
        "max_iterations": max_iterations,
        "report_md": "",
        "assets": [],
        "search_results": [],
        "provider": provider,
        "model": model,
        "profile_id": profile_id,
        "phase_models": profile_data.get("phase_models", {}) if profile_data else {},
    }
    # Generate initial process.md
    tmpl = (profile_data or {}).get("process_template", DEFAULT_PROCESS_MD)
    process_md = render_process_md(
        tmpl, query=query, model=model or "System Default",
        provider=provider or "System Default", sources=sources,
        target_score=target_score, max_iterations=max_iterations,
    )
    pdir = _proj_dir(pid)
    with open(os.path.join(pdir, "process.md"), "w", encoding="utf-8") as f:
        f.write(process_md)

    # ── Archive Warm-Start (HyperAgents-inspired) ──────────────────────────
    # Try to bootstrap from the best known strategy for this query's domain
    domain = classify_query_domain(query)
    proj["query_domain"] = domain
    prior = get_best_strategy_for_domain(domain, min_score=0.65)
    if prior:
        proj["warm_start"] = {
            "domain": domain,
            "prior_score": prior.get("final_score", 0),
            "prior_actions": prior.get("action_sequence", []),
        }
        logger.info(
            "Warm-start from archive: domain=%s prior_score=%.0f%%",
            domain, prior.get("final_score", 0) * 100,
        )
    # ───────────────────────────────────────────────────────────────────────

    return proj

def _list_projects() -> list:
    projects = []
    if not os.path.isdir(RESEARCH_DIR):
        return projects
    for name in os.listdir(RESEARCH_DIR):
        pf = os.path.join(RESEARCH_DIR, name, "project.json")
        if os.path.isfile(pf):
            try:
                with open(pf) as f:
                    d = json.load(f)
                pid = d["id"]
                status = d.get("status", "idle")
                # Bug 16: Cross-check with _running to fix phantom "running" statuses
                if status == "running":
                    task = _running.get(pid)
                    if task is None or task.done():
                        # No live task — correct the status on disk
                        status = "paused"
                        d["status"] = status
                        d["updated_at"] = datetime.now(timezone.utc).isoformat()
                        with open(pf, "w") as fw:
                            json.dump(d, fw, indent=2, ensure_ascii=False)
                        _running.pop(pid, None)
                projects.append({
                    "id": pid, "title": d.get("title", ""),
                    "status": status,
                    "current_score": d.get("current_score", 0),
                    "iteration_count": len(d.get("iterations", [])),
                    "updated_at": d.get("updated_at", ""),
                    "profile_id": d.get("profile_id", ""),
                })
            except Exception:
                pass
    return sorted(projects, key=lambda x: x.get("updated_at", ""), reverse=True)


def _read_process(pid: str) -> str:
    """Read the process.md for a project."""
    path = os.path.join(_proj_dir(pid), "process.md")
    if os.path.isfile(path):
        with open(path, encoding="utf-8") as f:
            return f.read()
    return ""


def _write_process(pid: str, content: str):
    """Write the process.md for a project."""
    path = os.path.join(_proj_dir(pid), "process.md")
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)


def _ensure_venv(pid: str) -> str:
    """Create a Python venv for a project sandbox if it doesn't exist."""
    venv_dir = os.path.join(_proj_dir(pid), ".venv")
    if not os.path.isdir(venv_dir):
        subprocess.run(
            [sys.executable, "-m", "venv", venv_dir],
            capture_output=True, timeout=60,
        )
    return venv_dir


def _venv_python(pid: str) -> str:
    """Return path to the venv's python executable."""
    venv = _ensure_venv(pid)
    return os.path.join(venv, "bin", "python")


async def _emit(pid: str, event_type: str, data: dict):
    q = _sse_queues.get(pid)
    if q:
        await q.put(json.dumps({"type": event_type, **data}))


# ── Research Loop Engine ──

async def _do_web_search(query: str, max_results: int = 50) -> list[dict]:
    """Search using Tavily, DuckDuckGo, Google Scholar, Reddit, Twitter/X, and News in parallel."""
    from config import TAVILY_API_KEY
    from app.tools_web import _scrape_scholar

    async def _tavily_search() -> list[dict]:
        if not TAVILY_API_KEY:
            return []
        try:
            from tavily import AsyncTavilyClient
            client = AsyncTavilyClient(api_key=TAVILY_API_KEY)
            response = await client.search(query, max_results=max_results)
            return [
                {"title": r.get("title", ""), "snippet": r.get("content", ""),
                 "url": r.get("url", ""), "score": r.get("score", 0), "source": "tavily"}
                for r in response.get("results", [])
            ]
        except Exception as e:
            logger.warning("Tavily search failed: %s", e)
            return []

    async def _ddg_search() -> list[dict]:
        try:
            from ddgs import DDGS
            results = await asyncio.to_thread(
                lambda: list(DDGS().text(query, max_results=max_results))
            )
            return [{"title": r.get("title",""), "snippet": r.get("body",""),
                     "url": r.get("href",""), "source": "duckduckgo"} for r in results]
        except Exception as e:
            logger.warning("DuckDuckGo search failed: %s", e)
            return []

    async def _scholar_search() -> list[dict]:
        try:
            res = await asyncio.to_thread(_scrape_scholar, query, max_results)
            res_labs = await asyncio.to_thread(
                _scrape_scholar, query, max_results,
                "https://scholar.google.com/scholar_labs/search",
            )
            combined = res + res_labs
            return [{"title": r.get("title",""), "snippet": r.get("snippet",""),
                     "url": r.get("url",""), "source": "scholar"} for r in combined]
        except Exception as e:
            logger.warning("Google Scholar search failed: %s", e)
            return []

    async def _reddit_search() -> list[dict]:
        """Search Reddit discussions via DuckDuckGo site-scoped query."""
        try:
            from ddgs import DDGS
            reddit_query = f"site:reddit.com {query}"
            results = await asyncio.to_thread(
                lambda: list(DDGS().text(reddit_query, max_results=min(max_results, 20)))
            )
            return [{"title": f"[Reddit] {r.get('title','')}", "snippet": r.get("body",""),
                     "url": r.get("href",""), "source": "reddit"} for r in results]
        except Exception as e:
            logger.warning("Reddit search failed: %s", e)
            return []

    async def _twitter_search() -> list[dict]:
        """Search Twitter/X using the existing tools_twitter module."""
        try:
            from app.tools_twitter import search_tweets
            tweets = await search_tweets(query, max_results=min(max_results, 20))
            return [{"title": f"[X] @{t.get('author_username','unknown')}: {t.get('text','')[:80]}",
                     "snippet": t.get("text", ""),
                     "url": t.get("url", ""),
                     "source": "twitter"} for t in tweets]
        except Exception as e:
            logger.warning("Twitter/X search failed: %s", e)
            return []

    async def _news_search() -> list[dict]:
        """Search news articles via DuckDuckGo News API."""
        try:
            from ddgs import DDGS
            results = await asyncio.to_thread(
                lambda: list(DDGS().news(query, max_results=min(max_results, 25)))
            )
            return [{"title": f"[News] {r.get('title','')}", "snippet": r.get("body",""),
                     "url": r.get("url",""), "source": "news",
                     "date": r.get("date", "")} for r in results]
        except Exception as e:
            logger.warning("News search failed: %s", e)
            return []

    # Run ALL sources in parallel
    (tavily_results, ddg_results, scholar_results,
     reddit_results, twitter_results, news_results) = await asyncio.gather(
        _tavily_search(), _ddg_search(), _scholar_search(),
        _reddit_search(), _twitter_search(), _news_search(),
    )

    # Merge & deduplicate by URL
    # Priority order: Tavily > News > Twitter > Reddit > Scholar > DDG
    seen_urls = set()
    merged = []
    for r in (tavily_results + news_results + twitter_results
              + reddit_results + scholar_results + ddg_results):
        url = r.get("url", "")
        if url and url not in seen_urls:
            seen_urls.add(url)
            merged.append(r)
    return merged


async def _do_web_search_with_retry(
    query: str, max_results: int = 50, max_retries: int = 2,
) -> list[dict]:
    """
    Search wrapper with retry + fallback (DeepResearch feature #5).
    If the primary search returns empty results, retries with a
    simplified query before giving up.
    """
    for attempt in range(max_retries + 1):
        try:
            results = await _do_web_search(query, max_results)
            if results:
                return results
            if attempt < max_retries:
                # Simplify query: remove quotes and special operators
                import re as _re
                simpler = _re.sub(r'["\(\)\+\-]', '', query).strip()
                logger.info("Search retry %d with simpler query: %s", attempt + 1, simpler[:60])
                query = simpler
        except Exception as e:
            if attempt < max_retries:
                logger.warning("Search attempt %d failed (%s), retrying...", attempt + 1, e)
                await asyncio.sleep(1)
            else:
                logger.error("Search exhausted retries: %s", e)
    return []


async def _scrape_url(url: str) -> str:
    try:
        import httpx
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
            resp = await client.get(url, headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) Chrome/120.0.0.0"
            })
            if resp.status_code != 200:
                return ""
            ct = resp.headers.get("content-type", "")
            if "html" in ct:
                from html.parser import HTMLParser
                class TextExtractor(HTMLParser):
                    def __init__(self):
                        super().__init__()
                        self.text = []
                        self._skip = False
                    def handle_starttag(self, tag, attrs):
                        if tag in ("script","style","nav","footer","header"):
                            self._skip = True
                    def handle_endtag(self, tag):
                        if tag in ("script","style","nav","footer","header"):
                            self._skip = False
                    def handle_data(self, data):
                        if not self._skip and data.strip():
                            # Append with a space to prevent words from sticking together
                            self.text.append(data.strip() + " ")
                parser = TextExtractor()
                parser.feed(resp.text)
                return "".join(parser.text).replace("  ", " ")[:8000]
            return resp.text[:8000]
    except Exception as e:
        logger.warning("Scrape failed for %s: %s", url, e)
        return ""


async def _download_asset(url: str, pid: str) -> dict | None:
    try:
        import httpx
        async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
            resp = await client.get(url)
            if resp.status_code != 200:
                return None
            ct = resp.headers.get("content-type", "")
            ext = ".bin"
            if "pdf" in ct: ext = ".pdf"
            elif "png" in ct: ext = ".png"
            elif "jpeg" in ct or "jpg" in ct: ext = ".jpg"
            elif "gif" in ct: ext = ".gif"
            elif "svg" in ct: ext = ".svg"
            elif "webp" in ct: ext = ".webp"
            name = f"asset_{uuid.uuid4().hex[:6]}{ext}"
            path = os.path.join(_proj_dir(pid), "assets", name)
            with open(path, "wb") as f:
                f.write(resp.content)
            return {"name": name, "path": path, "type": ext.lstrip("."),
                    "url": url, "size": len(resp.content)}
    except Exception as e:
        logger.warning("Download failed for %s: %s", url, e)
        return None


async def _save_text_asset(title: str, text: str, source_type: str, url_or_ref: str, pid: str) -> dict:
    name = f"asset_{uuid.uuid4().hex[:6]}.md"
    path = os.path.join(_proj_dir(pid), "assets", name)
    content = f"# {title}\n\n**Source**: {url_or_ref}\n**Type**: {source_type}\n\n{text}"
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    return {"name": name, "path": path, "type": "md", "url": url_or_ref, "size": len(content)}


# ── Feature #6: Semantic Scholar Academic Search ──────────────────────────────

async def _semantic_scholar_search(query: str, max_results: int = 10) -> list[dict]:
    """
    Search Semantic Scholar (free API, no key required) for academic papers.

    Inspired by DeepResearch's tool_scholar.py — treats academic results
    differently from web results: preserves DOI, abstract, citation count,
    venue, and year as structured metadata for reliable academic citations.

    Returns result dicts compatible with the main search_results store.
    """
    try:
        import httpx
        url = "https://api.semanticscholar.org/graph/v1/paper/search"
        params = {
            "query": query,
            "limit": min(max_results, 20),
            "fields": "title,abstract,year,authors,citationCount,externalIds,venue,url",
        }
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(url, params=params, headers={
                "User-Agent": "Clawzd-Research/1.0 (academic search)"
            })
            if resp.status_code != 200:
                logger.warning("Semantic Scholar API returned %d", resp.status_code)
                return []

            data = resp.json()
            papers = data.get("data", [])
            results = []
            for p in papers:
                doi = (p.get("externalIds") or {}).get("DOI", "")
                paper_url = p.get("url", "")
                if doi and not paper_url:
                    paper_url = f"https://doi.org/{doi}"

                authors = ", ".join(
                    a.get("name", "") for a in (p.get("authors") or [])[:3]
                )
                year = p.get("year", "")
                citations = p.get("citationCount", 0)
                venue = p.get("venue", "")
                abstract = (p.get("abstract") or "")[:400]

                snippet = (
                    f"{abstract} "
                    f"[{authors}{', ' + str(year) if year else ''}. "
                    f"{venue + '. ' if venue else ''}"
                    f"Cited by {citations}]"
                ).strip()

                results.append({
                    "title": p.get("title", "Untitled"),
                    "snippet": snippet,
                    "url": paper_url or f"https://www.semanticscholar.org/search?q={query}",
                    "source": "semantic_scholar",
                    "doi": doi,
                    "year": year,
                    "citations": citations,
                    "venue": venue,
                    "relevance_score": min(1.0, (citations or 0) / 1000 + 0.3),
                })

            logger.info(
                "Semantic Scholar: %d papers for '%s'",
                len(results), query[:60],
            )
            return results

    except Exception as e:
        logger.warning("Semantic Scholar search failed: %s — falling back to Scholar scrape", e)
        try:
            from app.tools_web import _scrape_scholar
            raw = await asyncio.to_thread(_scrape_scholar, query, max_results)
            return [{
                "title": r.get("title", ""),
                "snippet": r.get("snippet", ""),
                "url": r.get("url", ""),
                "source": "scholar_fallback",
                "relevance_score": 0.5,
            } for r in raw]
        except Exception:
            return []


# ── Feature #7: Context Budget Monitor ───────────────────────────────────────

def _log_context_budget(
    results: list[dict],
    report_draft: str,
    branch_summaries: str,
    iteration: int,
) -> None:
    """Log estimated context size; warns when approaching the budget ceiling."""
    from app.tools.research_condensation import (
        estimate_context_size, CONTEXT_BUDGET_TOKENS,
    )
    estimated = estimate_context_size(results, report_draft, branch_summaries)
    pct = estimated / CONTEXT_BUDGET_TOKENS * 100
    if pct > 85:
        logger.warning(
            "Context budget (iter %d): ~%d tokens (%.0f%% — near ceiling %d)",
            iteration, estimated, pct, CONTEXT_BUDGET_TOKENS,
        )
    else:
        logger.debug(
            "Context budget (iter %d): ~%d tokens (%.0f%%)",
            iteration, estimated, pct,
        )


async def _llm_call(messages: list[dict], provider: str = "", model: str = "", pid: str = "") -> str:

    from app.llm_provider import get_llm_provider
    from app.metrics import get_metrics
    import time
    prov = get_llm_provider(provider or None)
    kwargs = {}
    if model:
        kwargs["model"] = model

    # Estimate input tokens using tiktoken
    input_text = "".join(m.get("content", "") for m in messages)
    input_tokens = max(1, count_tokens(input_text, model or ""))

    result = ""
    t0 = time.time()
    async for token in prov.chat_stream(messages, **kwargs):
        result += token
    elapsed = time.time() - t0

    # Estimate output tokens using tiktoken
    output_tokens = max(1, count_tokens(result, model or ""))

    # Record in metrics
    get_metrics().record_llm_call(
        provider=provider or "default",
        model=model or "default",
        input_tokens=input_tokens,
        output_tokens=output_tokens,
        latency_s=elapsed,
        session_id=pid or "",
    )

    # Emit token usage via SSE if a project is active
    if pid:
        await _emit(pid, "token_usage", {
            "input_tokens": input_tokens,
            "output_tokens": output_tokens,
            "total_tokens": input_tokens + output_tokens,
            "latency_s": round(elapsed, 2),
        })

    return result


async def _research_loop(pid: str):
    from app.tools.task_manager import register_task, unregister_task
    from config import (
        RESEARCH_SUMMARIZATION_MODEL, RESEARCH_MAIN_MODEL,
        RESEARCH_COMPRESSION_MODEL, RESEARCH_REPORT_MODEL,
    )
    proj = _load(pid)
    if not proj:
        return
    proj["status"] = "running"
    _save(proj)
    await _emit(pid, "status", {"status": "running"})
    register_task(pid, "research", proj.get("title", pid)[:60])

    provider = proj.get("provider", "")
    model = proj.get("model", "")
    query = proj["query"]

    # Use differentiated models per pipeline role (open_deep_research-inspired)
    main_model = RESEARCH_MAIN_MODEL if RESEARCH_MAIN_MODEL != model else model
    compression_model = RESEARCH_COMPRESSION_MODEL if RESEARCH_COMPRESSION_MODEL != model else model
    report_model = RESEARCH_REPORT_MODEL if RESEARCH_REPORT_MODEL != model else model

    # ── GPT-Researcher-inspired: structured progress + cost tracking ──────────
    _progress = ResearchProgress(
        total_depth=2,
        total_breadth=5,
        max_iterations=proj.get("max_iterations", 10),
        target_score=proj.get("target_score", 0.7),
    )
    _cost_tracker = ResearchCostTracker()

    async def _emit_progress(phase: str = "", label: str = "", **extra):
        """Emit a structured progress event via SSE for precise UI progress bars."""
        if phase:
            _progress.set_phase(phase, label)
        progress_data = _progress.to_dict()
        progress_data["cost"] = _cost_tracker.get_summary()
        progress_data.update(extra)
        await _emit(pid, "progress", progress_data)
    # ─────────────────────────────────────────────────────────────────────────


    from app import tools_project as pm
    pm_pid = proj.get("pm_project_id")
    if not pm_pid:
        pm_proj = pm._new_project(name=f"Research: {query[:40]}", description=f"Automated tracking for research on: {query}")
        now = datetime.now(timezone.utc).isoformat()
        pm_proj["tasks"] = [
            {"id": "t1", "title": "Perspective Decomposition", "status": "To Do", "progress": 0, "order": 0, "created_at": now},
            {"id": "t2", "title": "Data Gathering & Search", "status": "To Do", "progress": 0, "order": 1, "created_at": now},
            {"id": "t3", "title": "Multi-criteria Evaluation", "status": "To Do", "progress": 0, "order": 2, "created_at": now},
            {"id": "t4", "title": "Final Report Generation", "status": "To Do", "progress": 0, "order": 3, "created_at": now},
        ]
        pm._save_proj(pm_proj)
        proj["pm_project_id"] = pm_proj["id"]
        _save(proj)
        pm_pid = pm_proj["id"]

    def _update_pm_task(idx, status, progress):
        try:
            p = pm._load_proj(pm_pid)
            if p and len(p.get("tasks", [])) > idx:
                p["tasks"][idx]["status"] = status
                p["tasks"][idx]["progress"] = progress
                pm._save_proj(p)
        except Exception:
            pass
    # -----------------------------------

    # Wrap _llm_call: binds project ID + feeds cost tracker automatically
    async def _project_llm_call(messages, prov="", mdl=""):
        result = await _llm_call(messages, prov, mdl, pid=pid)
        try:
            from app.core.tokens import count_tokens as _ct
            in_tok = max(1, _ct("" .join(m.get("content", "") for m in messages), mdl or ""))
            out_tok = max(1, _ct(result, mdl or ""))
            _cost_tracker.add_call(in_tok, out_tok, model=mdl or model, provider=prov or provider)
            _progress.total_input_tokens += in_tok
            _progress.total_output_tokens += out_tok
            _progress.total_cost_usd = _cost_tracker.get_total_usd()
        except Exception:
            pass
        return result

    # ── Research Brief Generation (open_deep_research write_research_brief) ──
    # Transform raw query into a structured brief BEFORE any searching begins.
    # This improves the relevance of all downstream steps.
    research_brief_data: dict = {}
    try:
        await _emit_progress(phase=ResearchPhase.BRIEF)
        await _emit(pid, "log", {"msg": "📋 Generating structured research brief..."})
        research_brief_data = await generate_research_brief(
            query=query,
            llm_call=_project_llm_call,
            provider=provider,
            model=main_model or model,
        )
        if research_brief_data.get("research_brief") and research_brief_data["research_brief"] != query:
            brief_preview = research_brief_data['research_brief']
            await _emit(pid, "log", {
                "msg": (
                    f"📋 Brief: {brief_preview[:120]}..."
                    if len(brief_preview) > 120
                    else f"📋 Brief: {brief_preview}"
                )
            })
            if research_brief_data.get("key_dimensions"):
                dims = ', '.join(research_brief_data['key_dimensions'][:4])
                await _emit(pid, "log", {"msg": f"   Dimensions: {dims}"})
        proj["research_brief"] = research_brief_data
        _save(proj)
    except Exception as _be:
        logger.warning("Research brief generation failed: %s", _be)
    # ─────────────────────────────────────────────────────────────────────────


    perspectives = []
    sub_questions = []
    uncovered_questions = []
    branch_summaries = ""
    last_eval = {}

    # ── New DeepResearch state tracking ────────────────────────────────────
    dynamic_outline: dict = {}          # WebWeaver-style evolving outline
    report_draft: str = ""              # IterResearch evolving report draft
    perspective_synthesis: str = ""     # Multi-agent synthesis result
    cumulative_findings: str = proj.get("cumulative_findings", "")  # Iteration merging
    # ───────────────────────────────────────────────────────────────────────

    async def _emit_log(msg: str, extra: dict | None = None):
        payload = {"msg": msg}
        if extra:
            payload.update(extra)
        await _emit(pid, "log", payload)

    async def _emit_new_results(results: list[dict]):
        """Emit newly found search results so the frontend can update in real-time."""
        # Send lightweight versions (no full_text) to keep SSE messages small
        lite = []
        for r in results:
            lite.append({
                "title": r.get("title", "")[:120],
                "snippet": r.get("snippet", "")[:300],
                "url": r.get("url", ""),
                "source": r.get("source", "web"),
                "score": r.get("score", 0),
            })
        if lite:
            await _emit(pid, "new_results", {"results": lite})

    async def _emit_new_asset(asset: dict):
        """Emit a newly downloaded/saved asset so the frontend can update in real-time."""
        await _emit(pid, "new_asset", {
            "name": asset.get("name", ""),
            "type": asset.get("type", ""),
            "size": asset.get("size", 0),
            "url": asset.get("url", ""),
        })

    try:
        # ── Phase 0: Perspective Decomposition (STORM-style) ──
        _update_pm_task(0, "In Progress", 50)
        await _emit_progress(phase=ResearchPhase.PERSPECTIVES)
        await _emit_log("🔭 Generating research perspectives...")
        perspectives = await generate_perspectives(
            query, _project_llm_call, provider, model,
        )
        if perspectives:
            proj["perspectives"] = [
                {"perspective": p["perspective"], "description": p["description"]}
                for p in perspectives
            ]
            _save(proj)
            names = [p["perspective"] for p in perspectives]
            await _emit_log(f"📐 {len(perspectives)} perspectives: {', '.join(names)}")

            sub_questions = await generate_sub_questions(
                query, perspectives, _project_llm_call, provider, model,
            )
            uncovered_questions = list(sub_questions)
            await _emit_log(f"❓ {len(sub_questions)} sub-questions generated")
            _update_pm_task(0, "Done", 100)

            # ── Parallel Perspective Research (DeepResearch #2) ──────────────
            # For deep_research profile: launch all perspective branches in parallel
            # immediately after decomposition (Research-Synthesis paradigm)
            profile_id = proj.get("profile_id", "")
            if profile_id == "deep_research" and len(perspectives) >= 2:
                await _emit_log(
                    f"🔬 Launching {len(perspectives)} parallel perspective branches "
                    "(DeepResearch multi-agent mode)..."
                )
                try:
                    persp_branches = await research_by_perspectives_parallel(
                        query=query,
                        perspectives=perspectives,
                        search_fn=_do_web_search_with_retry,
                        scrape_fn=_scrape_url,
                        llm_call=_project_llm_call,
                        provider=provider,
                        model=model,
                        depth=1,
                        breadth=2,
                        emit_fn=_emit_log,
                    )
                    # Collect all results from perspective branches
                    for b in persp_branches:
                        branch_res = flatten_branch_results(b)
                        proj["search_results"].extend(branch_res)
                        await _emit_new_results(branch_res)
                        branch_summaries += "\n\n" + flatten_branch_summaries(b)

                    # Synthesis Agent step
                    if persp_branches:
                        perspective_synthesis = await synthesize_perspective_branches(
                            query=query,
                            branches=persp_branches,
                            perspectives=perspectives,
                            llm_call=_project_llm_call,
                            provider=provider,
                            model=model,
                        )
                        await _emit_log(
                            f"🧬 Perspective synthesis complete "
                            f"({len(perspective_synthesis)} chars)"
                        )
                        _save(proj)
                except Exception as _pe:
                    logger.warning("Parallel perspective research failed: %s", _pe)
                    await _emit_log(f"⚠️ Parallel research partial failure: {_pe}")
            # ────────────────────────────────────────────────────────────────

        # ── Iteration Loop ──
        for iteration in range(len(proj["iterations"]), proj["max_iterations"]):
            if pid in _stop_requested:
                _stop_requested.discard(pid)
                proj["status"] = "paused"
                _save(proj)
                await _emit(pid, "status", {"status": "paused"})
                return

            iter_num = iteration + 1
            iter_data = {
                "num": iter_num,
                "started_at": datetime.now(timezone.utc).isoformat(),
                "actions": [], "score": 0, "evaluation": "",
                "scores_detail": {},
            }
            _progress.current_iteration = iter_num
            _progress.search_results_count = len(proj.get("search_results", []))
            await _emit_progress(phase=ResearchPhase.SEARCH, current_query=query)
            await _emit(pid, "iteration_start", {"iteration": iter_num})
            _update_pm_task(1, "In Progress", int((iter_num - 1) / proj["max_iterations"] * 100))

            # ── 1. Planning (reflection-guided after first iteration) ──
            actions = []
            if iter_num > 1 and last_eval:
                await _emit_log(f"🪞 Reflecting on iteration {iter_num - 1}...")
                reflection = await reflect_on_iteration(
                    query, iter_num, proj["max_iterations"],
                    last_eval, uncovered_questions,
                    _project_llm_call, provider, model,
                )
                actions = reflection.get("priority_actions", [])
                # Mark covered questions
                for cq in reflection.get("covered_questions", []):
                    if cq in uncovered_questions:
                        uncovered_questions.remove(cq)
                if not actions:
                    for fq in reflection.get("focus_queries", [])[:3]:
                        actions.append({"action": "web_search", "params": {"query": fq}})

            if not actions:
                # Standard planning with perspective context
                existing = json.dumps(proj["search_results"][-20:], ensure_ascii=False)[:3000]
                persp_ctx = "\n".join(
                    f"- {p['perspective']}: {p.get('description', '')}"
                    for p in perspectives
                ) if perspectives else "No specific perspectives"
                uncov_ctx = "\n".join(
                    f"- {q}" for q in uncovered_questions[:8]
                ) if uncovered_questions else "All questions covered"

                plan_prompt = [
                    {"role": "system", "content": (
                        "You are an autonomous research assistant. Plan the next actions.\n"
                        "Available actions: web_search, scrape_url, deep_dive, "
                        "download_asset, write_script, query_rag, smart_scrape, ask_model, academic_search, fetch_market_data\n"
                        "- web_search: searches Tavily, DDG, Scholar, Reddit, X, News in parallel\n"
                        "- deep_dive: recursive deep research on a sub-topic (params: {topic, depth, breadth})\n"
                        "- smart_scrape: scrape + LLM extraction of relevant content (params: {urls: [...]})\n"
                        "- ask_model: exploit the AI model's internal knowledge (params: {question: \"...\"})\n"
                        "- write_script: write and execute python code in a sandbox (params: {code: \"...\", description: \"...\"})\n"
                        "- academic_search: search Semantic Scholar for peer-reviewed papers (params: {query: \"...\", max_results: 10})\n"
                        "  Use academic_search for scientific topics, medical, technical, or when citations matter.\n"
                        "- fetch_market_data: fetch OHLCV market data from Binance (crypto), Yahoo (stocks), Dukascopy (forex).\n"
                        "  params: {symbol: \"BTCUSDT\", source: \"crypto\", interval: \"1d\", limit: 30}\n"
                        "  Use for financial analysis, market trends, price comparisons, trading research.\n"
                        "Return JSON array of actions.\n"
                        "Return ONLY valid JSON array, no markdown fences.\n"
                        + _get_dev_profile_summary()
                    )},
                    {"role": "user", "content": (
                        f"Research topic: {query}\n"
                        f"Iteration: {iter_num}/{proj['max_iterations']}\n"
                        f"Current score: {proj['current_score']}\n\n"
                        f"Perspectives to cover:\n{persp_ctx}\n\n"
                        f"Uncovered questions:\n{uncov_ctx}\n\n"
                        f"Previous results: {existing[:2000]}\n"
                        f"Plan 2-5 actions to improve research quality."
                    )},
                ]
                await _emit_log(f"🧠 Planning iteration {iter_num}...")
                plan_text = await _project_llm_call(plan_prompt, provider, model)
                try:
                    s = plan_text.find("[")
                    e = plan_text.rfind("]")
                    if s != -1 and e != -1:
                        raw_actions = json.loads(plan_text[s:e+1])
                        # Filter: LLM sometimes returns strings instead of dicts
                        actions = [a for a in raw_actions if isinstance(a, dict)]
                except Exception:
                    actions = [{"action": "web_search", "params": {"query": query}}]

            # ── 2. Execute actions ──
            for act in actions[:5]:
                # Bug 15: Check for stop request between each action
                if pid in _stop_requested:
                    _stop_requested.discard(pid)
                    proj["status"] = "paused"
                    _save(proj)
                    await _emit(pid, "status", {"status": "paused"})
                    return

                action_type = act.get("action", "")
                params = act.get("params", {})
                # Human-readable action label (not raw JSON)
                action_label = (
                    params.get("query")
                    or params.get("url")
                    or params.get("question")
                    or params.get("symbol")
                    or params.get("description")
                    or ""
                )[:120]
                await _emit_log(f"{action_type}: {action_label}")

                if action_type == "web_search":
                    results = await _do_web_search_with_retry(params.get("query", query))
                    proj["search_results"].extend(results)
                    await _emit_new_results(results)
                    urls = [r.get("url") for r in results if r.get("url")]
                    iter_data["actions"].append({"type": "web_search", "count": len(results), "urls": urls, "params": params, "timestamp": datetime.now(timezone.utc).isoformat()})
                    # ── Auto-save top web pages as assets ──
                    _saved_pages = 0
                    for r in results[:3]:
                        if _saved_pages >= 3:
                            break
                        page_url = r.get("url", "")
                        if page_url and not page_url.startswith(("rag://", "memory://", "sandbox://", "market://", "model://")):
                            try:
                                page_text = await _scrape_url(page_url)
                                if page_text and len(page_text) > 100:
                                    asset = await _save_text_asset(
                                        f"Web: {r.get('title', page_url[:40])[:60]}",
                                        page_text, "web_page", page_url, pid,
                                    )
                                    proj["assets"].append(asset)
                                    await _emit_new_asset(asset)
                                    _saved_pages += 1
                            except Exception:
                                pass
                    await _emit_log(f"   Found {len(results)} results", extra={"urls": urls})

                    # ── think_tool: Post-Search Reflection (open_deep_research) ──
                    # Pause and analyse what was found, identify gaps, decide next step.
                    try:
                        existing_summary = branch_summaries[-800:] if branch_summaries else ""
                        reflection = await reflect_after_search(
                            query=query,
                            new_results=results,
                            existing_findings=existing_summary,
                            iteration_num=iter_num,
                            llm_call=_project_llm_call,
                            provider=provider,
                            model=main_model or model,
                        )
                        # Store reflection in iteration data for transparency
                        iter_data.setdefault("reflections", []).append(reflection)

                        if reflection.get("found"):
                            await _emit_log(f"   💡 Found: {reflection['found'][:120]}")
                        if reflection.get("gaps"):
                            gaps_str = ', '.join(reflection['gaps'][:3])
                            await _emit_log(f"   🔍 Gaps: {gaps_str}")

                        # Early stopping: if the LLM is confident enough, skip
                        # remaining actions in this iteration and proceed to evaluation
                        if should_stop_early(reflection, min_confidence=0.87):
                            reason = reflection.get("stop_reason") or (
                                f"confidence={reflection.get('confidence', 0):.0%}"
                            )
                            await _emit_log(
                                f"   ✅ Early stop: {reason} — skipping further searches"
                            )
                            break  # Exit the actions loop for this iteration

                        # Use the reflection's refined query for the next planned action
                        if reflection.get("next_query") and reflection["next_query"] != query:
                            # Inject the refined query into any remaining web_search actions
                            for remaining_act in actions[actions.index(act) + 1:]:
                                if remaining_act.get("action") == "web_search":
                                    remaining_act.setdefault("params", {})["query"] = reflection["next_query"]
                                    break
                    except Exception as _re:
                        logger.warning("reflect_after_search step failed: %s", _re)
                    # ─────────────────────────────────────────────────────────

                elif action_type == "deep_dive":
                    topic = params.get("topic", query)
                    depth = min(int(params.get("depth", 3)), 4)
                    breadth = min(int(params.get("breadth", 4)), 5)
                    branch = await deep_research_branch(
                        topic, query, depth, breadth,
                        _do_web_search, _scrape_url, _project_llm_call,
                        provider, model, _emit_log,
                    )
                    branch_results = flatten_branch_results(branch)
                    proj["search_results"].extend(branch_results)
                    await _emit_new_results(branch_results)
                    branch_summaries += "\n\n" + flatten_branch_summaries(branch)
                    iter_data["actions"].append({
                        "type": "deep_dive", "topic": topic,
                        "branches": len(branch.get("sub_branches", [])),
                        "results": len(branch_results),
                        "timestamp": datetime.now(timezone.utc).isoformat()
                    })
                    await _emit_log(
                        f"   🌿 Deep dive: {len(branch_results)} results, "
                        f"{len(branch.get('sub_branches', []))} branches"
                    )

                elif action_type == "smart_scrape":
                    urls = params.get("urls", [])
                    if isinstance(urls, str):
                        urls = [urls]
                    if urls:
                        scraped = await batch_scrape(
                            urls[:5], query, _scrape_url, _project_llm_call,
                            provider, model,
                        )
                        for s in scraped:
                            sr_entry = {
                                "title": f"Smart-scraped: {s['url'][:60]}",
                                "snippet": s["relevant_extract"][:500],
                                "url": s["url"],
                                "full_text": s["relevant_extract"],
                                "key_facts": s.get("key_facts", []),
                                "source": "smart_scrape",
                            }
                            proj["search_results"].append(sr_entry)
                            await _emit_new_results([sr_entry])
                            asset = await _save_text_asset(f"Smart Scrape: {s['url'][:60]}", s["relevant_extract"], "smart_scrape", s["url"], pid)
                            proj["assets"].append(asset)
                            await _emit_new_asset(asset)
                        
                        scraped_urls = [s["url"] for s in scraped if s.get("url")]
                        iter_data["actions"].append({"type": "smart_scrape", "count": len(scraped), "urls": scraped_urls, "timestamp": datetime.now(timezone.utc).isoformat()})
                        await _emit_log(f"   🔍 Smart-scraped {len(scraped)} pages", extra={"urls": scraped_urls})

                elif action_type == "scrape_url":
                    url = params.get("url", "")
                    if url:
                        text = await _scrape_url(url)
                        if text:
                            sr_entry = {
                                "title": f"Scraped: {url[:60]}",
                                "snippet": text[:500], "url": url,
                                "full_text": text,
                            }
                            proj["search_results"].append(sr_entry)
                            await _emit_new_results([sr_entry])
                            asset = await _save_text_asset(f"Scraped: {url[:60]}", text, "scrape", url, pid)
                            proj["assets"].append(asset)
                            await _emit_new_asset(asset)
                            iter_data["actions"].append({"type": "scrape", "url": url, "timestamp": datetime.now(timezone.utc).isoformat()})
                            await _emit_log(f"   Scraped {len(text)} chars")

                elif action_type == "download_asset":
                    url = params.get("url", "")
                    if url:
                        asset = await _download_asset(url, pid)
                        if asset:
                            proj["assets"].append(asset)
                            await _emit_new_asset(asset)
                            iter_data["actions"].append({"type": "download", "name": asset["name"], "timestamp": datetime.now(timezone.utc).isoformat()})
                            await _emit_log(f"   Downloaded: {asset['name']}")

                elif action_type == "query_rag":
                    try:
                        from app.rag import explicit_rag_search
                        rag_q = params.get("query", query)
                        rag_ctx = explicit_rag_search(rag_q, k=3)
                        if rag_ctx:
                            sr_entry = {
                                "title": f"RAG: {rag_q[:50]}",
                                "snippet": rag_ctx[:500], "url": "rag://local",
                            }
                            proj["search_results"].append(sr_entry)
                            await _emit_new_results([sr_entry])
                            asset = await _save_text_asset(f"RAG Context: {rag_q[:50]}", rag_ctx, "rag_context", "local", pid)
                            proj["assets"].append(asset)
                            await _emit_new_asset(asset)
                            iter_data["actions"].append({"type": "rag", "query": rag_q, "timestamp": datetime.now(timezone.utc).isoformat()})
                            await _emit_log(f"   RAG returned context")
                    except Exception:
                        pass

                elif action_type == "write_script":
                    script_code = params.get("code", "")
                    if script_code:
                        try:
                            script_path = os.path.join(_proj_dir(pid), "temp_script.py")
                            with open(script_path, "w", encoding="utf-8") as f:
                                f.write(script_code)
                            python_exe = _venv_python(pid)
                            result = await asyncio.to_thread(
                                subprocess.run, [python_exe, script_path],
                                capture_output=True, text=True, timeout=30,
                                cwd=_proj_dir(pid),
                            )
                            output = (result.stdout + "\n" + result.stderr)[:2000].strip()
                            sr_entry = {
                                "title": f"🧪 Experiment: {params.get('description', 'script')[:50]}",
                                "snippet": output[:500],
                                "url": "sandbox://experiment",
                                "full_text": output,
                                "source": "experiment",
                                "code": script_code[:500],
                            }
                            proj["search_results"].append(sr_entry)
                            await _emit_new_results([sr_entry])
                            asset = await _save_text_asset(f"Script Experiment", f"Code:\n```python\n{script_code}\n```\n\nOutput:\n```\n{output}\n```", "script_experiment", "sandbox", pid)
                            proj["assets"].append(asset)
                            await _emit_new_asset(asset)
                            iter_data["actions"].append({"type": "script", "output": output, "code": script_code, "timestamp": datetime.now(timezone.utc).isoformat()})
                            await _emit_log(f"   🧪 Script output: {output[:200]}...", extra={"code": script_code, "output": output})
                        except Exception as e:
                            iter_data["actions"].append({"type": "script", "output": str(e), "code": script_code, "timestamp": datetime.now(timezone.utc).isoformat()})
                            await _emit_log(f"   ❌ Script error: {e}", extra={"code": script_code, "output": str(e)})

                elif action_type == "ask_model":
                    question = params.get("question", params.get("query", query))
                    knowledge_prompt = [
                        {"role": "system", "content": (
                            "You are a knowledgeable expert. Answer the question "
                            "using your training knowledge. Provide detailed, factual "
                            "information with specific data points, dates, names. "
                            "If uncertain, say so. Be thorough (500-1000 words)."
                        )},
                        {"role": "user", "content": (
                            f"Research context: {query}\n\n"
                            f"Question: {question}"
                        )},
                    ]
                    knowledge = await _project_llm_call(knowledge_prompt, provider, model)
                    if knowledge:
                        sr_entry = {
                            "title": f"🤖 Model Knowledge: {question[:60]}",
                            "snippet": knowledge[:500],
                            "url": "model://internal-knowledge",
                            "full_text": knowledge,
                            "source": "model_knowledge",
                        }
                        proj["search_results"].append(sr_entry)
                        await _emit_new_results([sr_entry])
                        asset = await _save_text_asset(f"Model Knowledge: {question[:60]}", knowledge, "model_knowledge", "internal", pid)
                        proj["assets"].append(asset)
                        await _emit_new_asset(asset)
                        iter_data["actions"].append({"type": "ask_model", "question": question[:100], "timestamp": datetime.now(timezone.utc).isoformat()})
                        await _emit_log(f"   🤖 Model knowledge: {len(knowledge)} chars")

                elif action_type == "academic_search":
                    # Feature #6: Dedicated Semantic Scholar search
                    academic_query = params.get("query", query)
                    await _emit_log(f"   📚 Semantic Scholar: {academic_query[:60]}")
                    try:
                        sem_results = await _semantic_scholar_search(
                            academic_query,
                            max_results=params.get("max_results", 10),
                        )
                        proj["search_results"].extend(sem_results)
                        await _emit_new_results(sem_results)
                        iter_data["actions"].append({
                            "type": "academic_search", "query": academic_query,
                            "count": len(sem_results),
                            "timestamp": datetime.now(timezone.utc).isoformat(),
                        })
                        await _emit_log(f"   📚 Found {len(sem_results)} academic papers")
                    except Exception as e:
                        await _emit_log(f"   ⚠️ Academic search failed: {e}")

                elif action_type == "fetch_market_data":
                    # Market data fetcher (crypto, stocks, forex)
                    symbol = params.get("symbol", "")
                    if symbol:
                        await _emit_log(f"   📈 Fetching market data: {symbol}")
                        try:
                            from app.tools_market import fetch_market_data
                            market_result = fetch_market_data(params)
                            if market_result.get("error"):
                                await _emit_log(f"   ⚠️ Market data error: {market_result['error']}")
                            else:
                                count = market_result.get("count", 0)
                                source = market_result.get("source", "unknown")
                                data_rows = market_result.get("data", [])
                                # Build a readable summary for the research context
                                cols = market_result.get("columns", [])
                                summary_lines = [f"OHLCV data for {symbol} ({source}, {count} candles):"]
                                summary_lines.append(f"Columns: {', '.join(cols)}")
                                # Show first and last 3 rows
                                for row in data_rows[:3]:
                                    summary_lines.append(str(row))
                                if len(data_rows) > 6:
                                    summary_lines.append("...")
                                for row in data_rows[-3:]:
                                    summary_lines.append(str(row))
                                summary = "\n".join(summary_lines)

                                sr_entry = {
                                    "title": f"📈 Market Data: {symbol} ({source})",
                                    "snippet": summary[:500],
                                    "url": f"market://{source}/{symbol}",
                                    "full_text": summary,
                                    "source": "market_data",
                                    "market_data": market_result,
                                }
                                proj["search_results"].append(sr_entry)
                                await _emit_new_results([sr_entry])
                                asset = await _save_text_asset(
                                    f"Market Data: {symbol}", summary,
                                    "market_data", f"{source}/{symbol}", pid,
                                )
                                proj["assets"].append(asset)
                                await _emit_new_asset(asset)
                                iter_data["actions"].append({
                                    "type": "fetch_market_data",
                                    "symbol": symbol,
                                    "source": source,
                                    "count": count,
                                    "timestamp": datetime.now(timezone.utc).isoformat(),
                                })
                                await _emit_log(f"   📈 {source}: {count} candles for {symbol}")
                        except Exception as e:
                            await _emit_log(f"   ⚠️ Market data fetch failed: {e}")

            # ── 3. Structured Evaluation ──
            _update_pm_task(1, "In Progress", int((iter_num - 0.5) / proj["max_iterations"] * 100))
            _update_pm_task(2, "In Progress", int((iter_num - 1) / proj["max_iterations"] * 100))
            await _emit_log("📊 Multi-criteria evaluation...")
            last_eval = await evaluate_structured(
                query, perspectives, proj["search_results"],
                len(proj["assets"]), _project_llm_call, provider, model,
            )
            proj["current_score"] = last_eval["overall"]
            iter_data["score"] = last_eval["overall"]
            iter_data["evaluation"] = last_eval.get("evaluation", "")
            iter_data["scores_detail"] = last_eval.get("scores", {})

            scores_display = " | ".join(
                f"{ax}: {v:.0%}" for ax, v in last_eval.get("scores", {}).items()
            )
            await _emit_log(f"   {scores_display} → Overall: {last_eval['overall']:.0%}")
            if last_eval.get("gaps"):
                await _emit_log(f"   Gaps: {', '.join(last_eval['gaps'][:3])}")

            # ── Feature #7: Context Budget Monitor ─────────────────────────────
            _log_context_budget(
                proj["search_results"], report_draft, branch_summaries, iter_num
            )
            # ───────────────────────────────────────────────────────────────

            iter_data["completed_at"] = datetime.now(timezone.utc).isoformat()
            proj["iterations"].append(iter_data)
            _save(proj)
            _update_pm_task(2, "In Progress", int(iter_num / proj["max_iterations"] * 100))

            # ── Iteration Merging: cumulative findings ─────────────────────
            try:
                cumulative_findings = await merge_iteration_findings(
                    cumulative_findings=cumulative_findings,
                    new_results=proj["search_results"][-30:],
                    new_branch_summaries=branch_summaries[-2000:],
                    query=query,
                    iteration_num=iter_num,
                    llm_call=_project_llm_call,
                    provider=provider,
                    model=model,
                )
                proj["cumulative_findings"] = cumulative_findings
                _save(proj)
                await _emit_log(f"🔗 Findings merged ({len(cumulative_findings)} chars cumulative)")
            except Exception as _mfe:
                logger.warning("Iteration merging failed: %s", _mfe)
            # ───────────────────────────────────────────────────────────────
            await _emit(pid, "iteration_end", {
                "iteration": iter_num, "score": proj["current_score"],
                "evaluation": iter_data.get("evaluation", ""),
                "scores_detail": iter_data.get("scores_detail", {}),
            })

            # ── Dynamic Outline Update (WebWeaver #3) ─────────────────────
            try:
                dynamic_outline = await update_dynamic_outline(
                    current_outline=dynamic_outline,
                    new_findings=proj["search_results"][-30:],
                    query=query,
                    iteration_num=iter_num,
                    llm_call=_project_llm_call,
                    provider=provider,
                    model=model,
                )
                if dynamic_outline.get("focus_next"):
                    await _emit_log(
                        f"📋 Outline updated — next focus: "
                        f"{', '.join(dynamic_outline['focus_next'][:3])}"
                    )
            except Exception as _oe:
                logger.warning("Dynamic outline update failed: %s", _oe)
            # ─────────────────────────────────────────────────────────────

            # ── IterResearch Condensation (DeepResearch #1) ───────────────
            if should_condense(iter_num, proj["search_results"], report_draft, branch_summaries):
                await _emit_log(
                    f"🗜️ Context condensation triggered (iter {iter_num}, "
                    f"{len(proj['search_results'])} results)..."
                )
                try:
                    # Use dedicated compression model for condensation (open_deep_research style)
                    condensed_results, core_findings, report_draft = \
                        await condense_research_context(
                            results=proj["search_results"],
                            report_draft=report_draft,
                            query=query,
                            iteration_num=iter_num,
                            eval_scores=last_eval.get("scores", {}),
                            llm_call=_project_llm_call,
                            provider=provider,
                            model=compression_model or model,
                            emit_fn=_emit_log,
                        )
                    proj["search_results"] = condensed_results
                    _save(proj)
                except Exception as _ce:
                    logger.warning("IterResearch condensation failed: %s", _ce)
            # ─────────────────────────────────────────────────────────────

            # ── Self-Referential Process Improvement (HyperAgents-inspired) ──
            # Every 2 iterations, let the LLM improve the process.md itself
            if iter_num % 2 == 0 and proj["current_score"] < 0.82:
                try:
                    current_process = _read_process(pid)
                    improved = await improve_process_md(
                        current_process=current_process,
                        query=query,
                        eval_result=last_eval,
                        iteration_num=iter_num,
                        max_iterations=proj["max_iterations"],
                        llm_call=_project_llm_call,
                        provider=provider,
                        model=model,
                    )
                    if improved:
                        _write_process(pid, improved)
                        await _emit_log(
                            f"🔄 Process self-improved (iter {iter_num}) "
                            f"— weakest axis: {last_eval.get('weakest_axis', '?')}",
                            extra={"process_updated": True},
                        )
                except Exception as _e:
                    logger.warning("Self-improvement step failed: %s", _e)
            # ─────────────────────────────────────────────────────────────────

            # ── 4. Check if done ──
            if proj["current_score"] >= proj["target_score"]:
                await _emit_log(
                    f"✅ Target score reached ({proj['current_score']:.0%}). "
                    "Generating report..."
                )
                _update_pm_task(1, "Done", 100)
                _update_pm_task(2, "Done", 100)
                break

            await asyncio.sleep(1)

        # ── Generate final report with citations ──
        _update_pm_task(1, "Done", 100)
        _update_pm_task(2, "Done", 100)
        _update_pm_task(3, "In Progress", 10)
        await _emit_progress(phase=ResearchPhase.REPORT)
        await _generate_report(
            pid, provider, model, perspectives, branch_summaries,
            dynamic_outline=dynamic_outline,
            report_draft=(
                (f"## Cumulative Research Findings (merged from {len(proj.get('iterations', []))} iterations)\n\n{cumulative_findings}\n\n" if cumulative_findings else "")
                + report_draft
            ),
            perspective_synthesis=perspective_synthesis,
        )
        _update_pm_task(3, "Done", 100)
        proj = _load(pid)
        proj["status"] = "completed"
        _save(proj)
        await _emit(pid, "status", {"status": "completed"})
        await _emit_log("🎉 Research completed!")

        # ── Archive Winning Strategy (HyperAgents-inspired) ────────────────
        try:
            scores_per_iter = [
                it.get("score", 0) for it in proj.get("iterations", [])
            ]
            all_actions = []
            for it in proj.get("iterations", []):
                all_actions.extend(it.get("actions", []))
            save_strategy(
                query=query,
                domain=proj.get("query_domain", classify_query_domain(query)),
                action_sequence=all_actions,
                scores_per_iteration=scores_per_iter,
                final_score=proj.get("current_score", 0),
                num_iterations=len(proj.get("iterations", [])),
                profile_id=proj.get("profile_id", ""),
            )
            await _emit_log("📚 Strategy saved to archive for future warm-starts")
        except Exception as _e:
            logger.warning("Failed to archive strategy: %s", _e)
        # ──────────────────────────────────────────────────────────────────

        unregister_task(pid)

    except asyncio.CancelledError:
        proj = _load(pid) or proj
        proj["status"] = "paused"
        _save(proj)
        await _emit(pid, "status", {"status": "paused"})
        unregister_task(pid)
    except Exception as e:
        logger.error("Research loop error for %s: %s", pid, e, exc_info=True)
        proj = _load(pid) or proj
        proj["status"] = "error"
        proj["error_msg"] = str(e)
        _save(proj)
        await _emit(pid, "status", {"status": "error"})
        await _emit(pid, "log", {"msg": f"❌ Error: {e}"})
        unregister_task(pid)
    finally:
        _running.pop(pid, None)
        _stop_requested.discard(pid)
        _sse_queues.pop(pid, None)


async def _generate_report(
    pid: str, provider: str = "", model: str = "",
    perspectives: list[dict] | None = None,
    branch_summaries: str = "",
    dynamic_outline: dict | None = None,
    report_draft: str = "",
    perspective_synthesis: str = "",
):
    proj = _load(pid)
    if not proj:
        return
    await _emit(pid, "log", {"msg": "📝 Generating report with citations..."})

    # Project-bound LLM call for token tracking
    async def _report_llm_call(messages, prov="", mdl=""):
        return await _llm_call(messages, prov, mdl, pid=pid)

    async def _emit_section_log(msg: str):
        await _emit(pid, "log", {"msg": msg})

    report = ""

    # ── GPT-Researcher-inspired: sectioned report if outline has ≥3 sections ──
    outline_sections = (dynamic_outline or {}).get("sections", [])
    use_sectioned = len(outline_sections) >= 3 and len(proj.get("search_results", [])) >= 10

    if use_sectioned:
        await _emit(pid, "log", {
            "msg": f"✍️ Generating sectioned report ({len(outline_sections)} sections, GPT-Researcher style)..."
        })
        try:
            report = await generate_sectioned_report(
                query=proj["query"],
                sections=outline_sections,
                sources=proj["search_results"],
                branch_summaries=branch_summaries,
                score=proj.get("current_score", 0),
                llm_call=_report_llm_call,
                provider=provider,
                model=model,
                perspective_synthesis=perspective_synthesis,
                emit_fn=_emit_section_log,
            )
        except Exception as _se:
            logger.warning("Sectioned report failed, falling back to monolithic: %s", _se)
            await _emit(pid, "log", {"msg": f"⚠️ Sectioned report failed ({_se}), using standard generator..."})
            report = ""

    # Fallback / standard monolithic report generator
    if not report:
        report = await generate_report_with_citations(
            query=proj["query"],
            results=proj["search_results"],
            assets=proj.get("assets", []),
            perspectives=perspectives or proj.get("perspectives", []),
            branch_summaries=branch_summaries,
            score=proj.get("current_score", 0),
            num_iterations=len(proj.get("iterations", [])),
            llm_call=_report_llm_call,
            provider=provider,
            model=model,
            dynamic_outline=dynamic_outline or {},
            report_draft=report_draft,
            perspective_synthesis=perspective_synthesis,
        )

    # ── Sanitize mermaid diagrams before saving ──
    try:
        report = sanitize_mermaid_in_report(report)
    except Exception as _sme:
        logger.warning("Mermaid sanitization failed: %s", _sme)

    proj["report_md"] = report
    proj["report_mode"] = "sectioned" if use_sectioned else "monolithic"
    report_path = os.path.join(_proj_dir(pid), "report.md")
    with open(report_path, "w", encoding="utf-8") as f:
        f.write(report)
    _save(proj)
    await _emit(pid, "report_ready", {"length": len(report), "mode": proj["report_mode"]})


# ── API Endpoints ──

@router.get("/archive/stats")
async def api_archive_stats():
    """Return global strategy archive statistics (HyperAgents-inspired)."""
    return get_archive_stats()


@router.get("/profiles")
async def api_list_profiles():
    return {"profiles": list_profiles()}

@router.post("/profiles")
async def api_create_profile(request: Request):
    data = await request.json()
    prof = create_profile(data)
    return {"status": "created", "profile": prof}

@router.delete("/profiles/{pid}")
async def api_delete_profile(pid: str):
    if delete_profile(pid):
        return {"status": "deleted"}
    raise HTTPException(400, "Cannot delete profile")

@router.get("/projects/{pid}/process")
async def api_get_process(pid: str):
    if not _load(pid):
        raise HTTPException(404, "Project not found")
    return {"process_md": _read_process(pid)}

@router.put("/projects/{pid}/process")
async def api_update_process(pid: str, request: Request):
    if not _load(pid):
        raise HTTPException(404, "Project not found")
    data = await request.json()
    _write_process(pid, data.get("process_md", ""))
    return {"status": "updated"}

@router.post("/projects/{pid}/sandbox")
async def api_sandbox_execute(pid: str, request: Request):
    if not _load(pid):
        raise HTTPException(404, "Project not found")
    data = await request.json()
    code = data.get("code", "")
    if not code:
        raise HTTPException(400, "Code is required")
    
    script_path = os.path.join(_proj_dir(pid), "sandbox_script.py")
    with open(script_path, "w", encoding="utf-8") as f:
        f.write(code)
    
    python_exe = _venv_python(pid)
    try:
        proc = await asyncio.create_subprocess_exec(
            python_exe, script_path,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=_proj_dir(pid)
        )
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=60.0)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            return {"status": "error", "stderr": "Execution timed out after 60 seconds", "stdout": "", "exit_code": 124}
            
        return {
            "status": "success" if proc.returncode == 0 else "error",
            "stdout": stdout.decode("utf-8", errors="replace"),
            "stderr": stderr.decode("utf-8", errors="replace"),
            "exit_code": proc.returncode
        }
    except Exception as e:
        return {"status": "error", "stderr": str(e), "stdout": "", "exit_code": -1}

@router.post("/projects/{pid}/install-deps")
async def api_sandbox_install(pid: str, request: Request):
    if not _load(pid):
        raise HTTPException(404, "Project not found")
    data = await request.json()
    packages = data.get("packages", [])
    if not packages:
        raise HTTPException(400, "Packages list is required")
    
    python_exe = _venv_python(pid)
    try:
        proc = await asyncio.create_subprocess_exec(
            python_exe, "-m", "pip", "install", *packages,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120.0)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            return {"status": "error", "stderr": "Installation timed out after 120 seconds", "stdout": ""}
            
        return {
            "status": "success" if proc.returncode == 0 else "error",
            "stdout": stdout.decode("utf-8", errors="replace"),
            "stderr": stderr.decode("utf-8", errors="replace")
        }
    except Exception as e:
        return {"status": "error", "stderr": str(e), "stdout": ""}

@router.get("/projects")
async def list_projects():
    return {"projects": _list_projects()}

@router.post("/projects")
async def create_project(request: Request):
    data = await request.json()
    query = data.get("query", "").strip()
    if not query:
        raise HTTPException(400, "Research query is required")
    proj = _new_project(
        query,
        data.get("provider", ""),
        data.get("model", ""),
        profile_id=data.get("profile_id", ""),
        profile_data=data.get("profile_data"),
    )
    # Override target_score / max_iterations if explicitly provided in the request
    if "target_score" in data:
        proj["target_score"] = float(data["target_score"])
    if "max_iterations" in data:
        proj["max_iterations"] = int(data["max_iterations"])
    _save(proj)
    return {"status": "created", "project": proj}

@router.get("/projects/{pid}")
async def get_project(pid: str):
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
        
    assets = proj.get("assets", [])
    known_names = {a.get("name") for a in assets}
    assets_dir = os.path.join(_proj_dir(pid), "assets")
    if os.path.isdir(assets_dir):
        for fname in os.listdir(assets_dir):
            if fname not in known_names:
                fpath = os.path.join(assets_dir, fname)
                if os.path.isfile(fpath):
                    ext = os.path.splitext(fname)[1].lstrip(".") or "bin"
                    size = os.path.getsize(fpath)
                    assets.append({
                        "name": fname,
                        "path": fpath,
                        "type": ext,
                        "url": "local",
                        "size": size,
                        "timestamp": datetime.now(timezone.utc).isoformat()
                    })
    proj["assets"] = assets
    
    return {"project": proj}

@router.delete("/projects/{pid}")
async def delete_project(pid: str):
    if pid in _running:
        task = _running.pop(pid, None)
        if task and not task.done():
            task.cancel()
            try:
                await asyncio.wait_for(task, timeout=2.0)
            except (asyncio.CancelledError, asyncio.TimeoutError, Exception):
                pass
    _stop_requested.discard(pid)
    _sse_queues.pop(pid, None)
    d = os.path.join(RESEARCH_DIR, pid)
    if os.path.isdir(d):
        shutil.rmtree(d)
    return {"status": "deleted"}

@router.post("/projects/{pid}/start")
async def start_research(pid: str):
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
    if pid in _running:
        task = _running[pid]
        if not task.done():
            return {"status": "already_running"}
        # Bug 14: Task is finished but not cleaned up — wait for it to
        # fully shut down before removing the entry and starting a new one
        old_task = _running.pop(pid, None)
        if old_task:
            try:
                await asyncio.wait_for(asyncio.shield(old_task), timeout=3.0)
            except Exception:
                pass

    # Set status to running immediately on disk (Bug 4)
    proj["status"] = "running"
    _save(proj)

    _stop_requested.discard(pid)
    _sse_queues[pid] = asyncio.Queue()
    task = asyncio.create_task(_research_loop(pid))
    _running[pid] = task
    return {"status": "started"}

@router.post("/projects/{pid}/stop")
async def stop_research(pid: str):
    from app.tools.task_manager import unregister_task
    # Signal the loop to stop cleanly at the next iteration check
    _stop_requested.add(pid)
    if pid in _running:
        _running[pid].cancel()
        _running.pop(pid, None)
    # Bug 13: Do NOT pop _sse_queues here — the CancelledError handler in
    # _research_loop needs the queue to emit the final "paused" status event.
    # The queue is cleaned up in the finally block of _research_loop.
    proj = _load(pid)
    if proj:
        proj["status"] = "paused"
        _save(proj)
    unregister_task(pid)
    return {"status": "stopped"}

@router.get("/projects/{pid}/status")
async def research_status_sse(pid: str):
    # Bug 11: Only create a queue if the project is actually running,
    # otherwise send a one-shot status event and close.
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")

    if pid not in _sse_queues:
        # No active queue — check if actually running
        if pid not in _running or _running[pid].done():
            # Send current status as a one-shot and close
            async def gen_oneshot():
                yield {"data": json.dumps({"type": "status", "status": proj.get("status", "idle")})}
            return EventSourceResponse(gen_oneshot())
        _sse_queues[pid] = asyncio.Queue()

    async def gen():
        while True:
            # Re-fetch queue reference each iteration in case it was replaced
            q = _sse_queues.get(pid)
            if q is None:
                # Queue was cleaned up — the research finished
                yield {"data": json.dumps({"type": "status", "status": "completed"})}
                break
            try:
                msg = await asyncio.wait_for(q.get(), timeout=30)
                yield {"data": msg}
            except asyncio.TimeoutError:
                # Check if the task is still alive before sending keepalive
                task = _running.get(pid)
                if task is None or task.done():
                    # Task ended but queue wasn't cleaned — send final status and close
                    final_proj = _load(pid)
                    final_status = final_proj.get("status", "idle") if final_proj else "idle"
                    yield {"data": json.dumps({"type": "status", "status": final_status})}
                    _sse_queues.pop(pid, None)
                    break
                yield {"data": json.dumps({"type": "ping"})}
            except Exception:
                break
    return EventSourceResponse(gen())

@router.post("/projects/{pid}/evaluate")
async def force_evaluate(pid: str):
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
    # Quick re-evaluation
    results_summary = "\n".join(
        f"- {r.get('title','')}: {r.get('snippet','')[:80]}"
        for r in proj["search_results"][-20:]
    )[:3000]
    eval_prompt = [
        {"role": "system", "content": "Evaluate research quality. Return JSON: {\"score\": 0.0-1.0, \"evaluation\": \"...\"} ONLY."},
        {"role": "user", "content": f"Topic: {proj['query']}\nResults:\n{results_summary}"}
    ]
    text = await _llm_call(eval_prompt, proj.get("provider",""), proj.get("model",""), pid=pid)
    try:
        s = text.find("{"); e = text.rfind("}")
        if s != -1 and e != -1:
            d = json.loads(text[s:e+1])
            proj["current_score"] = float(d.get("score", proj["current_score"]))
    except Exception:
        pass
    _save(proj)
    return {"score": proj["current_score"]}

@router.get("/projects/{pid}/report")
async def get_report(pid: str):
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
    return {"report": proj.get("report_md", ""), "title": proj.get("title", "")}

@router.post("/projects/{pid}/export")
async def export_report(pid: str, request: Request):
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
    data = await request.json()
    fmt = data.get("format", "md")
    report = proj.get("report_md", "")
    if not report:
        raise HTTPException(400, "No report generated yet")

    pdir = _proj_dir(pid)
    if fmt == "md":
        path = os.path.join(pdir, "report.md")
        with open(path, "w", encoding="utf-8") as f:
            f.write(report)
        return FileResponse(path, filename=f"research_{proj['title'][:30]}.md")
    elif fmt == "docx":
        try:
            from app.tools.research_export import export_docx
            path = await export_docx(proj, report, pdir)
            return FileResponse(path, filename=f"research_{proj['title'][:30]}.docx")
        except ImportError as e:
            raise HTTPException(500, f"DOCX export dependency missing: {e}")
        except Exception as e:
            logger.error("DOCX export failed: %s", e)
            raise HTTPException(500, f"DOCX export failed: {e}")
    elif fmt == "pdf":
        try:
            from app.tools.research_export import export_pdf
            path = await export_pdf(proj, report, pdir)
            return FileResponse(path, filename=f"research_{proj['title'][:30]}.pdf")
        except ImportError as e:
            raise HTTPException(500, f"PDF export dependency missing: {e}")
        except Exception as e:
            logger.error("PDF export failed: %s", e)
            raise HTTPException(500, f"PDF export failed: {e}")
    elif fmt == "pptx":
        return await _export_research_pptx(proj, pdir)
    raise HTTPException(400, f"Unsupported format: {fmt}")


async def _export_research_pptx(proj: dict, pdir: str) -> FileResponse:
    """Export research report as a professional PPTX presentation.

    Parses the markdown report into structured slides:
    - Title slide with research query and metadata
    - One slide per H2 section
    - Bullet points for content
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        raise HTTPException(500, "python-pptx not installed")

    report = proj.get("report_md", "")
    title = proj.get("title", "Research Report")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 30, 30)
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    p2 = tf.add_paragraph()
    score = proj.get("current_score", 0)
    iters = len(proj.get("iterations", []))
    p2.text = f"Score: {score:.0%} | {iters} iterations | Clawzd Research Studio"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.alignment = PP_ALIGN.CENTER

    # --- Content Slides (one per H2 section) ---
    sections = _split_report_sections(report)
    for section_title, section_body in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Section title
        txTitle = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1),
        )
        tf_title = txTitle.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = section_title
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(30, 30, 30)

        # Section content (bullet points)
        txBody = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12), Inches(5.5),
        )
        tf_body = txBody.text_frame
        tf_body.word_wrap = True

        lines = section_body.strip().split("\n")
        first_line = True
        for line in lines:
            # Skip marker lines (__CHART__, __TABLE__, etc.)
            if line.strip().startswith("__") and line.strip().endswith("__"):
                continue
            if not line.strip():
                continue

            if first_line:
                p = tf_body.paragraphs[0]
                first_line = False
            else:
                p = tf_body.add_paragraph()

            # Handle bullet points
            if line.strip().startswith("- "):
                p.text = "• " + line.strip()[2:]
                p.level = 1
            elif line.strip().startswith("### "):
                p.text = line.strip()[4:]
                p.font.bold = True
                p.font.size = Pt(18)
                continue
            else:
                p.text = line.strip()

            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_after = Pt(6)

    path = os.path.join(pdir, "report.pptx")
    def _save_pptx():
        prs.save(path)
    await asyncio.to_thread(_save_pptx)
    return FileResponse(
        path,
        filename=f"research_{proj['title'][:30]}.pptx",
    )


def _split_report_sections(report_md: str) -> list[tuple[str, str]]:
    """Split a markdown report into (title, body) sections at H2 boundaries."""
    sections = []
    current_title = ""
    current_body = []

    for line in report_md.split("\n"):
        if line.startswith("## "):
            if current_title or current_body:
                sections.append((current_title, "\n".join(current_body)))
            current_title = line[3:].strip()
            current_body = []
        elif line.startswith("# "):
            # H1 is the report title, skip it for slides
            continue
        else:
            current_body.append(line)

    if current_title or current_body:
        sections.append((current_title, "\n".join(current_body)))

    return sections


@router.post("/projects/{pid}/to-presentation")
async def research_to_presentation(pid: str, request: Request):
    """Convert research report directly to a professional presentation.

    Bridge between Research Studio and Presentation Studio.
    Takes the research report markdown and generates a PPTX file
    with structured slides, one per section.
    """
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
    report = proj.get("report_md", "")
    if not report:
        raise HTTPException(400, "No report available — run research first")

    pdir = _proj_dir(pid)
    return await _export_research_pptx(proj, pdir)

@router.get("/projects/{pid}/export-zip")
async def export_zip(pid: str):
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
    import zipfile, io
    pdir = _proj_dir(pid)
    # Ensure report.md exists
    report_path = os.path.join(pdir, "report.md")
    if proj.get("report_md") and not os.path.exists(report_path):
        with open(report_path, "w", encoding="utf-8") as f:
            f.write(proj["report_md"])
    buf = io.BytesIO()
    def _make_zip():
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for root, dirs, files in os.walk(pdir):
                for fname in files:
                    fpath = os.path.join(root, fname)
                    arcname = os.path.relpath(fpath, pdir)
                    zf.write(fpath, arcname)
    await asyncio.to_thread(_make_zip)
    buf.seek(0)
    safe_title = proj.get("title", "research")[:30].replace(" ", "_")
    return StreamingResponse(buf, media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="research_{safe_title}.zip"'})

@router.get("/projects/{pid}/assets")
async def list_assets(pid: str):
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
        
    assets = proj.get("assets", [])
    known_names = {a.get("name") for a in assets}
    
    assets_dir = os.path.join(_proj_dir(pid), "assets")
    if os.path.isdir(assets_dir):
        for fname in os.listdir(assets_dir):
            if fname not in known_names:
                fpath = os.path.join(assets_dir, fname)
                if os.path.isfile(fpath):
                    ext = os.path.splitext(fname)[1].lstrip(".") or "bin"
                    size = os.path.getsize(fpath)
                    assets.append({
                        "name": fname,
                        "path": fpath,
                        "type": ext,
                        "url": "local",
                        "size": size,
                        "timestamp": datetime.now(timezone.utc).isoformat()
                    })
                    
    return {"assets": assets}

@router.get("/projects/{pid}/iterations")
async def list_iterations(pid: str):
    proj = _load(pid)
    if not proj:
        raise HTTPException(404, "Project not found")
    return {"iterations": proj.get("iterations", [])}


@router.get("/analytics")
async def research_analytics():
    """Return aggregated research analytics for the dashboard."""
    from app.tools.research_analytics import get_research_analytics
    return get_research_analytics()
