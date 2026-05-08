"""
Clawzd — Document creation tools.
Supports generating PowerPoint, PDF, Markdown, Word, and Excel documents.
Includes a document translation endpoint for PPTX / DOCX files.
"""

from fastapi import APIRouter, Request, HTTPException, UploadFile, File, Form
import os
import uuid
from datetime import datetime
from config import DATA_DIR
import logging

logger = logging.getLogger("clawzd.tools_document")

router = APIRouter()

def get_documents_dir():
    docs_dir = os.path.join(DATA_DIR, "documents")
    os.makedirs(docs_dir, exist_ok=True)
    return docs_dir

def _generate_markdown(content: str, title: str, filepath: str):
    with open(filepath, "w", encoding="utf-8") as f:
        if title:
            f.write(f"# {title}\n\n")
        f.write(content)
    return filepath

def _generate_word(content: str, title: str, filepath: str):
    try:
        from docx import Document
        doc = Document()
        if title:
            doc.add_heading(title, 0)
        doc.add_paragraph(content)
        doc.save(filepath)
        return filepath
    except ImportError:
        raise Exception("python-docx is not installed.")

def _generate_excel(content: str, title: str, filepath: str):
    try:
        import xlsxwriter
        workbook = xlsxwriter.Workbook(filepath)
        worksheet = workbook.add_worksheet()
        if title:
            worksheet.write(0, 0, title)

        lines = content.split('\n')
        row = 1 if title else 0
        for line in lines:
            cols = line.split(',')
            for col_idx, col_val in enumerate(cols):
                worksheet.write(row, col_idx, col_val.strip())
            row += 1

        workbook.close()
        return filepath
    except ImportError:
        raise Exception("XlsxWriter is not installed.")

def _generate_powerpoint(content: str, title: str, filepath: str):
    try:
        from pptx import Presentation
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]

        if title:
            title_shape.text = title

        subtitle.text = content[:500] + ("..." if len(content) > 500 else "")
        prs.save(filepath)
        return filepath
    except ImportError:
        raise Exception("python-pptx is not installed.")

def _generate_pdf(content: str, title: str, filepath: str):
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)

        if title:
            pdf.set_font("Arial", 'B', size=16)
            pdf.cell(200, 10, txt=title.encode('latin-1', 'replace').decode('latin-1'), ln=1, align='C')
            pdf.set_font("Arial", size=12)

        # Write content handling newlines
        # convert content to latin-1 due to fpdf basic font limitations unless we load a unicode font
        encoded_content = content.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 10, txt=encoded_content)
        pdf.output(filepath)
        return filepath
    except ImportError:
        raise Exception("fpdf2 is not installed.")
    except Exception as e:
        raise Exception(f"PDF generation failed: {e}")

async def create_document_core(format_type: str, content: str, title: str = "") -> dict:
    format_type = format_type.lower()
    docs_dir = get_documents_dir()

    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    file_id = uuid.uuid4().hex[:6]

    if format_type not in ["markdown", "md", "word", "docx", "excel", "xlsx", "powerpoint", "pptx", "pdf"]:
        return {"error": f"Unsupported format: {format_type}"}

    ext = format_type
    if format_type == "markdown": ext = "md"
    if format_type == "word": ext = "docx"
    if format_type == "excel": ext = "xlsx"
    if format_type == "powerpoint": ext = "pptx"

    filename = f"doc_{timestamp}_{file_id}.{ext}"
    filepath = os.path.join(docs_dir, filename)

    try:
        if ext == "md":
            _generate_markdown(content, title, filepath)
        elif ext == "docx":
            _generate_word(content, title, filepath)
        elif ext == "xlsx":
            _generate_excel(content, title, filepath)
        elif ext == "pptx":
            _generate_powerpoint(content, title, filepath)
        elif ext == "pdf":
            _generate_pdf(content, title, filepath)

        return {"status": "ok", "filename": filename, "format": ext, "path": filepath, "url": f"/data/documents/{filename}"}
    except Exception as e:
        logger.error(f"Error generating document: {e}")
        return {"error": str(e)}


@router.post("/create")
async def create_document_endpoint(request: Request):
    """Create a document from text content via HTTP."""
    data = await request.json()
    format_type = data.get("format_type", "markdown")
    content = data.get("content", "")
    title = data.get("title", "")
    if not content.strip():
        raise HTTPException(400, "Content is required")
    result = await create_document_core(format_type, content, title)
    if "error" in result:
        raise HTTPException(500, result["error"])
    return result


@router.post("/translate-upload")
async def translate_document_upload(
    file: UploadFile = File(...),
    target_language: str = Form("French"),
    source_language: str = Form(""),
    provider: str = Form(""),
    model: str = Form(""),
):
    """Translate a full PPTX or DOCX document using the DocTranslatorSkill.

    Form fields:
        file: The .pptx or .docx file to translate.
        target_language: Desired output language (default: French).
        source_language: Optional source language hint.
        provider: LLM provider override (e.g. 'ollama', 'openai').
        model: LLM model override.
    """
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in (".pptx", ".docx"):
        raise HTTPException(400, "Only .pptx and .docx files are supported.")

    # Save the upload to a temp location in the documents dir
    docs_dir = get_documents_dir()
    tmp_id = uuid.uuid4().hex[:8]
    tmp_path = os.path.join(docs_dir, f"_upload_{tmp_id}{ext}")

    try:
        content_bytes = await file.read()
        with open(tmp_path, "wb") as fh:
            fh.write(content_bytes)

        # Execute via the skill registry
        from app.skill_registry import get_registry
        from app.skill_model import SkillContext

        registry = get_registry()
        context = SkillContext(
            user_message=f"Translate document to {target_language}",
            provider=provider or "ollama",
            model=model or "",
            data_dir=DATA_DIR,
        )
        params = {
            "file_path": tmp_path,
            "target_language": target_language,
        }
        if source_language:
            params["source_language"] = source_language

        result = await registry.execute("doc_translator", params, context)

        if not result.success:
            raise HTTPException(500, result.error or "Translation failed")

        return {
            "status": "ok",
            **result.data,
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error("translate-upload error: %s", e, exc_info=True)
        raise HTTPException(500, str(e))
    finally:
        # Clean up the temporary upload file
        try:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        except Exception:
            pass
